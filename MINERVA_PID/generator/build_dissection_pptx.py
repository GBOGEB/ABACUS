#!/usr/bin/env python3
"""
build_dissection_pptx.py
Builds MINERVA_PID_DISSECTION.pptx  -- a 14-slide visual dissection of the
v3 P&ID deliverable for the MINERVA CryoCell project (SCK CEN / Mott MacDonald).

Pure tooling: python-pptx only. Image assets are rendered beforehand into
deck_assets/ by /tmp/render_deck_pngs.py + /tmp/crops.py (cairosvg + PIL).
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from PIL import Image

ROOT = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
ASSETS = os.path.join(ROOT, "deck_assets")
OUT = os.path.join(ROOT, "output_v3", "MINERVA_PID_DISSECTION.pptx")

# ---- palette ----
NAVY   = RGBColor(0x0B, 0x21, 0x47)
BLUE   = RGBColor(0x00, 0x33, 0xCC)
RED    = RGBColor(0xE0, 0x00, 0x00)
CYAN   = RGBColor(0x00, 0xA6, 0xBD)
GREEN  = RGBColor(0x00, 0xA0, 0x00)
GREY   = RGBColor(0x55, 0x55, 0x55)
LGREY  = RGBColor(0xE8, 0xEC, 0xF1)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
DARK   = RGBColor(0x1A, 0x1A, 0x1A)
ORANGE = RGBColor(0xC0, 0x60, 0x00)

prs = Presentation()
prs.slide_width  = Inches(13.333)   # 16:9
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]

SW = prs.slide_width
SH = prs.slide_height


def _add_rect(slide, x, y, w, h, color, line=None):
    from pptx.enum.shapes import MSO_SHAPE
    sp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    sp.fill.solid(); sp.fill.fore_color.rgb = color
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line; sp.line.width = Pt(1)
    sp.shadow.inherit = False
    return sp


def _txt(slide, x, y, w, h, text, size=18, color=DARK, bold=False,
         align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, italic=False, font="Calibri"):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame; tf.word_wrap = True
    tf.vertical_anchor = anchor
    lines = text.split("\n")
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        r = p.add_run(); r.text = ln
        r.font.size = Pt(size); r.font.bold = bold; r.font.italic = italic
        r.font.color.rgb = color; r.font.name = font
    return tb


def _bullets(slide, x, y, w, h, items, size=15, color=DARK, gap=4):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame; tf.word_wrap = True
    for i, (lvl, txt, *opt) in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.level = lvl
        p.space_after = Pt(gap)
        bullet = "\u2022 " if lvl == 0 else "\u2013 "
        r = p.add_run(); r.text = bullet + txt
        r.font.size = Pt(size - lvl * 1)
        r.font.color.rgb = opt[0] if opt else color
        r.font.name = "Calibri"
        if opt and len(opt) > 1 and opt[1]:
            r.font.bold = True
    return tb


def _img_fit(slide, path, x, y, w, h, frame=True, frame_color=GREY):
    """Place image fitting inside (x,y,w,h) preserving aspect, centered."""
    im = Image.open(path); iw, ih = im.size
    ar = iw / ih; box_ar = w / h
    if ar > box_ar:
        nw = w; nh = int(w / ar)
    else:
        nh = h; nw = int(h * ar)
    nx = x + (w - nw) // 2; ny = y + (h - nh) // 2
    if frame:
        _add_rect(slide, nx - Emu(9144), ny - Emu(9144),
                  nw + Emu(18288), nh + Emu(18288), WHITE, line=frame_color)
    slide.shapes.add_picture(path, nx, ny, nw, nh)
    return (nx, ny, nw, nh)


def header(slide, idx, title, subtitle=None, accent=NAVY):
    _add_rect(slide, 0, 0, SW, Inches(0.95), accent)
    _add_rect(slide, 0, Inches(0.95), SW, Emu(38100), CYAN)
    _txt(slide, Inches(0.45), Inches(0.10), Inches(11.5), Inches(0.55),
         title, size=26, color=WHITE, bold=True, anchor=MSO_ANCHOR.MIDDLE)
    if subtitle:
        _txt(slide, Inches(0.47), Inches(0.62), Inches(11.5), Inches(0.30),
             subtitle, size=12, color=RGBColor(0xBF, 0xD3, 0xF2),
             anchor=MSO_ANCHOR.MIDDLE)
    # slide number badge
    _txt(slide, Inches(12.4), Inches(0.10), Inches(0.8), Inches(0.6),
         f"{idx:02d}", size=22, color=CYAN, bold=True,
         align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)


def footer(slide):
    _txt(slide, Inches(0.45), Inches(7.12), Inches(8), Inches(0.3),
         "MINERVA CryoCell  \u00b7  P&ID v3 Dissection  \u00b7  SCK CEN / Mott MacDonald",
         size=9, color=GREY)
    _txt(slide, Inches(9.0), Inches(7.12), Inches(3.9), Inches(0.3),
         "AD_01.16  \u00b7  ISO 10628  \u00b7  ISA 5.1", size=9, color=GREY,
         align=PP_ALIGN.RIGHT)


def asset(name):
    return os.path.join(ASSETS, name)

# =====================================================================
# SLIDE 1 -- Title
# =====================================================================
s = prs.slides.add_slide(BLANK)
_add_rect(s, 0, 0, SW, SH, NAVY)
_add_rect(s, 0, Inches(2.55), SW, Inches(0.06), CYAN)
_add_rect(s, 0, Inches(4.55), SW, Inches(0.02), RGBColor(0x33,0x4E,0x7A))
_txt(s, Inches(0.8), Inches(1.05), Inches(11.7), Inches(0.5),
     "MINERVA CryoCell \u2014 MYRRHA Phase 1", size=18, color=CYAN, bold=True)
_txt(s, Inches(0.8), Inches(1.55), Inches(11.7), Inches(1.1),
     "P&ID v3 \u2014 Drawing Dissection", size=44, color=WHITE, bold=True)
_txt(s, Inches(0.82), Inches(2.75), Inches(11.7), Inches(0.8),
     "Layer hierarchy \u00b7 piping classes \u00b7 instrumentation \u00b7 monochrome \u00b7 A3 legibility\n"
     "QCELL (=NA.PS01_PFB712)  &  RFCELL (=NA.PS01_PFB713)",
     size=16, color=RGBColor(0xCF,0xDD,0xF2))
# meta grid
meta = [
    ("Standard", "SCK CEN AD_01.16 \u00b7 ISO 10628 \u00b7 ISA 5.1"),
    ("Consultant", "Mott MacDonald \u2014 Bristol, UK"),
    ("Client", "SCK CEN \u2014 Boeretang 200, 2400 Mol, BE"),
    ("MMD Project", "411066"),
    ("Sheets", "4 sheets \u00d7 4 variants = 16 SVG + 16 PDF (A3 landscape)"),
    ("Status", "S2 \u2014 FOR ACCEPTANCE  \u00b7  RESTRICTED"),
]
y = Inches(4.75)
for i, (k, v) in enumerate(meta):
    ry = y + Inches(0.42) * i
    _txt(s, Inches(0.85), ry, Inches(2.4), Inches(0.4), k, size=13,
         color=CYAN, bold=True, anchor=MSO_ANCHOR.MIDDLE)
    _txt(s, Inches(3.35), ry, Inches(9.0), Inches(0.4), v, size=13,
         color=WHITE, anchor=MSO_ANCHOR.MIDDLE)
_txt(s, Inches(0.8), Inches(7.12), Inches(11), Inches(0.3),
     "Revision C1 \u00b7 2026-06 \u00b7 v3 refinement (layers / legibility / monochrome / default views)",
     size=10, color=RGBColor(0x9F,0xB4,0xD4))

# =====================================================================
# SLIDE 2 -- Deliverable map / layer hierarchy
# =====================================================================
s = prs.slides.add_slide(BLANK)
header(s, 2, "Deliverable Map & Layer Hierarchy",
       "22-layer hierarchical naming \u00b7 toggleable overlays \u00b7 consistent across all sheets")
# left: layer list
layers = [
    "00_Background_TitleBlock", "01_Grid_Match-lines", "02_Scope_Boundaries",
    "03_Equipment_Vessels", "04_Equipment_HX-Cryo", "05_Piping_40K",
    "06_Piping_4K5", "07_Piping_2K", "08_Piping_Water", "09_Piping_Infra",
    "10_Piping_OutOfScope", "11_Valves_Actuators", "12_Instruments_Field",
    "13_Instruments_Panel", "14_Signal_Pneumatic", "15_Signal_Electric_SW",
    "16_Legend_TOGGLEABLE", "17_Notes_TOGGLEABLE",
]
_add_rect(s, Inches(0.45), Inches(1.25), Inches(5.0), Inches(5.55), LGREY)
_txt(s, Inches(0.65), Inches(1.35), Inches(4.7), Inches(0.4),
     "Layer stack (bottom \u2192 top)", size=14, color=NAVY, bold=True)
tb = s.shapes.add_textbox(Inches(0.65), Inches(1.85), Inches(4.7), Inches(4.8))
tf = tb.text_frame; tf.word_wrap = True
for i, ly in enumerate(layers):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    p.space_after = Pt(2)
    r = p.add_run(); r.text = ly
    r.font.size = Pt(11.5); r.font.name = "Consolas"
    tog = ly.endswith("TOGGLEABLE")
    r.font.color.rgb = ORANGE if tog else DARK
    r.font.bold = tog
# right: bullets
_bullets(s, Inches(5.8), Inches(1.35), Inches(7.0), Inches(5.4), [
    (0, "Naming convention: NN_Group_Detail", NAVY, True),
    (1, "two-digit ordinal controls Z-order / print order", GREY),
    (1, "Group = Background / Piping / Instruments / Signal / Legend", GREY),
    (0, "Piping split by cryogenic class for selective visibility", NAVY, True),
    (1, "40K, 4K5, 2K, Water, Infra, OutOfScope each on own layer", GREY),
    (0, "Signals separated from process piping", NAVY, True),
    (1, "Pneumatic vs Electric/Software on dedicated layers", GREY),
    (0, "TOGGLEABLE overlays default to display:none", ORANGE, True),
    (1, "Legend (16) and Notes (17) toggle on demand \u2014 zero clutter by default", GREY),
    (0, "Identical layer tree across all 4 sheets & 4 variants", NAVY, True),
    (1, "see LAYER_NAMING_STANDARD.md for the controlled list", GREY),
])
footer(s)

# =====================================================================
# SLIDES 3-6 -- the four sheets
# =====================================================================
sheet_defs = [
    (3, "QCELL \u2014 Sheet 1 \u00b7 Cryogenic Circuits", "40 K / 4.5 K / 2 K + HX  \u00b7  =NA.PS01_PFB712", "q1_std.png",
     [(0,"Three cryogenic temperature classes on one sheet",NAVY,True),
      (1,"40 K (red) \u00b7 4.5 K (blue) \u00b7 2 K (cyan)",GREY),
      (0,"RF cavities CAV.1/CAV.2 + couplers + tuners",NAVY,True),
      (0,"DI-water cooling (green) & infrastructure runs",NAVY,True),
      (0,"Buffer-volume & scope hand-over notes",ORANGE,True)]),
    (4, "QCELL \u2014 Sheet 2 \u00b7 Instrumentation", "Signals / loops / ISA bubbles  \u00b7  =NA.PS01_PFB712", "q2_std.png",
     [(0,"Process greyed; instrumentation in foreground",NAVY,True),
      (0,"ISA 5.1 bubbles: field / panel / shared",GREY),
      (0,"Three distinct 0.25 mm signal patterns",NAVY,True),
      (1,"pneumatic (dash+//), electric (dotted), hydraulic (dash-dot)",GREY),
      (0,"Re-allocated piezo sensors PZ535 / PZ525",ORANGE,True)]),
    (5, "RFCELL \u2014 Sheet 1 \u00b7 Process", "RF distribution & cooling  \u00b7  =NA.PS01_PFB713", "r1_std.png",
     [(0,"RF power & cooling distribution cell",NAVY,True),
      (0,"Same class colour & line-weight system",GREY),
      (0,"Scope diamonds at the last-metre hand-over",ORANGE,True),
      (0,"Shares symbol library with QCELL sheets",NAVY,True)]),
    (6, "RFCELL \u2014 Sheet 2 \u00b7 Instrumentation", "Signals / loops / ISA bubbles  \u00b7  =NA.PS01_PFB713", "r2_std.png",
     [(0,"Instrumentation foreground / process backdrop",NAVY,True),
      (0,"Consistent signal-pattern legend",GREY),
      (0,"Panel vs field instrument differentiation",NAVY,True),
      (0,"Cross-references to RFCELL Sheet 1",GREY)]),
]
for idx, title, sub, img, bl in sheet_defs:
    s = prs.slides.add_slide(BLANK)
    header(s, idx, title, sub)
    _img_fit(s, asset(img), Inches(0.45), Inches(1.25), Inches(8.7), Inches(5.55))
    _add_rect(s, Inches(9.35), Inches(1.25), Inches(3.55), Inches(5.55), LGREY)
    _txt(s, Inches(9.55), Inches(1.38), Inches(3.2), Inches(0.4),
         "Highlights", size=14, color=NAVY, bold=True)
    _bullets(s, Inches(9.55), Inches(1.9), Inches(3.2), Inches(4.7), bl, size=12.5, gap=7)
    footer(s)

# =====================================================================
# SLIDE 7 -- Piping hierarchy
# =====================================================================
s = prs.slides.add_slide(BLANK)
header(s, 7, "Piping Hierarchy & Line Weights",
       "Visual weight communicates importance \u2014 readable at A3")
_img_fit(s, asset("crop_piping.png"), Inches(0.45), Inches(1.3), Inches(7.2), Inches(4.2))
_txt(s, Inches(0.45), Inches(5.6), Inches(7.2), Inches(0.4),
     "QCELL Sheet 1 \u2014 process-piping detail", size=11, color=GREY, italic=True)
# weight table
rows = [
    ("PRIMARY", "1.0 mm", "main headers per cryo class", BLUE),
    ("BRANCH", "0.7 mm", "take-offs / sub-runs per class", RED),
    ("SECONDARY", "0.5 mm", "DI-water & utility", GREEN),
    ("OUT-OF-SCOPE", "0.35 mm dashed", "services beyond hand-over", GREY),
    ("SIGNAL", "0.25 mm", "instrument signal lines", RGBColor(0x7a,0x00,0xa0)),
]
tx = Inches(8.0); ty = Inches(1.5)
_txt(s, tx, Inches(1.05), Inches(4.9), Inches(0.4),
     "Line-weight system", size=15, color=NAVY, bold=True)
for i, (nm, wt, desc, col) in enumerate(rows):
    ry = ty + Inches(1.02) * i
    _add_rect(s, tx, ry, Inches(0.85), Inches(0.12 + i*0.0), col)  # swatch bar
    sp = _add_rect(s, tx, ry + Inches(0.18), Inches(4.9), Inches(0.78), WHITE, line=LGREY)
    _txt(s, tx + Inches(0.1), ry + Inches(0.20), Inches(2.0), Inches(0.35),
         nm, size=13, color=col, bold=True)
    _txt(s, tx + Inches(2.1), ry + Inches(0.20), Inches(2.7), Inches(0.35),
         wt, size=12, color=DARK, bold=True, align=PP_ALIGN.RIGHT)
    _txt(s, tx + Inches(0.1), ry + Inches(0.54), Inches(4.7), Inches(0.32),
         desc, size=11, color=GREY)
footer(s)

# =====================================================================
# SLIDE 8 -- Instrumentation & ISA bubbles
# =====================================================================
s = prs.slides.add_slide(BLANK)
header(s, 8, "Instrumentation \u2014 ISA 5.1 Bubbles & Signals",
       "Location modifiers \u00b7 safety outlines \u00b7 three signal patterns")
_img_fit(s, asset("crop_bubbles.png"), Inches(0.45), Inches(1.3), Inches(6.0), Inches(5.3))
_bullets(s, Inches(6.8), Inches(1.35), Inches(6.1), Inches(5.4), [
    (0,"Bubble = ISA 5.1 instrument tag",NAVY,True),
    (1,"plain circle = field-mounted",GREY),
    (1,"single line = panel / main control",GREY),
    (1,"double line = shared display / DCS",GREY),
    (1,"dashed outline = safety / interlock function",GREY),
    (0,"Tag text \u2265 2.0 mm; main tags \u2265 2.5 mm @ A3",ORANGE,True),
    (0,"Signal lines kept at 0.25 mm but visually distinct:",NAVY,True),
    (1,"Pneumatic \u2014 dashed with // cross-ticks",RGBColor(0x7a,0x00,0xa0)),
    (1,"Electric \u2014 fine dotted",RGBColor(0x00,0x52,0x9b)),
    (1,"Hydraulic \u2014 dash-dot",RGBColor(0xa0,0x6a,0x00)),
    (0,"Signals on dedicated layers, separate from process piping",NAVY,True),
])
footer(s)

# =====================================================================
# SLIDE 9 -- Colour vs Monochrome
# =====================================================================
s = prs.slides.add_slide(BLANK)
header(s, 9, "Colour vs Monochrome Variants",
       "Every sheet ships full-colour + mono \u2014 mono survives B/W plotting & photocopy")
_img_fit(s, asset("crop_piping.png"), Inches(0.45), Inches(1.45), Inches(6.0), Inches(3.5))
_txt(s, Inches(0.45), Inches(5.0), Inches(6.0), Inches(0.35),
     "STANDARD (full colour)", size=12, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
_img_fit(s, asset("crop_mono_piping.png"), Inches(6.85), Inches(1.45), Inches(6.0), Inches(3.5))
_txt(s, Inches(6.85), Inches(5.0), Inches(6.0), Inches(0.35),
     "STANDARD_MONO (black / grey)", size=12, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
_bullets(s, Inches(0.6), Inches(5.5), Inches(12.2), Inches(1.4), [
    (0,"Mono removes class colour \u2014 hierarchy now carried by line weight + dash pattern alone",GREY),
    (0,"Marker arrowheads / junction dots recoloured to black; cavity fills switched to white for contrast",GREY),
    (0,"Guarantees legibility when printed on a mono A3 plotter or photocopied",GREY),
])
footer(s)

# =====================================================================
# SLIDE 10 -- Symbol library (legend)
# =====================================================================
s = prs.slides.add_slide(BLANK)
header(s, 10, "Symbol Library & Toggleable Legend",
       "On-sheet legend (layer 16) \u2014 hidden by default, one click to reveal")
_img_fit(s, asset("crop_legend.png"), Inches(0.6), Inches(1.3), Inches(4.6), Inches(5.4))
_bullets(s, Inches(5.6), Inches(1.4), Inches(7.3), Inches(5.4), [
    (0,"Legend embedded on the sheet, not a wasteful side panel",NAVY,True),
    (1,"lives on layer 16_Legend_TOGGLEABLE (display:none default)",GREY),
    (0,"Documents the full visual grammar:",NAVY,True),
    (1,"piping hierarchy (primary / branch / secondary / out-of-scope)",GREY),
    (1,"three 0.25 mm signal patterns",GREY),
    (1,"ISA 5.1 instrument bubble families",GREY),
    (1,"3-compartment scope diamond (TP / cat+num / next system)",GREY),
    (0,"AD_01.16 scope categories spelled out:",ORANGE,True),
    (1,"B=Building C=Civil E=Electrical G=Compressed gasses",GREY),
    (1,"H=HVAC L=Liquid waste S=Solid waste W=Water",GREY),
    (0,"Maximises usable drawing area \u2014 a core v3 goal",NAVY,True),
])
footer(s)

# =====================================================================
# SLIDE 11 -- Sensor re-allocations
# =====================================================================
s = prs.slides.add_slide(BLANK)
header(s, 11, "Sensor Re-allocations & New Tags",
       "Piezo / magnetic-coupler sensor rationalisation captured in the nomenclature master")
_add_rect(s, Inches(0.45), Inches(1.3), Inches(12.45), Inches(2.1), LGREY)
_txt(s, Inches(0.65), Inches(1.4), Inches(12), Inches(0.4),
     "Re-allocations", size=15, color=NAVY, bold=True)
realloc = [
    ("TT535", "PZ535", "coldest piezo (TT-CX) re-tagged as pressure/piezo element"),
    ("TT525", "PZ525", "warmest piezo (TT-PT100) re-tagged (applied if present in source)"),
    ("4 \u00d7 TT", "CX / PT100", "magnetic-coupler sensors redistributed across cryo stages"),
]
for i,(a,b,desc) in enumerate(realloc):
    ry = Inches(1.95) + Inches(0.46)*i
    _txt(s, Inches(0.7), ry, Inches(1.6), Inches(0.4), a, size=14, color=RED, bold=True, anchor=MSO_ANCHOR.MIDDLE)
    _txt(s, Inches(2.3), ry, Inches(0.6), Inches(0.4), "\u2192", size=14, color=GREY, bold=True, anchor=MSO_ANCHOR.MIDDLE)
    _txt(s, Inches(2.9), ry, Inches(1.7), Inches(0.4), b, size=14, color=GREEN, bold=True, anchor=MSO_ANCHOR.MIDDLE)
    _txt(s, Inches(4.7), ry, Inches(8.0), Inches(0.4), desc, size=12, color=DARK, anchor=MSO_ANCHOR.MIDDLE)
_bullets(s, Inches(0.6), Inches(3.7), Inches(12.3), Inches(3.2), [
    (0,"Full instrument & equipment register lives in MINERVA_NOMENCLATURE_MASTER.xlsx",NAVY,True),
    (1,"01_Instruments (297 rows) \u00b7 02_Equipment (44) \u00b7 03_Re-allocations \u00b7 04_New_components",GREY),
    (1,"05_Layer_map \u00b7 06_Colour_LineWeight \u00b7 07_Scope_categories",GREY),
    (0,"Every tag traceable from drawing \u2192 register \u2192 layer \u2192 colour/weight",NAVY,True),
    (0,"Re-allocations flagged so reviewers can audit the change set quickly",ORANGE,True),
])
footer(s)

# =====================================================================
# SLIDE 12 -- Scope boundaries / diamonds
# =====================================================================
s = prs.slides.add_slide(BLANK)
header(s, 12, "Scope Boundaries \u2014 AD_01.16 Diamonds",
       "3-compartment hand-over markers at the last-metre boundary")
_img_fit(s, asset("crop_diamonds.png"), Inches(0.45), Inches(1.35), Inches(6.4), Inches(5.0))
_bullets(s, Inches(7.1), Inches(1.4), Inches(5.8), Inches(5.4), [
    (0,"Diamond = scope hand-over point (TP)",NAVY,True),
    (0,"Three compartments:",NAVY,True),
    (1,"top \u2014 TP marker",GREY),
    (1,"middle \u2014 category letter + number",GREY),
    (1,"bottom \u2014 next system / package",GREY),
    (0,"AD_01.16 categories:",ORANGE,True),
    (1,"B Building \u00b7 C Civil \u00b7 E Electrical",GREY),
    (1,"G Compressed gasses \u00b7 H HVAC",GREY),
    (1,"L Liquid waste \u00b7 S Solid waste \u00b7 W Water",GREY),
    (0,"Out-of-scope piping drawn 0.35 mm dashed for instant recognition",NAVY,True),
    (0,"Boundaries isolated on layer 02_Scope_Boundaries",GREY),
])
footer(s)

# =====================================================================
# SLIDE 13 -- Default views
# =====================================================================
s = prs.slides.add_slide(BLANK)
header(s, 13, "Default Views \u2014 Embedded Presets",
       "5 named view presets stored as SVG <metadata> for the viewer")
views = [
    ("01 Overview", "fit whole sheet \u00b7 all process layers \u00b7 legend off"),
    ("02 Cryogenic", "isolate 40K/4K5/2K piping + cryo equipment"),
    ("03 Instrumentation", "instruments + signals foreground, process greyed"),
    ("04 Scope review", "scope boundaries + diamonds + out-of-scope piping"),
    ("05 Print (mono)", "monochrome-ready, legend + notes on, A3 frame"),
]
_add_rect(s, Inches(0.45), Inches(1.3), Inches(12.45), Inches(0.55), NAVY)
_txt(s, Inches(0.6), Inches(1.34), Inches(4), Inches(0.45), "Preset", size=13, color=WHITE, bold=True, anchor=MSO_ANCHOR.MIDDLE)
_txt(s, Inches(4.6), Inches(1.34), Inches(8), Inches(0.45), "What it shows", size=13, color=WHITE, bold=True, anchor=MSO_ANCHOR.MIDDLE)
for i,(nm,desc) in enumerate(views):
    ry = Inches(1.95) + Inches(0.72)*i
    bg = WHITE if i%2==0 else LGREY
    _add_rect(s, Inches(0.45), ry, Inches(12.45), Inches(0.66), bg, line=RGBColor(0xDD,0xDD,0xDD))
    _txt(s, Inches(0.6), ry, Inches(4.0), Inches(0.66), nm, size=14, color=NAVY, bold=True, anchor=MSO_ANCHOR.MIDDLE)
    _txt(s, Inches(4.6), ry, Inches(8.1), Inches(0.66), desc, size=12.5, color=DARK, anchor=MSO_ANCHOR.MIDDLE)
_txt(s, Inches(0.5), Inches(6.55), Inches(12.3), Inches(0.5),
     "Presets are machine-readable metadata (layer visibility + zoom box) \u2014 see DEFAULT_VIEWS_GUIDE.md.",
     size=12, color=GREY, italic=True)
footer(s)

# =====================================================================
# SLIDE 14 -- A3 printing / closing
# =====================================================================
s = prs.slides.add_slide(BLANK)
_add_rect(s, 0, 0, SW, SH, NAVY)
_add_rect(s, 0, Inches(1.0), SW, Emu(38100), CYAN)
_txt(s, Inches(0.6), Inches(0.25), Inches(12), Inches(0.7),
     "A3 Printing & Deliverable Summary", size=28, color=WHITE, bold=True)
_txt(s, Inches(12.4), Inches(0.2), Inches(0.8), Inches(0.6), "14", size=22, color=CYAN, bold=True, align=PP_ALIGN.RIGHT)
_bullets(s, Inches(0.7), Inches(1.4), Inches(7.0), Inches(5.3), [
    (0,"Drawing frame = A3 landscape 420 \u00d7 297 mm",WHITE,True),
    (1,"compact full-width bottom title block (~30 mm)",RGBColor(0xCF,0xDD,0xF2)),
    (1,"no wasteful right-hand panel \u2014 max drawing area",RGBColor(0xCF,0xDD,0xF2)),
    (0,"Text sizing tuned for true A3 legibility",WHITE,True),
    (1,"main tags \u2265 2.5 mm \u00b7 bubbles 2.0 mm \u00b7 callouts 2.2 mm \u00b7 legend 1.8 mm",RGBColor(0xCF,0xDD,0xF2)),
    (0,"16 PDF exported at exact A3 (1190.55 \u00d7 841.89 pt)",WHITE,True),
    (0,"Print mono variants for B/W plotters & copies",WHITE,True),
    (0,"See PRINTING_GUIDE_A3.pdf for plotter setup",CYAN,True),
])
# deliverables card
_add_rect(s, Inches(8.0), Inches(1.4), Inches(4.8), Inches(5.3), RGBColor(0x12,0x2C,0x55))
_txt(s, Inches(8.25), Inches(1.55), Inches(4.4), Inches(0.4), "Deliverables", size=16, color=CYAN, bold=True)
_bullets(s, Inches(8.25), Inches(2.1), Inches(4.4), Inches(4.5), [
    (0,"16 SVG (4 sheets \u00d7 4 variants)",WHITE),
    (0,"16 PDF (A3, print-ready)",WHITE),
    (0,"MINERVA_NOMENCLATURE_MASTER.xlsx",WHITE),
    (0,"MINERVA_PID_DISSECTION.pptx",WHITE),
    (0,"LAYER_NAMING_STANDARD.md/.pdf",WHITE),
    (0,"DEFAULT_VIEWS_GUIDE.md/.pdf",WHITE),
    (0,"PRINTING_GUIDE_A3.pdf",WHITE),
], size=13, gap=9)
_txt(s, Inches(0.7), Inches(7.05), Inches(12), Inches(0.35),
     "MINERVA CryoCell \u00b7 SCK CEN / Mott MacDonald \u00b7 AD_01.16 / ISO 10628 / ISA 5.1 \u00b7 Rev C1 2026-06",
     size=10, color=RGBColor(0x9F,0xB4,0xD4))

prs.save(OUT)
print("WROTE", OUT, "slides:", len(prs.slides._sldIdLst))
