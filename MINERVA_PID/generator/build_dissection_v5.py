#!/usr/bin/env python3
"""
build_dissection_v5.py
Builds MINERVA_PID_DISSECTION_v5.pptx -- the v5 visual dissection deck for the
MINERVA CryoCell P&ID (SCK CEN / Mott MacDonald).

v5 focus (vs the v3/v4 dissection deck):
  * Slide 3  "Line Classification & Colour Coding" -- REVISED to the new
             cryogenic-focus palette (cold header blue/cyan on top, thermal
             shield orange/red, warm WPS green/lime/olive on the bottom) with
             temperatures, pressures and flow rates.
  * Slide 3.5 "Line Nomenclature & Branching" -- NEW slide explaining the
             [LINE]-[SIZE]-[MOC] tag scheme and the MAIN -> BRANCH (A -> A')
             hierarchy with size reduction & material of construction.

Pure tooling: python-pptx + PIL.  Image assets live in ../deck_assets_v5/.
"""
import os

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from PIL import Image

import line_spec_data as LSD

ROOT = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
ASSETS = os.path.join(ROOT, "deck_assets_v5")
OUT = os.path.join(ROOT, "output_v5", "MINERVA_PID_DISSECTION_v5.pptx")

# ---- palette (deck chrome) ----
NAVY   = RGBColor(0x0B, 0x21, 0x47)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
DARK   = RGBColor(0x1A, 0x1A, 0x1A)
GREY   = RGBColor(0x55, 0x55, 0x55)
LGREY  = RGBColor(0xE8, 0xEC, 0xF1)
CYANUI = RGBColor(0x00, 0xA6, 0xBD)
ORANGE = RGBColor(0xC0, 0x60, 0x00)

# ---- line colours (must match the P&ID v5 palette) ----
C_A  = RGBColor(0x00, 0x00, 0xFF)
C_Ap = RGBColor(0x00, 0x00, 0x80)
C_B  = RGBColor(0x00, 0xC0, 0xD0)   # cyan, darkened a touch for white text legibility
C_Bp = RGBColor(0x00, 0x8B, 0x8B)
C_D  = RGBColor(0xFF, 0x80, 0x00)
C_Dp = RGBColor(0xFF, 0xB3, 0x66)
C_E  = RGBColor(0xFF, 0x00, 0x00)
C_Ep = RGBColor(0xCC, 0x00, 0x00)
C_W  = RGBColor(0x00, 0xC0, 0x00)   # green, darkened for white text
C_S  = RGBColor(0x8F, 0xBF, 0x00)   # lime, darkened
C_U  = RGBColor(0x80, 0x80, 0x00)
C_OS = RGBColor(0x80, 0x80, 0x80)


def hx(s):
    s = s.lstrip("#")
    return RGBColor(int(s[0:2], 16), int(s[2:4], 16), int(s[4:6], 16))


prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]
SW, SH = prs.slide_width, prs.slide_height


def _rect(slide, x, y, w, h, color, line=None, line_w=1.0):
    sp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    sp.fill.solid(); sp.fill.fore_color.rgb = color
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line; sp.line.width = Pt(line_w)
    sp.shadow.inherit = False
    return sp


def _txt(slide, x, y, w, h, text, size=18, color=DARK, bold=False,
         align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, italic=False, font="Calibri"):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame; tf.word_wrap = True
    tf.vertical_anchor = anchor
    for i, ln in enumerate(text.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        r = p.add_run(); r.text = ln
        r.font.size = Pt(size); r.font.bold = bold; r.font.italic = italic
        r.font.color.rgb = color; r.font.name = font
    return tb


def _bullets(slide, x, y, w, h, items, size=15, color=DARK, gap=5):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame; tf.word_wrap = True
    for i, (lvl, txt, *opt) in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.level = lvl
        p.space_after = Pt(gap)
        bullet = "\u2022 " if lvl == 0 else "\u2013 "
        r = p.add_run(); r.text = bullet + txt
        r.font.size = Pt(size - lvl)
        r.font.color.rgb = opt[0] if opt else color
        r.font.name = "Calibri"
        if opt and len(opt) > 1 and opt[1]:
            r.font.bold = True
    return tb


def _img_fit(slide, path, x, y, w, h, frame=True, frame_color=GREY):
    im = Image.open(path); iw, ih = im.size
    ar = iw / ih; box_ar = w / h
    if ar > box_ar:
        nw = w; nh = int(w / ar)
    else:
        nh = h; nw = int(h * ar)
    nx = x + (w - nw) // 2; ny = y + (h - nh) // 2
    if frame:
        _rect(slide, nx - Emu(9144), ny - Emu(9144),
              nw + Emu(18288), nh + Emu(18288), WHITE, line=frame_color)
    slide.shapes.add_picture(path, nx, ny, nw, nh)
    return (nx, ny, nw, nh)


def header(slide, idx, title, subtitle=None):
    _rect(slide, 0, 0, SW, Inches(0.95), NAVY)
    _rect(slide, 0, Inches(0.95), SW, Emu(38100), CYANUI)
    _txt(slide, Inches(0.45), Inches(0.10), Inches(11.4), Inches(0.55),
         title, size=25, color=WHITE, bold=True, anchor=MSO_ANCHOR.MIDDLE)
    if subtitle:
        _txt(slide, Inches(0.47), Inches(0.62), Inches(11.4), Inches(0.30),
             subtitle, size=12, color=RGBColor(0xBF, 0xD3, 0xF2),
             anchor=MSO_ANCHOR.MIDDLE)
    _txt(slide, Inches(12.3), Inches(0.10), Inches(0.9), Inches(0.6),
         idx, size=20, color=CYANUI, bold=True,
         align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)


def footer(slide):
    _txt(slide, Inches(0.45), Inches(7.13), Inches(8.5), Inches(0.3),
         "MINERVA CryoCell  \u00b7  P&ID v5 Dissection  \u00b7  SCK CEN / Mott MacDonald",
         size=9, color=GREY)
    _txt(slide, Inches(9.0), Inches(7.13), Inches(3.9), Inches(0.3),
         "AD_01.16  \u00b7  ISO 10628  \u00b7  ISA 5.1-2022", size=9, color=GREY,
         align=PP_ALIGN.RIGHT)


def asset(n):
    return os.path.join(ASSETS, n)


# =====================================================================
# SLIDE 1 -- Title
# =====================================================================
s = prs.slides.add_slide(BLANK)
_rect(s, 0, 0, SW, SH, NAVY)
_rect(s, 0, Inches(2.55), SW, Inches(0.06), CYANUI)
_txt(s, Inches(0.8), Inches(1.15), Inches(11.7), Inches(1.1),
     "MINERVA CryoCell \u2014 P&ID", size=42, color=WHITE, bold=True)
_txt(s, Inches(0.8), Inches(2.7), Inches(11.7), Inches(0.8),
     "Revision v5  \u00b7  Cryogenic-Focus Colour Scheme & Line Nomenclature",
     size=22, color=CYANUI, bold=True)
_bullets(s, Inches(0.85), Inches(3.7), Inches(11.6), Inches(2.6), [
    (0, "New cryogenic-focus palette: cold header (A/B) on top, thermal "
        "shield (D/E), warm WPS (W/S/U) on the bottom", WHITE),
    (0, "[LINE]-[SIZE]-[MOC] nomenclature labels + flow-direction arrows on "
        "every classified run", WHITE),
    (0, "Line specification table, MAIN-LINES-ONLY schematic & Line W "
        "temperature-gradient annotation", WHITE),
    (0, "16 production sheets (colour + mono) + 2 single-line schematics + "
        "LINE_SPECIFICATION_MASTER.xlsx", WHITE),
], size=15, gap=10)
_txt(s, Inches(0.8), Inches(6.7), Inches(11.7), Inches(0.5),
     "SCK CEN  \u00b7  MYRRHA / MINERVA Phase 1  \u00b7  Mott MacDonald (Bristol, UK)  "
     "\u00b7  MMD 411066  \u00b7  S2 - FOR ACCEPTANCE  \u00b7  RESTRICTED",
     size=12, color=RGBColor(0xBF, 0xD3, 0xF2))

# =====================================================================
# SLIDE 2 -- What changed in v5
# =====================================================================
s = prs.slides.add_slide(BLANK)
header(s, "02", "What Changed in v5",
       "Ten-phase revision of the v4 deliverable \u2014 cryogenic clarity first")
_bullets(s, Inches(0.55), Inches(1.35), Inches(6.1), Inches(5.6), [
    (0, "Cryogenic-focus colour scheme", NAVY, True),
    (1, "A 4.5 K blue \u00b7 B 2 K cyan \u00b7 D 40 K orange \u00b7 E 60 K red", GREY),
    (1, "warm WPS: W green \u00b7 S lime \u00b7 U olive", GREY),
    (0, "Spatial reorganisation (zone bands)", NAVY, True),
    (1, "cold header top \u2192 thermal shield \u2192 equipment \u2192 warm bottom", GREY),
    (0, "[LINE]-[SIZE]-[MOC] line labels (layer 04D)", NAVY, True),
    (0, "Flow-direction arrows on every run (layer 04G)", NAVY, True),
    (0, "Line specification table embedded in the legend", NAVY, True),
    (0, "MAIN-LINES-ONLY preset + standalone schematic", ORANGE, True),
    (0, "Mono = line-weight + dash class differentiation", NAVY, True),
    (0, "Line W temperature gradient (4.5 K \u2192 300 K)", ORANGE, True),
], size=14, gap=6)
_img_fit(s, asset("ml_q.png"), Inches(6.9), Inches(1.35), Inches(6.0), Inches(5.5))
_txt(s, Inches(6.9), Inches(6.75), Inches(6.0), Inches(0.3),
     "QCELL \u2014 MAIN-LINES-ONLY single-line schematic (new in v5)",
     size=10, color=GREY, italic=True, align=PP_ALIGN.CENTER)
footer(s)

# =====================================================================
# SLIDE 3 -- Line Classification & Colour Coding  (REVISED)
# =====================================================================
s = prs.slides.add_slide(BLANK)
header(s, "03", "Line Classification & Colour Coding  \u2014  revised v5 palette",
       "Cold header (top) \u00b7 thermal shield \u00b7 warm piping system (bottom)")

# three zone groups, each a coloured band of line rows
groups = [
    ("COLD HEADER", hx("#0000FF"), hx("#eaf0ff"), [
        ("A",  C_A,  "4.5 K primary He",  "4.5 K",  "3 bar",   "~50 g/s",  "DN50", "SS316L"),
        ("A'", C_Ap, "4.5 K branches",    "4.5 K",  "3 bar",   "varies",   "DN25", "SS316L"),
        ("B",  C_B,  "2 K primary He",    "2 K",    "27 mbar", "~47.5 g/s","DN40", "SS316L"),
        ("B'", C_Bp, "2 K branches",      "2 K",    "27 mbar", "varies",   "DN25", "SS316L"),
    ]),
    ("THERMAL SHIELD", hx("#FF8000"), hx("#fff1e3"), [
        ("D",  C_D,  "40 K shield inlet",  "40 K", "14 bar", "TBD",    "DN32", "Cu"),
        ("D'", C_Dp, "40 K branches",      "40 K", "14 bar", "varies", "DN20", "Cu"),
        ("E",  C_E,  "60 K shield outlet", "60 K", "13 bar", "TBD",    "DN32", "Cu"),
        ("E'", C_Ep, "60 K branches",      "60 K", "13 bar", "varies", "DN20", "Cu"),
    ]),
    ("WARM PIPING SYSTEM (WPS)", hx("#00a000"), hx("#ecf8e8"), [
        ("W", C_W, "WPS warm return",        "4.5 K-300 K", "6 bar",    "~2.5 g/s", "DN20", "SS304"),
        ("S", C_S, "WPS service / safety",   "2-292 K",     "1.05 bar", "TBD",      "DN15", "SS304"),
        ("U", C_U, "WPS GHe supply inlet",   "292 K",       "14 bar",   "TBD",      "DN15", "SS304"),
    ]),
]
# column header
colx = [Inches(0.55), Inches(1.55), Inches(4.7), Inches(6.5), Inches(8.0),
        Inches(9.5), Inches(11.0), Inches(12.0)]
hdr = ["LINE", "Description", "Temp", "Press", "Flow", "DN", "MOC", ""]
_txt(s, colx[0], Inches(1.18), Inches(1.0), Inches(0.3), "LINE", size=11,
     color=NAVY, bold=True)
_txt(s, colx[1], Inches(1.18), Inches(3.0), Inches(0.3), "Description", size=11,
     color=NAVY, bold=True)
_txt(s, colx[2], Inches(1.18), Inches(1.6), Inches(0.3), "Temp", size=11,
     color=NAVY, bold=True)
_txt(s, colx[3], Inches(1.18), Inches(1.4), Inches(0.3), "Pressure", size=11,
     color=NAVY, bold=True)
_txt(s, colx[4], Inches(1.18), Inches(1.4), Inches(0.3), "Flow", size=11,
     color=NAVY, bold=True)
_txt(s, colx[5], Inches(1.18), Inches(1.4), Inches(0.3), "DN", size=11,
     color=NAVY, bold=True)
_txt(s, colx[6], Inches(1.18), Inches(1.4), Inches(0.3), "MOC", size=11,
     color=NAVY, bold=True)

y = Inches(1.55)
rh = Inches(0.34)
for gname, gedge, gfill, rows in groups:
    _rect(s, Inches(0.45), y, Inches(12.45), Inches(0.30), gedge)
    _txt(s, Inches(0.55), y + Emu(10000), Inches(8.0), Inches(0.28), gname,
         size=12, color=WHITE, bold=True, anchor=MSO_ANCHOR.MIDDLE)
    y = y + Inches(0.34)
    for i, (lid, col, desc, temp, press, flow, dn, moc) in enumerate(rows):
        if i % 2 == 0:
            _rect(s, Inches(0.45), y, Inches(12.45), rh, gfill)
        # swatch line
        sw = _rect(s, Inches(0.55), y + Inches(0.14), Inches(0.85), Emu(40000), col)
        _txt(s, colx[0], y + Emu(8000), Inches(0.95), rh, lid, size=12,
             color=col, bold=True, anchor=MSO_ANCHOR.MIDDLE)
        _txt(s, colx[1], y, Inches(3.0), rh, desc, size=11.5, color=DARK,
             anchor=MSO_ANCHOR.MIDDLE)
        _txt(s, colx[2], y, Inches(1.6), rh, temp, size=11.5, color=DARK,
             anchor=MSO_ANCHOR.MIDDLE)
        _txt(s, colx[3], y, Inches(1.4), rh, press, size=11.5, color=DARK,
             anchor=MSO_ANCHOR.MIDDLE)
        _txt(s, colx[4], y, Inches(1.4), rh, flow, size=11.5, color=DARK,
             anchor=MSO_ANCHOR.MIDDLE)
        _txt(s, colx[5], y, Inches(1.4), rh, dn, size=11.5, color=DARK,
             anchor=MSO_ANCHOR.MIDDLE)
        _txt(s, colx[6], y, Inches(1.4), rh, moc, size=11.5, color=DARK,
             anchor=MSO_ANCHOR.MIDDLE)
        y = y + rh
    y = y + Inches(0.06)
_txt(s, Inches(0.55), y + Inches(0.02), Inches(12.0), Inches(0.3),
     "Temperature increases top \u2192 bottom (4.5 K \u2192 300 K).  Primary lines are "
     "saturated; branches (\u2032) use the darker / lighter shade of the same hue.",
     size=10.5, color=GREY, italic=True)
footer(s)

# =====================================================================
# SLIDE 3.5 -- Line Nomenclature & Branching  (NEW)
# =====================================================================
s = prs.slides.add_slide(BLANK)
header(s, "3.5", "Line Nomenclature & Branching  \u2014  new in v5",
       "[LINE]-[SIZE]-[MOC] tags  \u00b7  MAIN \u2192 BRANCH (A \u2192 A\u2032) hierarchy")

# left: tag anatomy
_rect(s, Inches(0.45), Inches(1.3), Inches(6.0), Inches(2.55), LGREY)
_txt(s, Inches(0.65), Inches(1.42), Inches(5.6), Inches(0.4),
     "Tag anatomy", size=15, color=NAVY, bold=True)
_txt(s, Inches(0.65), Inches(1.95), Inches(5.7), Inches(0.7),
     "A - DN50 - SS316L", size=30, color=C_A, bold=True, font="Consolas")
_bullets(s, Inches(0.65), Inches(2.7), Inches(5.7), Inches(1.1), [
    (0, "[LINE]   class designator (A, B, D, E, W, S, U)", GREY),
    (0, "[SIZE]   nominal bore (DN50, DN40, DN32 \u2026)", GREY),
    (0, "[MOC]    material of construction (SS316L, Cu, SS304)", GREY),
], size=12, gap=4)

# left-bottom: branching hierarchy
_rect(s, Inches(0.45), Inches(4.05), Inches(6.0), Inches(2.85), LGREY)
_txt(s, Inches(0.65), Inches(4.17), Inches(5.6), Inches(0.4),
     "Main \u2192 branch hierarchy", size=15, color=NAVY, bold=True)
pairs = [
    ("A  DN50", C_A, "A\u2032 DN25", C_Ap, "4.5 K take-offs"),
    ("B  DN40", C_B, "B\u2032 DN25", C_Bp, "2 K take-offs"),
    ("D  DN32", C_D, "D\u2032 DN20", C_Dp, "40 K shield branches"),
    ("E  DN32", C_E, "E\u2032 DN20", C_Ep, "60 K shield branches"),
]
yy = Inches(4.7)
for main, mc, br, bc, note in pairs:
    _txt(s, Inches(0.65), yy, Inches(1.7), Inches(0.32), main, size=13,
         color=mc, bold=True, font="Consolas")
    _txt(s, Inches(2.35), yy, Inches(0.6), Inches(0.32), "\u2192", size=13,
         color=GREY, bold=True)
    _txt(s, Inches(2.95), yy, Inches(1.6), Inches(0.32), br, size=13,
         color=bc, bold=True, font="Consolas")
    _txt(s, Inches(4.6), yy, Inches(1.8), Inches(0.32), note, size=10.5,
         color=GREY, anchor=MSO_ANCHOR.MIDDLE)
    yy = yy + Inches(0.5)

# right: bullets explaining rules + image
_bullets(s, Inches(6.7), Inches(1.35), Inches(6.2), Inches(3.0), [
    (0, "Primary trunk vs branch decided by run length", NAVY, True),
    (1, "55th-percentile threshold per class splits primary / branch", GREY),
    (1, "primary on layers 04/05/06; branches on 04B/05B/06B", GREY),
    (0, "Branch = same hue, reduced bore, same fluid", NAVY, True),
    (1, "A\u2032/B\u2032 darker shade; D\u2032/E\u2032 lighter shade", GREY),
    (0, "Size reduces along the distribution chain", NAVY, True),
    (1, "DN50/DN40 headers \u2192 DN25/DN20 branches \u2192 instrument tube", GREY),
    (0, "Labels render on layer 04D, offset along the pipe normal", ORANGE, True),
], size=12.5, gap=5)
_img_fit(s, asset("ml_q.png"), Inches(6.7), Inches(4.4), Inches(6.2), Inches(2.55))
_txt(s, Inches(6.7), Inches(6.95), Inches(6.2), Inches(0.3),
     "Branch drops A\u2032/B\u2032 into the cryomodule (schematic detail)",
     size=10, color=GREY, italic=True, align=PP_ALIGN.CENTER)
footer(s)

# =====================================================================
# SLIDE 4 -- Spatial reorganisation (zone bands) on the live sheet
# =====================================================================
s = prs.slides.add_slide(BLANK)
header(s, "04", "Spatial Reorganisation \u2014 Cold Top / Warm Bottom",
       "Zone bands group the sheet by temperature class (layer 02C)")
_img_fit(s, asset("q1_std.png"), Inches(0.45), Inches(1.25), Inches(8.8), Inches(5.6))
_add = _rect(s, Inches(9.45), Inches(1.25), Inches(3.45), Inches(5.6), LGREY)
_txt(s, Inches(9.6), Inches(1.38), Inches(3.2), Inches(0.4), "Reading order",
     size=14, color=NAVY, bold=True)
_bullets(s, Inches(9.6), Inches(1.9), Inches(3.2), Inches(4.8), [
    (0, "COLD HEADER (top)", C_A, True),
    (1, "A 4.5 K / B 2 K", GREY),
    (0, "THERMAL SHIELD", ORANGE, True),
    (1, "D 40 K in / E 60 K out", GREY),
    (0, "EQUIPMENT", GREY, True),
    (1, "cavities / vessels / HX", GREY),
    (0, "WARM PIPING (bottom)", C_W, True),
    (1, "W / S / U \u2192 QRB handover", GREY),
    (0, "Bands are reference overlays \u2014 extracted geometry is preserved, "
        "not re-routed.", DARK),
], size=12, gap=6)
footer(s)

# =====================================================================
# SLIDE 5 -- Specification table + MAIN-LINES-ONLY
# =====================================================================
s = prs.slides.add_slide(BLANK)
header(s, "05", "Line Specification Table & MAIN-LINES-ONLY View",
       "Embedded legend table (layer 16) + VIEW_MAINLINES_ONLY preset")
_img_fit(s, asset("crop_spectable.png"), Inches(0.5), Inches(1.3), Inches(5.2),
         Inches(5.5))
_txt(s, Inches(0.5), Inches(6.85), Inches(5.2), Inches(0.3),
     "Specification table embedded in the on-sheet legend",
     size=10, color=GREY, italic=True, align=PP_ALIGN.CENTER)
_img_fit(s, asset("ml_q.png"), Inches(6.0), Inches(1.4), Inches(6.9), Inches(4.0))
_txt(s, Inches(6.0), Inches(5.45), Inches(6.9), Inches(0.3),
     "Standalone QCELL-MAINLINES_VIEW.svg single-line schematic",
     size=10, color=GREY, italic=True, align=PP_ALIGN.CENTER)
_bullets(s, Inches(6.1), Inches(5.8), Inches(6.8), Inches(1.2), [
    (0, "Table lists every line: colour, temp, pressure, DN, MOC", GREY),
    (0, "MAIN-LINES preset hides branches & instrumentation for a clean "
        "distribution overview", GREY),
], size=11.5, gap=4)
footer(s)

# =====================================================================
# SLIDE 6 -- Colour vs Mono class differentiation
# =====================================================================
s = prs.slides.add_slide(BLANK)
header(s, "06", "Colour vs Monochrome \u2014 Class by Weight & Dash",
       "Mono survives B/W plotting: cold = heavy solid, warm = dashed")
_img_fit(s, asset("crop_spectable.png"), Inches(0.5), Inches(1.4), Inches(5.9),
         Inches(5.0))
_txt(s, Inches(0.5), Inches(6.5), Inches(5.9), Inches(0.3),
     "STANDARD (full colour)", size=12, color=NAVY, bold=True,
     align=PP_ALIGN.CENTER)
_img_fit(s, asset("crop_spectable_mono.png"), Inches(6.9), Inches(1.4),
         Inches(5.9), Inches(5.0))
_txt(s, Inches(6.9), Inches(6.5), Inches(5.9), Inches(0.3),
     "STANDARD_MONO (line-weight + dash)", size=12, color=NAVY, bold=True,
     align=PP_ALIGN.CENTER)
footer(s)

# =====================================================================
# SLIDE 7 -- Line W temperature gradient
# =====================================================================
s = prs.slides.add_slide(BLANK)
header(s, "07", "Line W \u2014 Temperature Gradient Annotation",
       "WPS warm return: 4.5 K (QCELL side) \u2192 300 K (USER side / QRB)")
_img_fit(s, asset("crop_gradient.png"), Inches(0.5), Inches(1.5), Inches(8.3),
         Inches(2.6))
_bullets(s, Inches(0.55), Inches(4.4), Inches(12.3), Inches(2.6), [
    (0, "Line W carries cold helium that is progressively warmed back to "
        "ambient before the QRB handover", NAVY, True),
    (1, "cold end 4.5 K (QCELL / cold box side)", GREY),
    (1, "electrical heater / ambient gain along the run", GREY),
    (1, "warm end ~300 K at the USER side (NA.CP03 QRB)", GREY),
    (0, "Flow rate ~2.5 g/s \u00b7 DN20 \u00b7 SS304", ORANGE, True),
    (0, "Gradient bar drawn on layer 15_Temperature_Gradient and reproduced "
        "in the MAIN-LINES schematic", GREY),
], size=14, gap=6)
footer(s)

# =====================================================================
# SLIDE 8 -- Deliverables / closing
# =====================================================================
s = prs.slides.add_slide(BLANK)
header(s, "08", "v5 Deliverables",
       "Everything ships under output_v5/")
_bullets(s, Inches(0.55), Inches(1.4), Inches(6.1), Inches(5.4), [
    (0, "16 production sheets", NAVY, True),
    (1, "QCELL + RFCELL \u00d7 2 sheets \u00d7 {STANDARD, CONTROL-CENTRIC} "
        "\u00d7 {colour, mono}", GREY),
    (1, "each as SVG + A3 PDF + review PNG", GREY),
    (0, "2 MAIN-LINES-ONLY schematics (QCELL + RFCELL)", NAVY, True),
    (0, "LINE_SPECIFICATION_MASTER.xlsx", NAVY, True),
    (1, "01_Line_Specification \u00b7 02_Main_vs_Branch \u00b7 03_Source_Circuits",
        GREY),
    (0, "MINERVA_PID_DISSECTION_v5.pptx (this deck)", ORANGE, True),
    (0, "COLOR_SCHEME_REVISION_GUIDE.pdf", NAVY, True),
    (0, "TEMPERATURE_GRADIENT_VISUALIZATION.pdf", NAVY, True),
], size=14, gap=7)
_img_fit(s, asset("ml_r.png"), Inches(6.9), Inches(1.4), Inches(6.0), Inches(4.6))
_txt(s, Inches(6.9), Inches(6.05), Inches(6.0), Inches(0.3),
     "RFCELL \u2014 MAIN-LINES-ONLY schematic", size=10, color=GREY,
     italic=True, align=PP_ALIGN.CENTER)
footer(s)

prs.save(OUT)
print("wrote", OUT, "slides:", len(prs.slides._sldIdLst))
