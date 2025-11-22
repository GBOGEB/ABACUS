"""Generate a QPLANT/SCK CEN aligned PowerPoint template.

This script creates a professional PowerPoint presentation that follows the
SCK CEN corporate style guide. The generated template includes:

- Consistent blue title bars with white text
- Calibri typography throughout
- Numbered footers with customizable document codes
- Six pre-configured slide types for common presentation needs
- 16:9 aspect ratio optimized for modern displays

The template uses a cohesive color palette:
- Deep blue (#003F72) for headers and accents
- Cyan (#0086C0) for highlights and links
- Neutral sand (#E8E7E2) for content backgrounds
- Dark gray (#222222) for body text

Usage:
    python qplant_sckcen_template.py --doc-number "SCK CEN/1427"
    python qplant_sckcen_template.py --output custom/path.pptx

The generated .pptx file can be used directly or as a Pandoc reference template.
"""

import argparse
from pathlib import Path
from pptx import Presentation
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

# SCK CEN Corporate Color Palette
# These colors define the visual identity for official presentations
PRIMARY_COLOR = RGBColor(0, 63, 114)    # Deep blue used for title bars and headers
ACCENT_COLOR = RGBColor(0, 134, 192)    # Cyan accent for callouts and links
SAND_COLOR = RGBColor(232, 231, 226)    # Light neutral background for content
TEXT_COLOR = RGBColor(34, 34, 34)       # Dark gray for readable body text

# Typography constants - Calibri is the SCK CEN standard font
TITLE_FONT = "Calibri"
BODY_FONT = "Calibri"

# Slide dimensions for 16:9 aspect ratio (standard for modern displays)
SLIDE_WIDTH = Inches(13.33)
SLIDE_HEIGHT = Inches(7.5)


def _text_box(slide, left, top, width, height, text="", font_size=18, bold=False,
              color=TEXT_COLOR, align=PP_ALIGN.LEFT):
    """Create a formatted text box on a slide.
    
    This is a helper function that simplifies adding text to slides with
    consistent formatting. It handles font styling, sizing, and alignment.
    
    Args:
        slide: The slide object to add the text box to
        left: Distance from left edge (in Inches)
        top: Distance from top edge (in Inches)
        width: Width of the text box (in Inches)
        height: Height of the text box (in Inches)
        text: Text content to display (default: empty string)
        font_size: Font size in points (default: 18)
        bold: Whether to make text bold (default: False)
        color: RGB color for the text (default: TEXT_COLOR)
        align: Text alignment (default: left-aligned)
    
    Returns:
        TextFrame object that can be further customized if needed
    """
    box = slide.shapes.add_textbox(left, top, width, height)
    frame = box.text_frame
    frame.clear()
    p = frame.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.name = BODY_FONT
    p.font.color.rgb = color
    p.alignment = align
    return frame


def _add_footer(slide, doc_number, slide_idx, total_slides):
    """Add a standard footer with document code and page numbers.
    
    Creates a thin blue line above the footer area, then adds the document
    code on the left and page numbers (X/Y format) on the right. This footer
    appears on every slide for consistency.
    
    Args:
        slide: The slide object to add the footer to
        doc_number: Document code (e.g., "SCK CEN/1427")
        slide_idx: Current slide number (1-based)
        total_slides: Total number of slides in the presentation
    """
    width = SLIDE_WIDTH
    height = SLIDE_HEIGHT

    # Thin horizontal line separator above the footer
    line = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        Inches(0.3),
        height - Inches(0.55),
        width - Inches(0.6),
        Inches(0.02),
    )
    line.fill.solid()
    line.fill.fore_color.rgb = PRIMARY_COLOR
    line.line.fill.background()

    # Document code on the left side
    _text_box(
        slide,
        Inches(0.3),
        height - Inches(0.5),
        width - Inches(1.6),
        Inches(0.35),
        text=f"Document: {doc_number}",
        font_size=12,
        color=PRIMARY_COLOR,
    )
    
    # Page numbers on the right side (e.g., "3/6")
    _text_box(
        slide,
        width - Inches(1.3),
        height - Inches(0.5),
        Inches(1.0),
        Inches(0.35),
        text=f"{slide_idx}/{total_slides}",
        font_size=12,
        align=PP_ALIGN.RIGHT,
        color=PRIMARY_COLOR,
    )


def _title_block(slide, title, subtitle=""):
    """Create the blue header band at the top of a slide.
    
    This creates the distinctive SCK CEN slide header with a deep blue
    background, white title text, and optional lighter subtitle text.
    
    Args:
        slide: The slide object to add the title block to
        title: Main title text (large, bold, white)
        subtitle: Optional subtitle text (smaller, light blue)
    """
    # Blue rectangle spanning the full width at the top
    header = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        Inches(0),
        Inches(0),
        SLIDE_WIDTH,
        Inches(1.5),
    )
    header.fill.solid()
    header.fill.fore_color.rgb = PRIMARY_COLOR
    header.line.fill.background()  # No border

    # Main title in white
    title_frame = _text_box(
        slide,
        Inches(0.6),
        Inches(0.35),
        SLIDE_WIDTH - Inches(1.2),
        Inches(0.8),
        text=title,
        font_size=38,
        bold=True,
        color=RGBColor(255, 255, 255),
    )
    title_frame.paragraphs[0].font.name = TITLE_FONT

    # Optional subtitle in light blue
    if subtitle:
        subtitle_frame = _text_box(
            slide,
            Inches(0.6),
            Inches(1.1),
            SLIDE_WIDTH - Inches(1.2),
            Inches(0.5),
            text=subtitle,
            font_size=18,
            color=RGBColor(220, 235, 247),
        )
        subtitle_frame.paragraphs[0].font.name = BODY_FONT


def _agenda_slide(prs, doc_number):
    """Generate the opening agenda slide.
    
    Creates a slide with a numbered list of typical presentation sections.
    Users can customize this list for their specific presentation needs.
    
    Args:
        prs: Presentation object
        doc_number: Document code for the footer
    
    Returns:
        The created slide object
    """
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _title_block(slide, "Agenda", "Reusable section outline")

    content = slide.shapes.add_textbox(Inches(0.8), Inches(1.8), prs.slide_width - Inches(1.6), Inches(4.5))
    tf = content.text_frame
    tf.text = "1. Executive summary\n2. System scope and constraints\n3. Integration timeline\n4. Interfaces and dependencies\n5. Risks and mitigations\n6. Next actions"
    for p in tf.paragraphs:
        p.font.name = BODY_FONT
        p.font.size = Pt(22)
        p.font.color.rgb = TEXT_COLOR
    return slide


def _status_slide(prs, doc_number):
    """Generate a two-column status update slide.
    
    Creates a slide with "Highlights" on the left (recent achievements)
    and "KPI snapshot" on the right (key metrics). Both sections use
    colored backgrounds to visually separate the content.
    
    Args:
        prs: Presentation object
        doc_number: Document code for the footer
    
    Returns:
        The created slide object
    """
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _title_block(slide, "Project status", "Drop in current data for weekly reviews")

    # Left panel: Highlights with sand background
    left_block = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        Inches(0.6),
        Inches(1.8),
        prs.slide_width / 2 - Inches(0.9),
        Inches(3.5),
    )
    left_block.fill.solid()
    left_block.fill.fore_color.rgb = SAND_COLOR
    left_block.line.color.rgb = PRIMARY_COLOR

    left_tf = left_block.text_frame
    left_tf.text = "Highlights"  # heading
    left_tf.paragraphs[0].font.size = Pt(20)
    left_tf.paragraphs[0].font.bold = True
    left_tf.paragraphs[0].font.name = TITLE_FONT
    left_tf.paragraphs[0].font.color.rgb = PRIMARY_COLOR
    
    # Sample highlight bullets (replace with actual project highlights)
    for bullet in [
        "Cryogenic loop sizing validated against SCK CEN purge parameters.",
        "Control cabinet design frozen; vendor RFQs issued.",
        "Integration test window confirmed with MYRRHA operations.",
    ]:
        p = left_tf.add_paragraph()
        p.text = bullet
        p.level = 1
        p.font.name = BODY_FONT
        p.font.size = Pt(16)
        p.font.color.rgb = TEXT_COLOR

    # Right panel: KPIs with light blue background
    right_block = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        prs.slide_width / 2 + Inches(0.3),
        Inches(1.8),
        prs.slide_width / 2 - Inches(0.9),
        Inches(3.5),
    )
    right_block.fill.solid()
    right_block.fill.fore_color.rgb = RGBColor(240, 247, 252)
    right_block.line.color.rgb = PRIMARY_COLOR

    right_tf = right_block.text_frame
    right_tf.text = "KPI snapshot"
    right_tf.paragraphs[0].font.size = Pt(20)
    right_tf.paragraphs[0].font.bold = True
    right_tf.paragraphs[0].font.name = TITLE_FONT
    right_tf.paragraphs[0].font.color.rgb = PRIMARY_COLOR
    
    # Sample KPI metrics (replace with actual project metrics)
    for label, value in [
        ("Design maturity", "82% complete"),
        ("Interface agreements", "5/7 signed"),
        ("Risk exposure", "Low (2 open items)"),
    ]:
        p = right_tf.add_paragraph()
        p.text = f"{label}: {value}"
        p.level = 1
        p.font.name = BODY_FONT
        p.font.size = Pt(16)
        p.font.color.rgb = TEXT_COLOR

    return slide


def _architecture_slide(prs, doc_number):
    """Generate a system architecture placeholder slide.
    
    Creates a slide with a large canvas area and sample component blocks.
    Users should replace these with actual architecture diagrams (Visio,
    draw.io exports, etc.).
    
    Args:
        prs: Presentation object
        doc_number: Document code for the footer
    
    Returns:
        The created slide object
    """
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _title_block(slide, "System architecture", "Replace shapes with updated diagrams")

    # Large rounded rectangle as the canvas background
    canvas = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(0.8),
        Inches(1.8),
        prs.slide_width - Inches(1.6),
        Inches(4.0),
    )
    canvas.fill.solid()
    canvas.fill.fore_color.rgb = RGBColor(245, 249, 252)  # Very light blue
    canvas.line.color.rgb = PRIMARY_COLOR

    # Sample component blocks - replace with actual system components
    for idx, title in enumerate(["Helium storage", "Compression train", "Cold box", "Distribution to cryomodules"]):
        block = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
            Inches(1.0 + idx * 1.8),  # Horizontal spacing
            Inches(2.3),
            Inches(1.5),
            Inches(1.2),
        )
        block.fill.solid()
        block.fill.fore_color.rgb = RGBColor(220, 235, 247)  # Light blue
        block.line.color.rgb = ACCENT_COLOR
        
        # Add component label
        tf = block.text_frame
        tf.text = title
        tf.paragraphs[0].font.size = Pt(14)
        tf.paragraphs[0].font.name = BODY_FONT
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].font.color.rgb = TEXT_COLOR

    # Instructional caption for users
    caption = _text_box(
        slide,
        Inches(0.9),
        Inches(4.3),
        prs.slide_width - Inches(1.8),
        Inches(0.8),
        text="Use this slide to paste Visio exports or draw lightweight flows linking helium storage, compression, cold box, and distribution headers.",
        font_size=14,
        color=TEXT_COLOR,
    )
    caption.paragraphs[0].font.name = BODY_FONT
    return slide


def _timeline_slide(prs, doc_number):
    """Generate a schedule and milestones slide.
    
    Creates a slide with quarterly milestones and a hyperlink to dependencies.
    The milestone format shows quarter labels on the left with descriptions
    on the right.
    
    Args:
        prs: Presentation object
        doc_number: Document code for the footer
    
    Returns:
        The created slide object
    """
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _title_block(slide, "Schedule & milestones", "Quarterly checkpoints with owner and dependencies")

    # Sample quarterly milestones - replace with actual project timeline
    milestones = [
        ("Q1", "Finalize P&ID freeze and HAZOP"),
        ("Q2", "Factory acceptance for compression skids"),
        ("Q3", "Cold box delivery to SCK CEN"),
        ("Q4", "Site integration and performance test"),
    ]

    # Create milestone entries with quarter labels and descriptions
    for idx, (label, desc) in enumerate(milestones):
        top = Inches(2.0 + idx * 0.9)
        
        # Quarter label (e.g., "Q1")
        _text_box(
            slide,
            Inches(0.8),
            top,
            Inches(1.0),
            Inches(0.4),
            text=label,
            font_size=16,
            bold=True,
            color=PRIMARY_COLOR,
        )
        
        # Milestone description
        _text_box(
            slide,
            Inches(1.8),
            top,
            prs.slide_width - Inches(2.6),
            Inches(0.4),
            text=desc,
            font_size=16,
        )

    # Add a note with a clickable hyperlink to dependencies
    note = _text_box(
        slide,
        Inches(0.8),
        Inches(5.0),
        prs.slide_width - Inches(1.6),
        Inches(0.6),
        text="Dependencies documented in the SCK CEN interface matrix (clickable link).",
        font_size=14,
        color=ACCENT_COLOR,
    )
    
    # Add hyperlink to the note
    run = note.paragraphs[0].add_run()
    run.text = "  https://sharepoint.sckcen.be/qplant/interfaces"
    run.hyperlink.address = "https://sharepoint.sckcen.be/qplant/interfaces"
    run.font.color.rgb = ACCENT_COLOR
    run.font.name = BODY_FONT
    run.font.size = Pt(14)
    return slide


def _risk_slide(prs, doc_number):
    """Generate a risk and mitigation register slide.
    
    Creates a three-column table showing risks, their impacts, and
    mitigation strategies. The table format makes it easy to track
    and communicate risk management activities.
    
    Args:
        prs: Presentation object
        doc_number: Document code for the footer
    
    Returns:
        The created slide object
    """
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _title_block(slide, "Risk and mitigation register", "Keep the table; replace rows with your current items")

    # Sample risk entries - replace with actual project risks
    rows = [
        ("Cryogenic valve lead time", "Vendor lead time exceeds site window", "Place advance order; keep dual suppliers"),
        ("Interface firmware", "PLC safety libraries misaligned", "Synchronize versions with SCK CEN validation bench"),
        ("Site readiness", "Utility hook-ups delayed", "Stage temporary utilities; pre-commission compressors"),
    ]

    # Create table with header row + data rows
    table = slide.shapes.add_table(
        len(rows) + 1,  # +1 for header row
        3,              # 3 columns: Risk, Impact, Mitigation
        Inches(0.6),
        Inches(1.9),
        prs.slide_width - Inches(1.2),
        Inches(3.6)
    ).table
    
    # Set column widths for balanced display
    table.columns[0].width = Inches(2.2)
    table.columns[1].width = Inches(3.0)
    table.columns[2].width = Inches(3.0)

    # Format header row with bold blue text
    headers = ["Risk", "Impact", "Mitigation"]
    for idx, header in enumerate(headers):
        cell = table.cell(0, idx)
        cell.text = header
        para = cell.text_frame.paragraphs[0]
        para.font.bold = True
        para.font.name = TITLE_FONT
        para.font.size = Pt(16)
        para.font.color.rgb = PRIMARY_COLOR

    # Fill in data rows
    for row_idx, row in enumerate(rows, start=1):
        for col_idx, value in enumerate(row):
            cell = table.cell(row_idx, col_idx)
            cell.text = value
            para = cell.text_frame.paragraphs[0]
            para.font.name = BODY_FONT
            para.font.size = Pt(14)
            para.font.color.rgb = TEXT_COLOR
    return slide


def _closing_slide(prs, doc_number):
    """Generate a next actions checklist slide.
    
    Creates the final slide with action items, owners, and due dates.
    This helps ensure clear follow-up from presentations.
    
    Args:
        prs: Presentation object
        doc_number: Document code for the footer
    
    Returns:
        The created slide object
    """
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _title_block(slide, "Next actions", "Update dates, owners, and status for handover")

    # Sample action items - replace with actual next steps
    checklist = [
        ("Update purge parameter alignment", "Owner: Process team", "Due: 15 Feb"),
        ("Freeze interface signals with controls", "Owner: PLC lead", "Due: 22 Feb"),
        ("Publish installation work-pack", "Owner: Site engineering", "Due: 01 Mar"),
    ]

    # Create bullet list with action items, owners, and dates
    for idx, (item, owner, due) in enumerate(checklist):
        _text_box(
            slide,
            Inches(0.8),
            Inches(1.9 + idx * 0.9),
            prs.slide_width - Inches(1.6),
            Inches(0.5),
            text=f"• {item} ({owner}) — {due}",
            font_size=18,
            color=TEXT_COLOR,
        )

    # Add reminder about PDF export for external distribution
    note = _text_box(
        slide,
        Inches(0.8),
        Inches(4.8),
        prs.slide_width - Inches(1.6),
        Inches(0.8),
        text="Reminder: export PDF with embedded fonts before circulating outside SCK CEN.",
        font_size=14,
        color=PRIMARY_COLOR,
    )
    note.paragraphs[0].font.bold = True
    return slide


def build_template(output_path: Path, doc_number: str = "SCK CEN/0000") -> None:
    """Build the complete PowerPoint template.
    
    Creates a new presentation with all standard slides and applies consistent
    formatting, colors, and footers throughout. This is the main entry point
    for template generation.
    
    Args:
        output_path: Path where the .pptx file should be saved
        doc_number: Document code to display in footers (default: "SCK CEN/0000")
    
    The generated presentation includes:
        1. Agenda slide
        2. Project status slide (highlights & KPIs)
        3. System architecture slide (placeholder diagram)
        4. Schedule & milestones slide
        5. Risk & mitigation register
        6. Next actions checklist
    """
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    # Update global dimensions to match presentation size
    global SLIDE_WIDTH, SLIDE_HEIGHT
    SLIDE_WIDTH = prs.slide_width
    SLIDE_HEIGHT = prs.slide_height

    # Generate all standard slides
    slides = [
        _agenda_slide(prs, doc_number),
        _status_slide(prs, doc_number),
        _architecture_slide(prs, doc_number),
        _timeline_slide(prs, doc_number),
        _risk_slide(prs, doc_number),
        _closing_slide(prs, doc_number),
    ]

    # Add footer to every slide with page numbers
    for idx, slide in enumerate(slides, start=1):
        _add_footer(slide, doc_number, idx, len(slides))

    prs.save(output_path)


if __name__ == "__main__":
    parser = argparse.ArgumentParser(description="Generate the QPLANT/SCK CEN slide template")
    parser.add_argument(
        "--doc-number",
        default="SCK CEN/0000",
        help="Document number to display in the footer on every slide",
    )
    parser.add_argument(
        "--output",
        type=Path,
        default=Path(__file__).with_name("qplant_sckcen_template.pptx"),
        help="Where to write the generated PowerPoint deck",
    )
    args = parser.parse_args()

    build_template(args.output, doc_number=args.doc_number)
    print(f"Generated {args.output}")
