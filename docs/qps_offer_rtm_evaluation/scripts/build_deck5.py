"""
build_deck5.py -- Phase 5 amendment to QPS_MTBF_WCS_DMAIC_v4.pptx.

Adds the requested energy-consumption breakdown (HP Compressors vs PVPS vs
Total plant load) to Slide 6 (Configuration Baseline), sourced from the
newly-available canonical contract data (RTM-395, Table 19 "Compressor
Room (CCB) electrical supply and load constraints" and Table 20 "Cold Box
Room (AUB) electrical supply & load constraints", QPS_Contract_mirror_DOCX.pdf).

This is an AMENDMENT, not a new slide -- deliberately, to avoid a full
deck-wide slide-renumbering pass (the deck has ~30 internal "Slide N"
cross-references; inserting a new slide before Slide 38 would touch most of
them). Slide 6 already states the "1.4 MW total installed" HP-compressor
figure and has unused vertical space in its right column, making it the
natural, lowest-risk home for the plant-wide electrical mix.

Also blends turquoise and light-purple/pink hues into the palette per GBO's
request, via this new chart's colour choices, without re-touching the
already-QA'd v4 element colours.
"""
from pptx import Presentation
from pptx.util import Emu, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

SRC = "QPS_MTBF_WCS_DMAIC_v4.pptx"
OUT = "QPS_MTBF_WCS_DMAIC_v5.pptx"

# Same figures used to build energy_mix_donut.png -- kept here too so the
# on-slide labels and the chart image can never drift out of sync.
HP_COMPRESSORS = 4 * 356
PVPS = 150
REST = (3 * 42) + 65 + 7.5 + 1.5 + 3 + 3
TOTAL = HP_COMPRESSORS + PVPS + REST

COLORS = {
    "HP": RGBColor(0x56, 0x28, 0x73),   # matches this slide's own DEFINE purple
    "PVPS": RGBColor(0x1F, 0xA7, 0xA0),  # new turquoise
    "REST": RGBColor(0xE0, 0xA9, 0xD6),  # new light purple/pink
}

prs = Presentation(SRC)
slide = prs.slides[5]  # Slide 6

# donut image -- fills the unused gap between the "WHY IT MATTERS" box
# (bottom 4500000) and the "KEY TAKEAWAY" bar (top 5950000)
pic = slide.shapes.add_picture(
    "energy_mix_donut.png",
    left=Emu(7650000), top=Emu(4560000),
    width=Emu(1330000), height=Emu(1330000),
)

# label column, right of the donut
labels_box = slide.shapes.add_textbox(
    Emu(9150000), Emu(4560000), Emu(2670000), Emu(1350000)
)
tf = labels_box.text_frame
tf.word_wrap = True
rows = [
    ("HP Compressors", f"{HP_COMPRESSORS:,} kW (80%)", COLORS["HP"]),
    ("PVPS", f"{PVPS} kW (8%)", COLORS["PVPS"]),
    ("Rest of plant", f"{REST:.0f} kW (12%)", COLORS["REST"]),
]
for i, (label, value, color) in enumerate(rows):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    p.space_after = Pt(2)
    p.space_before = Pt(0)
    p.line_spacing = 1.0
    r1 = p.add_run()
    r1.text = "■ "
    r1.font.size = Pt(11.5)
    r1.font.color.rgb = color
    r1.font.name = "Aptos"
    r2 = p.add_run()
    r2.text = f"{label} — {value}"
    r2.font.size = Pt(11.5)
    r2.font.bold = False
    r2.font.color.rgb = RGBColor(0x33, 0x33, 0x3D)
    r2.font.name = "Aptos"

p = tf.add_paragraph()
p.space_before = Pt(4)
r = p.add_run()
r.text = "RTM-395, Table 19–20 (excl. back-up diesel)"
r.font.size = Pt(8)
r.font.italic = True
r.font.color.rgb = RGBColor(0x70, 0x70, 0x78)
r.font.name = "Aptos"

prs.save(OUT)
print(f"saved {OUT}")
print(f"HP={HP_COMPRESSORS}kW PVPS={PVPS}kW REST={REST}kW TOTAL={TOTAL}kW")
