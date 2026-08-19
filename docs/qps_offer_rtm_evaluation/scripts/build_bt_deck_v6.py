"""
build_bt_deck_v6.py -- Phase 6 addition to BT_Method_Evaluation_v5.pptx.

GBO's ask: "the plots for deck of BT could be blown up to make these 1x2 or
1x1 diagrams in own right - maybe with some grids as 'master or summery'
slides".

Adds 6 new slides at the end (14-19), appended after EVAL-S06 (slide 13):
  - CHART-S00: a master/summary grid of all 5 charts at a glance, for
    orientation and as a printable one-pager.
  - CHART-S01..S05: each of the 5 charts blown up to near-full-slide size
    (1x1), one per slide, with a short caption reused from its source
    slide's "board read" line so each stands on its own without forcing the
    audience back to the original slide.

Chosen NOT to touch slides 1-13 at all -- every cross-reference in this deck
uses METHOD-Sxx/EVAL-Sxx codes and slide-number placeholders (both verified
during the v5 restyle to auto-update / not be hardcoded "Slide N" text), so
purely appending new slides at the end carries no renumbering risk.

Source images are extracted byte-for-byte from the v5 deck's own embedded
media (not regenerated), so what's blown up here is pixel-identical to what
appears small on slides 5/8/9/10/11.
"""
import warnings
warnings.filterwarnings("ignore")
from pptx import Presentation
from pptx.util import Emu, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from PIL import Image
import copy

IN = "BT_Method_Evaluation_v5.pptx"
OUT = "BT_Method_Evaluation_v6.pptx"

HEADER = RGBColor(0x44, 0x1F, 0x63)
FOOTER_BG = RGBColor(0xF3, 0xEC, 0xF8)
CALLOUT_BG = RGBColor(0xEF, 0xE5, 0xF5)
ACCENT = RGBColor(0x0B, 0x3D, 0x5C)
BODY = RGBColor(0x22, 0x22, 0x22)
FOOTER_TXT = RGBColor(0x55, 0x55, 0x55)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
CODE_TXT = RGBColor(0xBF, 0xD6, 0xE6)

SLIDE_W = Emu(12192000)
SLIDE_H = Emu(6858000)
HEADER_H = Emu(868680)
FOOTER_TOP = Emu(6547104)
LEFT_MARGIN = Emu(320040)
CONTENT_RIGHT = SLIDE_W - LEFT_MARGIN

CHARTS = [
    {"src_slide": 5, "src_code": "METHOD-S04", "new_code": "CHART-S01",
     "title": "From scores to ranking: S → win % → λ (enlarged)",
     "caption": "Quote S to defend a ranking, win % to brief, λ on request — all three are deterministic transforms of the same weighted scores.",
     "img": "slide5_chart.png"},
    {"src_slide": 8, "src_code": "EVAL-S01", "new_code": "CHART-S02",
     "title": "OFFER importance result (enlarged)",
     "caption": "The two leaders (OFFER-21/22) are both LOOP/recovery-strategy items — see EVAL-S01 for the full top-6 driver breakdown.",
     "img": "slide8_chart.png"},
    {"src_slide": 9, "src_code": "EVAL-S02", "new_code": "CHART-S03",
     "title": "Robustness: weight sensitivity (enlarged)",
     "caption": "528 scenarios — one-at-a-time ±25%/±50% per dimension, plus combinatorial sweeps. See EVAL-S02 for the full scenario methodology.",
     "img": "slide9_chart.png"},
    {"src_slide": 10, "src_code": "EVAL-S03", "new_code": "CHART-S04",
     "title": "RTM requirement importance (enlarged)",
     "caption": "722 requirements ranked on identical weights — see EVAL-S03 for tier definitions and the P1-critical cutoff.",
     "img": "slide10_chart.png"},
    {"src_slide": 11, "src_code": "EVAL-S04", "new_code": "CHART-S05",
     "title": "Where the contract weight sits (enlarged)",
     "caption": "Use for review staffing — put your strongest reviewers on the highest-weight sections. See EVAL-S04 for the full section ranking.",
     "img": "slide11_chart.png"},
]

prs = Presentation(IN)
blank_layout = prs.slide_layouts[6]  # "Blank"
assert blank_layout.name == "Blank"


def add_header_footer(slide, title, code, footer_note):
    rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Emu(0), Emu(0), SLIDE_W, HEADER_H)
    rect.fill.solid(); rect.fill.fore_color.rgb = HEADER
    rect.line.fill.background()
    rect.shadow.inherit = False

    tb = slide.shapes.add_textbox(Emu(411480), Emu(109728), Emu(9875520), Emu(685800))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    r = p.add_run(); r.text = title
    r.font.size = Pt(25); r.font.bold = True; r.font.color.rgb = WHITE; r.font.name = "Aptos"

    cb = slide.shapes.add_textbox(Emu(10607040), Emu(292608), Emu(1371600), Emu(320040))
    ctf = cb.text_frame
    cp = ctf.paragraphs[0]
    cp.alignment = PP_ALIGN.RIGHT
    cr = cp.add_run(); cr.text = code
    cr.font.size = Pt(11); cr.font.bold = True; cr.font.color.rgb = CODE_TXT; cr.font.name = "Aptos"

    frect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Emu(0), FOOTER_TOP, SLIDE_W, Emu(310896))
    frect.fill.solid(); frect.fill.fore_color.rgb = FOOTER_BG
    frect.line.fill.background()
    frect.shadow.inherit = False

    fb = slide.shapes.add_textbox(Emu(411480), Emu(6565392), Emu(10515600), Emu(274320))
    ftf = fb.text_frame
    fp = ftf.paragraphs[0]
    fr = fp.add_run(); fr.text = footer_note
    fr.font.size = Pt(8.5); fr.font.color.rgb = FOOTER_TXT; fr.font.name = "Aptos"

    return rect, frect


def add_caption(slide, top, text, lead_in=None):
    box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, LEFT_MARGIN, top, CONTENT_RIGHT - LEFT_MARGIN, Emu(548640))
    box.fill.solid(); box.fill.fore_color.rgb = CALLOUT_BG
    box.line.fill.background()
    box.shadow.inherit = False
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(137160); tf.margin_top = Emu(64008); tf.margin_right = Emu(137160)
    p = tf.paragraphs[0]
    if lead_in:
        r0 = p.add_run(); r0.text = lead_in
        r0.font.bold = True; r0.font.size = Pt(12.5); r0.font.color.rgb = ACCENT; r0.font.name = "Aptos"
    r1 = p.add_run(); r1.text = text
    r1.font.size = Pt(12.5); r1.font.color.rgb = BODY; r1.font.name = "Aptos"
    return box


# ============================================================ CHART-S00 ===
# master/summary grid -- all 5 charts at a glance
slide = prs.slides.add_slide(blank_layout)
add_header_footer(
    slide, "Chart appendix — all 5 diagrams at a glance", "CHART-S00",
    "BT METHOD — Engineering Baseline (canonical RTM-722) · CHART-S00 | v3.6",
)
sub = slide.shapes.add_textbox(Emu(411480), Emu(900000), Emu(11000000), Emu(280000))
stf = sub.text_frame
sp = stf.paragraphs[0]
sr = sp.add_run()
sr.text = "Master index — each thumbnail below is enlarged on its own slide (CHART-S01…S05) right after this one."
sr.font.size = Pt(12); sr.font.italic = True; sr.font.color.rgb = FOOTER_TXT; sr.font.name = "Aptos"

# 3-then-2 grid layout, cells sized to the content band between header and footer
GRID_TOP = Emu(1280000)
GRID_BOTTOM = Emu(6420000)
row_h = (GRID_BOTTOM - GRID_TOP) // 2
col_w3 = (CONTENT_RIGHT - LEFT_MARGIN) // 3
col_w2 = (CONTENT_RIGHT - LEFT_MARGIN) // 2
CELL_PAD = Emu(90000)
LABEL_H = Emu(430000)

def place_thumb(cell_left, cell_top, cell_w, cell_h, chart):
    img = Image.open(f"/tmp/bt_charts/{chart['img']}")
    ar = img.size[0] / img.size[1]
    avail_w = cell_w - 2 * CELL_PAD
    avail_h = cell_h - 2 * CELL_PAD - LABEL_H
    if avail_w / avail_h > ar:
        h = avail_h
        w = int(h * ar)
    else:
        w = avail_w
        h = int(w / ar)
    left = cell_left + (cell_w - w) // 2
    top = cell_top + CELL_PAD
    slide.shapes.add_picture(f"/tmp/bt_charts/{chart['img']}", left, top, width=w, height=h)
    lab = slide.shapes.add_textbox(cell_left, cell_top + cell_h - LABEL_H, cell_w, LABEL_H)
    ltf = lab.text_frame
    ltf.word_wrap = True
    lp = ltf.paragraphs[0]
    lp.alignment = PP_ALIGN.CENTER
    lr = lp.add_run(); lr.text = chart["title"].replace(" (enlarged)", "")
    lr.font.size = Pt(11); lr.font.bold = True; lr.font.color.rgb = HEADER; lr.font.name = "Aptos"
    lp2 = ltf.add_paragraph()
    lp2.alignment = PP_ALIGN.CENTER
    lr2 = lp2.add_run(); lr2.text = f"{chart['new_code']}  ·  from {chart['src_code']}"
    lr2.font.size = Pt(9); lr2.font.color.rgb = FOOTER_TXT; lr2.font.name = "Aptos"

top_row = CHARTS[:3]
bottom_row = CHARTS[3:]
for i, c in enumerate(top_row):
    place_thumb(LEFT_MARGIN + i * col_w3, GRID_TOP, col_w3, row_h, c)
bottom_offset = LEFT_MARGIN + (CONTENT_RIGHT - LEFT_MARGIN - 2 * col_w2) // 2
for i, c in enumerate(bottom_row):
    place_thumb(bottom_offset + i * col_w2, GRID_TOP + row_h, col_w2, row_h, c)

print("CHART-S00 (master grid) added")

# ======================================================= CHART-S01..S05 ===
IMG_TOP = Emu(1050000)
CAPTION_H = Emu(548640)
CAPTION_GAP = Emu(90000)
IMG_BOTTOM_MAX = FOOTER_TOP - Emu(90000) - CAPTION_H - CAPTION_GAP

for c in CHARTS:
    slide = prs.slides.add_slide(blank_layout)
    add_header_footer(
        slide, c["title"], c["new_code"],
        f"BT METHOD — Engineering Baseline (canonical RTM-722) · {c['new_code']} (enlarged from {c['src_code']}) | v3.6",
    )
    img = Image.open(f"/tmp/bt_charts/{c['img']}")
    ar = img.size[0] / img.size[1]
    avail_w = CONTENT_RIGHT - LEFT_MARGIN
    avail_h = IMG_BOTTOM_MAX - IMG_TOP
    if avail_w / avail_h > ar:
        h = avail_h
        w = int(h * ar)
    else:
        w = avail_w
        w = int(avail_w)
        h = int(w / ar)
    left = LEFT_MARGIN + (avail_w - w) // 2
    slide.shapes.add_picture(f"/tmp/bt_charts/{c['img']}", left, IMG_TOP, width=w, height=h)
    caption_top = IMG_TOP + h + CAPTION_GAP
    if caption_top + CAPTION_H > FOOTER_TOP - Emu(60000):
        caption_top = FOOTER_TOP - Emu(60000) - CAPTION_H
    add_caption(slide, caption_top, c["caption"], lead_in="Board read:  ")
    print(f"{c['new_code']} (enlarged from {c['src_code']}) added")

prs.save(OUT)
print(f"saved {OUT} -- {len(prs.slides.__iter__.__self__._sldIdLst)} slides" if False else f"saved {OUT} -- {len(prs.slides._sldIdLst)} slides")
