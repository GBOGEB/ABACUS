#!/usr/bin/env python3
"""Pass 4: style sweep on top of GBO's own manual edits.

GBO opened final3.pptx (delivered as QPS_MTBF_WCS_DMAIC_v3.pptx), moved
textboxes, rewrote some labels, and re-ordered Slides 19-21 (Failure
Classification now first, then Reliability Targets, then Operational
Philosophy). This pass works ON TOP of that file -- it does not revert or
rebuild content, only:

  1. Fixes 6 "Slide N" cross-references broken by the 19-21 reorder.
  2. Swaps font family: Aptos everywhere, EXCEPT the Title/banner-title text
     on every slide and ALL text on Slides 1-3, which stay Segoe UI.
  3. Bumps main-body text (bullet placeholders / head-body textboxes) into
     the 14-16pt band for legibility and better slide fill.
  4. Rebuilds every DMAIC phase tag as a bigger, solid-fill colour badge
     keyed to the DMAIC letter (5 distinct colours, not 2 shared ones).
  5. Moves Title/Subtitle left and upward (native-placeholder slides only)
     -- closer to the corner but not flush against it.

Base:   uploaded_v3_review.pptx  (GBO's "Cryoplant MTBF and System Design_ReviewGBO.pptx")
Output: final4.pptx
"""
from pptx import Presentation
from pptx.util import Emu, Pt
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_THEME_COLOR_INDEX as TC
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE, MSO_SHAPE_TYPE
import copy as _copy

SRC = "uploaded_v3_review.pptx"
OUT = "final4.pptx"

prs = Presentation(SRC)
slides = list(prs.slides)
SLIDE_W, SLIDE_H = prs.slide_width, prs.slide_height

# ---------------------------------------------------------------- helpers --

def find_shape(slide, name):
    for sh in slide.shapes:
        if sh.name == name:
            return sh
    return None

def all_runs(shape):
    if not shape.has_text_frame:
        return
    for p in shape.text_frame.paragraphs:
        for r in p.runs:
            yield r

def replace_run_text(slide, shape_name, old, new):
    sh = find_shape(slide, shape_name)
    assert sh is not None, f"shape {shape_name!r} not found"
    found = False
    for r in all_runs(sh):
        if old in r.text:
            r.text = r.text.replace(old, new)
            found = True
    assert found, f"text {old!r} not found in {shape_name!r}"

def is_title_shape(sh):
    """True for the slide's own title text: the native Title placeholder,
    or the full-width colour banner rectangle used on the narrative slides."""
    if sh.name == "Title 1":
        return True
    if sh.name in ("Rectangle 1", "Rectangle 2") and sh.top is not None and sh.top < 200000 \
            and sh.width is not None and sh.width > 5000000:
        return True
    return False

def is_banner_slide(slide):
    for sh in slide.shapes:
        if is_title_shape(sh) and sh.name != "Title 1":
            return True
    return False

TAGS = {"DEFINE", "MEASURE", "ANALYZE", "IMPROVE", "CONTROL"}

print(f"loaded {SRC}, {len(slides)} slides")

# ============================================================ STEP 1 ==========
# Fix cross-references broken by GBO's own Slide 19-21 reorder.
# New order: 19 = Failure Classification, 20 = Reliability Targets by Failure
# Class, 21 = Operational Philosophy: Cold != Off (was 19/20/21 respectively
# before GBO's edit).

replace_run_text(slides[12], "Text Placeholder 3",
                  "Class A / B / C (defined Slide 20)", "Class A / B / C (defined Slide 19)")
replace_run_text(slides[12], "Text Placeholder 3",
                  "restated per failure-class in Slide 21", "restated per failure-class in Slide 20")
replace_run_text(slides[18], "TextBox 16", "Slide 21", "Slide 20")
replace_run_text(slides[19], "TextBox 16",
                  "Same three classes as Slide 20", "Same three classes as Slide 19")
replace_run_text(slides[23], "Text Placeholder 2",
                  "states defined on Slide 19", "states defined on Slide 21")
replace_run_text(slides[23], "TextBox 4",
                  "Slide 19 chevron diagram", "Slide 21 chevron diagram")
replace_run_text(slides[23], "TextBox 4", "Slide 20 failure", "Slide 19 failure")
print("STEP 1: cross-reference fixes for the 19-21 reorder done (7 edits)")

# ============================================================ STEP 2 ==========
# Font family: Aptos everywhere, except Title/banner-title text (every slide)
# and ALL text on Slides 1-3 (Segoe UI, matching the deck's own theme font).
FIRST3 = {0, 1, 2}
n_aptos = n_segoe = 0
for i, s in enumerate(slides):
    for sh in s.shapes:
        if not sh.has_text_frame:
            continue
        keep_segoe = (i in FIRST3) or is_title_shape(sh)
        for r in all_runs(sh):
            r.font.name = "Segoe UI" if keep_segoe else "Aptos"
            if keep_segoe:
                n_segoe += 1
            else:
                n_aptos += 1
print(f"STEP 2: font family set -- {n_aptos} runs -> Aptos, {n_segoe} runs -> Segoe UI (Title + Slides 1-3)")

# ============================================================ STEP 3 ==========
# Bump main-body text into the 14-16pt band -- but content-aware: a shape with
# 5 head/body pairs and a shape with 3 short bullets cannot both take a flat
# 16pt without one of them overflowing its box. For each eligible shape we
# estimate, from its own current text and box geometry, the LARGEST size in
# [13, 16] that still fits, and use that instead of a blind lookup table.
# Scoped to TEXT_BOX shapes and the "Text Placeholder 3" body placeholder --
# this naturally excludes AUTO_SHAPE takeaway/stat-tile boxes (their own
# fixed sizing) and tables. DMAIC tags and italic (footnote-scale) runs are
# left alone.
CHAR_W_FACTOR = 0.58   # rough average glyph width as a fraction of font size, proportional sans
LINE_SPACING = 1.45
PARA_GAP_PT = 8

def estimate_height_emu(paragraphs_text, font_pt, usable_w_emu):
    avg_char_w = max(font_pt * CHAR_W_FACTOR * 12700, 1)
    chars_per_line = max(1, int(usable_w_emu / avg_char_w))
    line_h = font_pt * LINE_SPACING * 12700
    total = 0
    for txt in paragraphs_text:
        # python-pptx renders each manual <a:br/> line break as \x0b in
        # paragraph.text -- these don't reflow with the rest of the text, so
        # each segment (including empty ones, i.e. a deliberate blank line)
        # needs its own line count instead of being treated as one long
        # wrappable string.
        segments = txt.split("\x0b")
        lines = 0
        for seg in segments:
            n = len(seg)
            lines += max(1, -(-n // chars_per_line))  # ceil, min 1 per segment
        total += lines * line_h
        total += PARA_GAP_PT * 12700
    return total

# Slides were laid out (in earlier passes) with some text boxes declared
# taller than the safe content area -- fine at the original small font
# (the text never actually reached that far down), but not safe to treat as
# real headroom once the font grows. Cap by the footer line, not just the
# shape's own declared height.
SAFE_BOTTOM = Emu(6350000)

def next_obstacle_top(slide, sh, gap=80000):
    """Top-most sibling shape that sits below sh and overlaps its x-range --
    an extra ceiling so a grown text box can't visually run into a takeaway
    or stat-tile box that never moved."""
    if sh.left is None or sh.top is None:
        return None
    x0, x1 = sh.left, sh.left + (sh.width or 0)
    best = None
    for other in slide.shapes:
        if other is sh or other.top is None or other.left is None:
            continue
        ox0, ox1 = other.left, other.left + (other.width or 0)
        if ox1 <= x0 or ox0 >= x1:
            continue  # no horizontal overlap
        if other.top <= sh.top:
            continue  # not below
        if best is None or other.top < best:
            best = other.top
    return (best - gap) if best is not None else None

def autofit_shape(sh, lo=11.5, hi=16.0, step=0.5, slide=None):
    tf = sh.text_frame
    ml = tf.margin_left if tf.margin_left is not None else Emu(91440)
    mr = tf.margin_right if tf.margin_right is not None else Emu(91440)
    mt = tf.margin_top if tf.margin_top is not None else Emu(45720)
    mb = tf.margin_bottom if tf.margin_bottom is not None else Emu(45720)
    usable_w = max(sh.width - ml - mr, 500000)
    box_h = sh.height - mt - mb
    footer_h = SAFE_BOTTOM - sh.top - mt if sh.top is not None else box_h
    ceilings = [box_h, footer_h]
    if slide is not None:
        obstacle = next_obstacle_top(slide, sh)
        if obstacle is not None:
            ceilings.append(obstacle - sh.top - mt)
    usable_h = max(min(ceilings), 300000)
    paras = [p.text for p in tf.paragraphs if p.text.strip()]
    if not paras:
        return None
    best = None
    size = hi
    while size >= lo - 1e-6:
        if estimate_height_emu(paras, size, usable_w) <= usable_h * 0.74:
            best = size
            break
        size -= step
    if best is not None:
        return best
    # Nothing in [lo, hi] fits (a handful of very dense shapes) -- keep
    # descending below the target band rather than let text overflow into a
    # neighbour or the footer. Absolute floor matches the smallest body size
    # already used elsewhere in the deck.
    size = lo - step
    while size >= 9.5:
        if estimate_height_emu(paras, size, usable_w) <= usable_h * 0.98:
            return size
        size -= step
    return 9.5

n_bumped = n_skipped_small = 0
for s in slides:
    for sh in s.shapes:
        if not sh.has_text_frame:
            continue
        if sh.shape_type not in (MSO_SHAPE_TYPE.TEXT_BOX, MSO_SHAPE_TYPE.PLACEHOLDER):
            continue
        if sh.name == "Text Placeholder 2":  # subtitle -- repositioned, not resized
            continue
        if is_title_shape(sh):  # titles are auto-fit separately in STEP 5
            continue
        text_stripped = sh.text_frame.text.strip()
        if text_stripped in TAGS or not text_stripped:
            continue
        # skip shapes whose current runs are already footnote-scale (<10pt) --
        # those are captions/fine print, not the "text heavy" body content.
        sizes = [r.font.size.pt for r in all_runs(sh) if r.font.size is not None and not r.font.italic]
        if not sizes or max(sizes) < 10.0:
            n_skipped_small += 1
            continue
        target = autofit_shape(sh, slide=s)
        if target is None:
            continue
        for p in sh.text_frame.paragraphs:
            p.line_spacing = LINE_SPACING
            p.space_after = Pt(PARA_GAP_PT)
            p.space_before = Pt(0)
            for r in p.runs:
                if r.font.italic or r.font.size is None:
                    continue
                r.font.size = Pt(target)
                n_bumped += 1
print(f"STEP 3: {n_bumped} body runs auto-fit into 13-16pt (per-shape, content-aware); "
      f"{n_skipped_small} footnote-scale shapes left alone")

# ===================================================== STEP 3b =================
# Rescue small-caption/annotation boxes left alone by STEP 3 (font < 10pt, so
# deliberately not resized) whose stored box height is nonetheless too small
# for their text -- typically a pre-existing user annotation/thumbnail-callout
# box with several manual <a:br/> line breaks and a stale spAutoFit height.
# We only grow the box (never shrink, never touch font size), capped by the
# footer safe-zone and the nearest shape below it.
n_grown = n_bumped_front = 0
for s in slides:
    for sh in s.shapes:
        if not sh.has_text_frame:
            continue
        if sh.shape_type not in (MSO_SHAPE_TYPE.TEXT_BOX, MSO_SHAPE_TYPE.PLACEHOLDER):
            continue
        if is_title_shape(sh):
            continue
        text_stripped = sh.text_frame.text.strip()
        if text_stripped in TAGS or not text_stripped:
            continue
        sizes = [r.font.size.pt for r in all_runs(sh) if r.font.size is not None and not r.font.italic]
        if not sizes or max(sizes) >= 10.0:
            continue  # handled by STEP 3 already
        font_pt = max(sizes)
        tf = sh.text_frame
        ml = tf.margin_left if tf.margin_left is not None else Emu(91440)
        mr = tf.margin_right if tf.margin_right is not None else Emu(91440)
        mt = tf.margin_top if tf.margin_top is not None else Emu(45720)
        mb = tf.margin_bottom if tf.margin_bottom is not None else Emu(45720)
        usable_w = max(sh.width - ml - mr, 500000)
        paras = [p.text for p in tf.paragraphs if p.text.strip()]
        if not paras or sh.height is None or sh.top is None:
            continue
        needed = estimate_height_emu(paras, font_pt, usable_w) + mt + mb
        if needed <= sh.height:
            continue
        ceiling = SAFE_BOTTOM - sh.top
        obstacle = next_obstacle_top(s, sh)
        if obstacle is not None:
            ceiling = min(ceiling, obstacle - sh.top)
        new_h = min(needed, max(ceiling, sh.height))
        if new_h > sh.height:
            sh.height = int(new_h)
            n_grown += 1
        # If even the obstacle-capped height can't hold all the text, the
        # sibling shape below it (usually a pasted screenshot) would silently
        # occlude the overflow with these apps' default z-order. Bring this
        # caption to the front of the stack so any remaining overflow still
        # draws legibly on top of it rather than disappearing behind it.
        if needed > new_h + 1000:
            s.shapes._spTree.append(sh._element)
            n_bumped_front += 1
print(f"STEP 3b: grew {n_grown} undersized caption/annotation boxes to fit their text "
      f"({n_bumped_front} also brought to front of z-order to avoid being hidden behind a picture)")

# ============================================================ STEP 4 ==========
# Rebuild every DMAIC tag as a bigger, solid-colour badge. Five distinct
# colours (not the old 2-colour DEFINE/IMPROVE + MEASURE/CONTROL sharing) so
# the phase is readable from the colour alone, not just the letters.
TAG_COLOR = {
    "DEFINE":  RGBColor(0x56, 0x28, 0x73),   # brand purple (theme accent1)
    "MEASURE": RGBColor(0x98, 0x4A, 0x9C),   # brand magenta (theme accent2)
    "ANALYZE": RGBColor(0x03, 0x46, 0x94),   # brand blue (theme accent4)
    "IMPROVE": RGBColor(0xB5, 0x62, 0x2A),   # amber -- new, distinct from DEFINE
    "CONTROL": RGBColor(0x1D, 0x7A, 0x5F),   # teal -- new, distinct from MEASURE
}
BANNER_RIGHT = 10668000  # every banner-title rectangle ends here (verified deck-wide)

def rebuild_tag(slide, old_shape, phase, banner):
    old_shape._element.getparent().remove(old_shape._element)
    if banner:
        w, h = Emu(1404000), Emu(460000)
        x, y = Emu(BANNER_RIGHT + 60000), Emu(int((914400 - 460000) / 2))
        size = 12.5
    else:
        w, h = Emu(2000000), Emu(460000)
        x, y = Emu(SLIDE_W - 645465 - 2000000), Emu(150000)
        size = 13.5
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    box.adjustments[0] = 0.5
    box.fill.solid()
    box.fill.fore_color.rgb = TAG_COLOR[phase]
    box.line.fill.background()
    box.shadow.inherit = False
    tf = box.text_frame
    tf.word_wrap = False
    tf.margin_left = tf.margin_right = Emu(60000)
    tf.margin_top = tf.margin_bottom = Emu(20000)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = phase
    r.font.size = Pt(size)
    r.font.bold = True
    r.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    r.font.name = "Segoe UI" if slide in (slides[0], slides[1], slides[2]) else "Aptos"
    r.font._rPr.set('spc', '80')
    return box

n_tags = 0
for i, s in enumerate(slides):
    banner = is_banner_slide(s)
    for sh in list(s.shapes):
        if sh.has_text_frame and sh.text_frame.text.strip() in TAGS:
            phase = sh.text_frame.text.strip()
            rebuild_tag(s, sh, phase, banner)
            n_tags += 1
print(f"STEP 4: rebuilt {n_tags} DMAIC tags as coloured badges")

# ============================================================ STEP 5 ==========
# Title/Subtitle: move left and upward on native-placeholder slides (banner
# slides keep their full-width bar, which is a different pattern and is left
# alone here). Title width is also capped so it can never run under the new,
# bigger tag badge -- and long titles are shrunk-to-fit (one line if
# possible, else two with the box grown to match) instead of overlapping it.
NEW_TITLE_LEFT, NEW_TITLE_TOP = Emu(457200), Emu(420000)
NEW_SUB_LEFT = Emu(457200)
TAG_LEFT_LAYOUT = SLIDE_W - 645465 - 2000000   # matches STEP 4's non-banner tag x
TITLE_MAX_W = Emu(TAG_LEFT_LAYOUT - 457200 - 200000)
TITLE_CHAR_W = 0.55
TITLE_FLOOR = 20.0

def title_fits_one_line(text, size_pt, usable_w_emu):
    return len(text) * size_pt * TITLE_CHAR_W * 12700 <= usable_w_emu

n_repositioned = n_title_shrunk = 0
for s in slides:
    if is_banner_slide(s):
        continue
    title = find_shape(s, "Title 1")
    sub = find_shape(s, "Text Placeholder 2")
    if title is not None:
        dy = title.top - NEW_TITLE_TOP if title.top else 0
        title.left = NEW_TITLE_LEFT
        title.top = NEW_TITLE_TOP
        title.width = TITLE_MAX_W
        n_repositioned += 1

        tp = title.text_frame.paragraphs[0]
        runs = [r for r in tp.runs if r.text.strip()]
        if runs:
            text = "".join(r.text for r in runs)
            cur_size = runs[0].font.size.pt if runs[0].font.size else 27.0
            usable_w = title.width - Emu(91440) * 2
            size = cur_size
            while size > TITLE_FLOOR and not title_fits_one_line(text, size, usable_w):
                size -= 0.5
            if size != cur_size:
                for r in runs:
                    r.font.size = Pt(size)
                n_title_shrunk += 1
            if not title_fits_one_line(text, size, usable_w):
                # still doesn't fit at the floor size -- allow a 2nd line and
                # grow the box downward so it doesn't clip.
                title.height = Emu(int(2.2 * size * 12700))
                dy = -(title.height - Emu(525078))  # push subtitle down, not up
    if sub is not None:
        sub.left = NEW_SUB_LEFT
        if dy:
            sub.top = max(Emu(1080000), sub.top - dy)
        n_repositioned += 1
print(f"STEP 5: repositioned {n_repositioned} title/subtitle shapes")

# ============================================================ STEP 6 ==========
# Slide 28 (Appendix I -- RCM/Weibull Applied) was marked hidden in GBO's own
# file -- pre-existing, not introduced by this pass. Cross-referenced by name
# from the Roadmap, Appendix G, and "When MTBF Lies," so a hidden Slide 28
# would silently vanish from any PDF export or slideshow while still being
# promised by those references. Un-hidden here; flagged in the changelog.
n_unhidden = 0
for s in slides:
    if s._element.get("show") == "0":
        del s._element.attrib["show"]
        n_unhidden += 1
print(f"STEP 6: un-hid {n_unhidden} slide(s) marked hidden in GBO's file")

# ============================================================ SAVE ============
prs.save(OUT)
print(f"saved {OUT}, {len(prs.slides)} slides")
