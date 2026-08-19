"""
fix_mtbf_rtm_numbers.py -- reconciles obsolete/mismatched RTM-### citations
in the MTBF/DMAIC deck's Appendix D/F/G/H/I reliability-modelling slides
against the current canonical RTM register (RTM-001...RTM-722 in
QPS_OFFER_Evaluation_FULL_v6.xlsx).

Finding (documented in full in the Phase 7 changelog addendum): this deck's
reliability appendices were built against an older/different RTM numbering
that predates the canonical-set correction. None of RTM-030/031/032/034/036,
RTM-047, or the range RTM-0237-0252 match their current canonical content.
Verified replacements by keyword/content search against the canonical text:

  RTM-034/035/036 (claimed: 3 failure classes)   -> RTM-061 (actual Class A/B/C definition, one item)
  RTM-05 / RTM-05a (claimed: SAE annual rate cap) -> RTM-062 (actual 90-day/12-month/5-year continuity numbers; the
                                                      SAE/yr figure is this deck's own derived conversion, not itself
                                                      a verbatim RTM -- relabelled as "derived", not "NORMATIVE")
  RTM-030/031/032 (claimed: MTBF/availability %)  -> RTM-062 (same real source as above)
  RTM-047 (claimed: VFD <=65 Hz)                  -> RTM-470 (actual VFD frequency limit, and the real canonical
                                                      figure is <=60 Hz, not 65 -- number corrected too)
  RTM-0237-0252 (claimed: CIS/MCS autonomy)       -> RTM-330/331/332 (actual CIS<->MCS real-time exchange reqs;
                                                      RTM-237-252 are unrelated helium-storage-vessel items)

BT_Method_Evaluation deck was checked too and needs no fix -- it was built
directly from the canonical data and every RTM/OFFER citation in it already
resolves correctly.
"""
import warnings
warnings.filterwarnings("ignore")
from pptx import Presentation

IN = "QPS_MTBF_WCS_DMAIC_v5.pptx"
OUT = "QPS_MTBF_WCS_DMAIC_v6.pptx"

prs = Presentation(IN)
n_fixed = 0

def set_run(slide_idx, shape_name, p_idx, r_idx, new_text, old_text_expected=None):
    global n_fixed
    slide = prs.slides[slide_idx - 1]
    for shp in slide.shapes:
        if shp.name == shape_name and shp.has_text_frame:
            p = shp.text_frame.paragraphs[p_idx]
            r = p.runs[r_idx]
            if old_text_expected is not None and r.text != old_text_expected:
                raise AssertionError(f"slide {slide_idx} {shape_name} p{p_idx}r{r_idx}: expected {old_text_expected!r}, got {r.text!r}")
            r.text = new_text
            n_fixed += 1
            return
    raise AssertionError(f"shape {shape_name!r} not found on slide {slide_idx}")

# ---------------------------------------------------------------- slide 12
set_run(12, "TextBox 9", 1, 1, "", old_text_expected="¨{")
set_run(12, "TextBox 9", 1, 2, "RTM-330/331/332 pass", old_text_expected="RTM-0237–0252 pass")
set_run(12, "TextBox 9", 1, 3, "", old_text_expected="}")

# ---------------------------------------------------------------- slide 14
set_run(14, "Text Placeholder 3", 0, 0, "RTM-330 / 331 / 332", old_text_expected="RTM 0237–0252")
set_run(14, "Text Placeholder 3", 1, 0,
        "CIS Autonomy / MCS Exchange obligation (reconciled to canonical RTM-330/331/332, MCS Integration §4.6.5).",
        old_text_expected="CIS Autonomy / MCS Exchange obligation (Addendum II verbatim extract).")
set_run(14, "Text Placeholder 3", 5, 0,
        "Fig. 9, Addendum II §3.5.4, RTM-470 (VFD target ≤ 60 Hz nominal, per IEC 60034).",
        old_text_expected="Fig. 9, Addendum II §3.5.4, RTM-047 (VFD target ≤ 65 Hz nominal).")

# ---------------------------------------------------------------- slide 19
set_run(19, "Rectangle 3", 0, 0, "RTM-061 — three consequence-based classes (Class A/B/C)",
        old_text_expected="RTM-034 to RTM-036 — three consequence-based classes")

# ---------------------------------------------------------------- slide 22
set_run(22, "Rounded Rectangle 5", 0, 1,
        "Appendix D (2/2) applies this to three specific cases; Appendix F uses it to derive a RTM-062-based 0.26 events/year limit; Appendix I applies the same logic with a Weibull (age-dependent) hazard instead of a constant one.",
        old_text_expected="Appendix D (2/2) applies this to three specific cases; Appendix F uses it to derive RTM-05's 0.26 events/year limit; Appendix I applies the same logic with a Weibull (age-dependent) hazard instead of a constant one.")

# ---------------------------------------------------------------- slide 23
set_run(23, "Rounded Rectangle 6", 0, 1,
        "Class-A target MTBF = 5 y (λ = 0.20/yr): 90-day P₀ ≈ 95% (RTM-062 compliant), 1-year ≈ 82%, 5-year ≈ 37%. This is the number RTM-",
        old_text_expected="Class-A target MTBF = 5 y (λ = 0.20/yr): 90-day P₀ ≈ 95% (RTM-031 compliant), 1-year ≈ 82%, 5-year ≈ 37%. This is the number RTM-")
set_run(23, "Rounded Rectangle 6", 0, 2, "", old_text_expected='"')
set_run(23, "Rounded Rectangle 6", 0, 3, "062 (Appendix F) is built on.",
        old_text_expected="05 (Appendix F) is built on.")
set_run(23, "Rounded Rectangle 7", 0, 1,
        "Reading “≥99% success” as a literal requirement over the full mission window (not per 90-day campaign) implies MTBF ≈ 387 years — impossible for a single-train cryoplant. This is why the RTM-062-derived cap is written as an annual event-rate limit, not a lifetime success probability.",
        old_text_expected="Reading “≥99% success” as a literal requirement over the full mission window (not per 90-day campaign) implies MTBF ≈ 387 years — impossible for a single-train cryoplant. This is why RTM-05 is written as an annual event-rate limit, not a lifetime success probability.")

# ---------------------------------------------------------------- slide 25
set_run(25, "Text Placeholder 2", 0, 0,
        "How the 90-day continuity requirement (RTM-062) becomes the annual SAE rate cap",
        old_text_expected="How the 90-day continuity requirement becomes RTM-05's annual rate cap")
set_run(25, "Rounded Rectangle 5", 0, 1,
        "Reading the success target as ≤1% failure probability over the entire multi-year mission (not per 90-day campaign) implies MTBF ≈ 387 years — see Appendix D (2/2), example 3. This RTM-062-derived cap is deliberately written as an annual rate, not a lifetime probability.",
        old_text_expected="Reading the success target as ≤1% failure probability over the entire multi-year mission (not per 90-day campaign) implies MTBF ≈ 387 years — see Appendix D (2/2), example 3. RTM-05 is deliberately written as an annual rate, not a lifetime probability.")
set_run(25, "Rounded Rectangle 6", 0, 0, "DERIVED CAP (from RTM-062)  ",
        old_text_expected="RTM-05 (NORMATIVE)  ")
set_run(25, "Rounded Rectangle 7", 0, 1,
        "Compliance is assessed with the same Poisson event-rate model (Appendix D), applied to the declared operational duty cycle. Worked check: λ ≈ 12/46.7 ≈ 0.257 SAE/yr → P(N=0, 90 d) ≈ e^(−0.257×0.2466) ≈ 0.94 — satisfies the RTM-062-derived cap.",
        old_text_expected="Compliance is assessed with the same Poisson event-rate model (Appendix D), applied to the declared operational duty cycle. Worked check: λ ≈ 12/46.7 ≈ 0.257 SAE/yr → P(N=0, 90 d) ≈ e^(−0.257×0.2466) ≈ 0.94 — satisfies RTM-05a.")

# ---------------------------------------------------------------- slide 27
set_run(27, "TextBox 7", 2, 0, "This deck's RTM-062-derived figure",
        old_text_expected="This deck's RTM-032")
set_run(27, "Rounded Rectangle 8", 0, 1,
        "The gap isn't a contradiction — it means this RTM-062-derived figure should be read either with a lower annual % target, a reduced planned-time basis, or (the best fit with what users ",
        old_text_expected="The gap isn't a contradiction — it means RTM-032 should be read either with a lower annual % target, a reduced planned-time basis, or (the best fit with what users ")

# ---------------------------------------------------------------- slide 28
set_run(28, "TextBox 4", 3, 0,
        "Single-stage MTBF = 105,000 h (Appendix G) ⇒ MTBF_train = 105,000 / 3 ≈ 35,000 h ≈ 4.0 years — the number used for Class A failure frequency and RTM-062 aggregation.",
        old_text_expected="Single-stage MTBF = 105,000 h (Appendix G) ⇒ MTBF_train = 105,000 / 3 ≈ 35,000 h ≈ 4.0 years — the number used for Class A failure frequency and RTM-030/031 aggregation.")
set_run(28, "Rounded Rectangle 5", 0, 1,
        "This is the correct number for Class A failure frequency and the RTM-062 90-day campaign probability — not the single-compressor 105,000 h figure.",
        old_text_expected="This is the correct number for Class A failure frequency, RTM-030 aggregation, and the RTM-031 90-day campaign probability — not the single-compressor 105,000 h figure.")

print(f"{n_fixed} runs fixed")
prs.save(OUT)
print(f"saved {OUT}")
