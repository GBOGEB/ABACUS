#!/usr/bin/env python3
"""DMAIC + narrative rebuild of QPS_MTBF_WCS.pptx (step1.pptx -> final.pptx)."""
from pptx import Presentation
from pptx.util import Emu, Pt
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_THEME_COLOR_INDEX
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import copy

SRC = "step1.pptx"
OUT = "final.pptx"

prs = Presentation(SRC)
slides = list(prs.slides)  # 1-indexed access via slides[i-1]
SLIDE_W = prs.slide_width
SLIDE_H = prs.slide_height

# ---------------------------------------------------------------- helpers --

def set_run(run, text=None, size=None, bold=None, italic=None, color_theme=None,
            brightness=None, rgb=None, font=None):
    if text is not None:
        run.text = text
    f = run.font
    if size is not None:
        f.size = Pt(size)
    if bold is not None:
        f.bold = bold
    if italic is not None:
        f.italic = italic
    if font is not None:
        f.name = font
    if color_theme is not None:
        f.color.theme_color = color_theme
        if brightness is not None:
            f.color.brightness = brightness
    elif rgb is not None:
        f.color.rgb = RGBColor(*rgb)


def clear_and_write(tf, lines, size=13, color_theme=MSO_THEME_COLOR_INDEX.TEXT_1,
                     brightness=0.0, bold_lead=False, align=PP_ALIGN.LEFT,
                     line_spacing=1.12, space_after=8, bullet=False, font="Calibri"):
    """lines: list[str] -> one paragraph each."""
    tf.word_wrap = True
    # first paragraph
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line_spacing
        p.space_after = Pt(space_after)
        r = p.add_run()
        set_run(r, text=line, size=size, bold=bold_lead and i == 0,
                color_theme=color_theme, brightness=brightness, font=font)


def add_takeaway(slide, x, y, cx, cy, label, text, size=12):
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, cx, cy)
    box.adjustments[0] = 0.06
    box.fill.solid()
    box.fill.fore_color.theme_color = MSO_THEME_COLOR_INDEX.ACCENT_1
    box.fill.fore_color.brightness = 0.90
    box.line.fill.background()
    box.shadow.inherit = False
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(137160)
    tf.margin_right = Emu(137160)
    tf.margin_top = Emu(80000)
    tf.margin_bottom = Emu(80000)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.line_spacing = 1.1
    r1 = p.add_run()
    set_run(r1, text=label + "  ", size=size, bold=True,
            color_theme=MSO_THEME_COLOR_INDEX.ACCENT_1)
    r2 = p.add_run()
    set_run(r2, text=text, size=size, bold=False,
            color_theme=MSO_THEME_COLOR_INDEX.ACCENT_1)
    return box


def add_narrative(slide, x, y, cx, cy, paragraphs, size=13):
    box = slide.shapes.add_textbox(x, y, cx, cy)
    tf = box.text_frame
    tf.word_wrap = True
    clear_and_write(tf, paragraphs, size=size, color_theme=MSO_THEME_COLOR_INDEX.TEXT_1,
                     brightness=0.15, line_spacing=1.15, space_after=9)
    return box


def add_tag(slide, label, style="layout"):
    """Small DMAIC phase tag. style='layout' -> top-right white strip (Text Slide
    layout, no banner). style='banner' -> right margin beside the blue/purple banner."""
    if style == "layout":
        x, y, cx, cy = Emu(9950000), Emu(160000), Emu(1750000), Emu(320000)
    else:
        x, y, cx, cy = Emu(10760000), Emu(300000), Emu(1330000), Emu(340000)
    box = slide.shapes.add_textbox(x, y, cx, cy)
    tf = box.text_frame
    tf.word_wrap = False
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    r = p.add_run()
    set_run(r, text=label, size=10.5, bold=True,
            color_theme=MSO_THEME_COLOR_INDEX.ACCENT_1, font="Calibri")
    r.font._rPr.set('spc', '120')  # slight letter spacing for a "label" feel
    return box


def bbox_of(shapes):
    xs0 = [s.left for s in shapes]
    ys0 = [s.top for s in shapes]
    xs1 = [s.left + s.width for s in shapes]
    ys1 = [s.top + s.height for s in shapes]
    return min(xs0), min(ys0), max(xs1) - min(xs0), max(ys1) - min(ys0)


def rescale_group(shapes, tx, ty, tcx, tcy):
    bx, by, bcx, bcy = bbox_of(shapes)
    scale = min(tcx / bcx, tcy / bcy)
    new_w, new_h = bcx * scale, bcy * scale
    off_x = tx + (tcx - new_w) / 2
    off_y = ty + (tcy - new_h) / 2
    for s in shapes:
        rel_x = (s.left - bx) * scale
        rel_y = (s.top - by) * scale
        s.left = Emu(int(off_x + rel_x))
        s.top = Emu(int(off_y + rel_y))
        s.width = Emu(int(s.width * scale))
        s.height = Emu(int(s.height * scale))


def find_banner(slide):
    for sh in slide.shapes:
        if sh.shape_type == 1 and sh.name.startswith("Rectangle"):
            try:
                if sh.fill.fore_color.theme_color == MSO_THEME_COLOR_INDEX.ACCENT_4:
                    return sh
            except Exception:
                pass
    return None


def set_banner(slide, new_title=None, theme=MSO_THEME_COLOR_INDEX.ACCENT_1):
    banner = find_banner(slide)
    if banner is None:
        return None
    banner.fill.fore_color.theme_color = theme
    if new_title is not None:
        # replace text preserving run formatting of first run
        tf = banner.text_frame
        p = tf.paragraphs[0]
        # remove extra runs beyond first
        for extra in p.runs[1:]:
            extra._r.getparent().remove(extra._r)
        p.runs[0].text = new_title
    return banner


def restyle_callout(shape, size=12):
    """Turn an old yellow/red 'note-to-self' box into a clean purple callout."""
    shape.fill.solid()
    shape.fill.fore_color.theme_color = MSO_THEME_COLOR_INDEX.ACCENT_1
    shape.fill.fore_color.brightness = 0.90
    shape.line.fill.background()
    shape.shadow.inherit = False
    for p in shape.text_frame.paragraphs:
        for r in p.runs:
            set_run(r, size=size, color_theme=MSO_THEME_COLOR_INDEX.ACCENT_1, bold=False)

print("helpers loaded OK")

# --------------------------------------------------------------- DMAIC tags --
# key = original (pre-insertion) slide number; value = phase label
TAGS = {
    2: "DEFINE", 3: "DEFINE", 4: "DEFINE",
    5: "MEASURE", 6: "MEASURE", 7: "MEASURE", 8: "MEASURE",
    9: "DEFINE", 10: "DEFINE", 11: "DEFINE",
    12: "MEASURE", 13: "MEASURE", 14: "MEASURE",
    15: "DEFINE", 16: "DEFINE",
    17: "ANALYZE", 18: "ANALYZE", 19: "ANALYZE", 20: "ANALYZE", 21: "ANALYZE",
    22: "ANALYZE", 23: "ANALYZE", 24: "ANALYZE", 25: "ANALYZE", 26: "ANALYZE",
    27: "IMPROVE", 28: "IMPROVE",
    29: "ANALYZE", 30: "ANALYZE", 31: "ANALYZE",
}

for orig_num, label in TAGS.items():
    new_pos = orig_num + 1  # +1 because the Roadmap slide was inserted at position 2
    slide = slides[new_pos - 1]
    style = "layout" if orig_num <= 14 else "banner"
    add_tag(slide, label, style=style)

print("Phase tags added to", len(TAGS), "baseline slides")

# ------------------------------------------------------------- Roadmap slide --
roadmap = slides[1]
ph_title = roadmap.shapes[0]      # idx 0 title
ph_sub = roadmap.shapes[1]        # idx 14 subtitle
ph_body = roadmap.shapes[2]       # idx 15 body

ph_title.text_frame.paragraphs[0].runs[0].text = "DMAIC Roadmap"

sub_p = ph_sub.text_frame.paragraphs[0]
sub_p.runs[0].text = "How to read this master deck — compendium structure for full-system reliability (QPLANT / WCS)"

# give the body placeholder a bit more room (footer starts ~6.48M EMU)
ph_body.left, ph_body.top = Emu(766763), Emu(2120000)
ph_body.width, ph_body.height = Emu(10607253), Emu(4300000)

roadmap_items = [
    ("DEFINE", "Slides 3–5, 10–12, 16–17",
     "What the system must achieve, and by what rule — scope, RTM obligations, the operating-state hierarchy, and failure-class definitions with their MTBF targets."),
    ("MEASURE", "Slides 6–9, 13–15",
     "What the fleet can actually deliver — OEM flow data, N−1 capability, acceptance-test criteria, and liquid-buffer coverage time."),
    ("ANALYZE", "Slides 18–32",
     "Why the numbers behave the way they do — the Poisson/exponential campaign-success mathematics, the 94–95% ceiling, the architecture trade-offs, and the Weibull wear-out mechanism that limits any MTBF-only argument."),
    ("IMPROVE", "Slides 33–38",
     "What changes as a result — why MTBF matters, when it stops being trustworthy, and the component evidence (compressors, PVPS, turbines) that proves it."),
    ("CONTROL", "Slides 39–42",
     "How the plant stays reliable going forward — predictable, usage-based service replacement; the reset-to-new renewal principle; and the governance loop that keeps the system inside the assumptions MTBF requires."),
]

tf = ph_body.text_frame
tf.word_wrap = True
for i, (phase, rng, desc) in enumerate(roadmap_items):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    p.line_spacing = 1.08
    p.space_after = Pt(11)
    r1 = p.add_run()
    set_run(r1, text=f"{phase}", size=15, bold=True, color_theme=MSO_THEME_COLOR_INDEX.ACCENT_1)
    r2 = p.add_run()
    set_run(r2, text=f"  ({rng})  — ", size=13, italic=True,
            color_theme=MSO_THEME_COLOR_INDEX.ACCENT_2)
    r3 = p.add_run()
    set_run(r3, text=desc, size=13, color_theme=MSO_THEME_COLOR_INDEX.TEXT_1, brightness=0.1)

# closing modularity note
p = tf.add_paragraph()
p.space_before = Pt(6)
p.line_spacing = 1.08
r = p.add_run()
set_run(r, text="This is a MASTER slide set: each phase above is self-contained and can be lifted "
              "into its own focused sub-deck (e.g. an ANALYZE-only briefing or a CONTROL-only "
              "maintenance-policy pack) without losing narrative continuity.",
        size=11.5, italic=True, color_theme=MSO_THEME_COLOR_INDEX.TEXT_1, brightness=0.35)

print("Roadmap slide populated")

# ------------------------------------------------------ Narrative rebuild (33-42) --
RH_X, RH_Y, RH_CX, RH_CY = Emu(6450000), Emu(1050000), Emu(5300000), Emu(5350000)
LT_X, LT_CX = Emu(550000), Emu(5650000)
LT_NARR_Y, LT_NARR_CY = Emu(1050000), Emu(3550000)
LT_TAKE_Y, LT_TAKE_CY = Emu(4700000), Emu(1500000)

def build_right_half_slide(pos, title, narrative, takeaway_label, takeaway_text,
                            evidence_names, restyle_textbox_names=()):
    slide = slides[pos - 1]
    set_banner(slide, title)
    evid = [sh for sh in slide.shapes if sh.name in evidence_names]
    rescale_group(evid, RH_X, RH_Y, RH_CX, RH_CY)
    for sh in slide.shapes:
        if sh.name in restyle_textbox_names:
            restyle_callout(sh, size=11)
    add_narrative(slide, LT_X, LT_NARR_Y, LT_CX, LT_NARR_CY, narrative, size=13)
    add_takeaway(slide, LT_X, LT_TAKE_Y, LT_CX, LT_TAKE_CY, takeaway_label, takeaway_text, size=12)
    add_tag(slide, "IMPROVE" if pos <= 38 else "CONTROL", style="banner")
    return slide

# ---- Slide 33 (orig 32): Why MTBF Matters -------------------------------------
s = slides[32]
set_banner(s, "Why MTBF Matters")
lead = [sh for sh in s.shapes if sh.name == "Text Placeholder 2"][0]
tf = lead.text_frame
tf.word_wrap = True
p0 = tf.paragraphs[0]
for extra in p0.runs[1:]:
    extra._r.getparent().remove(extra._r)
set_run(p0.runs[0],
        text="Every number in this deck — the 5-year Class A target, the N−1 margins, the "
             "94–95% campaign ceiling — reduces to one shared quantity: Mean Time Between Failures.",
        size=13, bold=False, color_theme=MSO_THEME_COLOR_INDEX.ACCENT_2)
p1 = tf.add_paragraph()
p1.line_spacing = 1.1
p1.space_before = Pt(4)
# copy the no-bullet pPr from paragraph 0 so this line doesn't inherit a layout bullet
import copy as _copy
src_pPr = p0._p.find('{http://schemas.openxmlformats.org/drawingml/2006/main}pPr')
if src_pPr is not None:
    p1._p.insert(0, _copy.deepcopy(src_pPr))
r = p1.add_run()
set_run(r, text="The BCR013 derivation opposite shows that budget in its rawest form: the "
               "accelerator's own 250 h MTBF caps how much unreliability QCELL/QPS/NF may add.",
        size=13, color_theme=MSO_THEME_COLOR_INDEX.ACCENT_2)
yellow = [sh for sh in s.shapes if sh.name == "TextBox 16"][0]
restyle_callout(yellow, size=11.5)
ytf = yellow.text_frame
yp = ytf.paragraphs[0]
for extra in yp.runs[1:]:
    extra._r.getparent().remove(extra._r)
r0 = yp.runs[0]
set_run(r0, text="KEY TAKEAWAY  ", size=11.5, bold=True, color_theme=MSO_THEME_COLOR_INDEX.ACCENT_1)
r1 = yp.add_run()
set_run(r1, text="MTBF is a shared 0.35-events/year budget, not a component property — every "
               "claim later in this deck draws on the same account.",
        size=11.5, bold=False, color_theme=MSO_THEME_COLOR_INDEX.ACCENT_1)
add_tag(s, "IMPROVE", style="banner")
print("slide 33 done")

# ---- Slide 34 (orig 33): The Hidden Assumption --------------------------------
build_right_half_slide(
    34, "The Hidden Assumption",
    narrative=[
        "MTBF only means what we want it to mean if failures arrive at random, at a "
        "constant rate — the Poisson/exponential assumption used throughout the ANALYZE "
        "section (Slides 18–32). Under that assumption, a component with MTBF = 5 years "
        "is exactly as likely to fail in month one as in month fifty.",
        "That is a convenient mathematical fiction, not a description of real machinery. "
        "It holds well for random, external trips — it says nothing about wear. The chart "
        "opposite is the same Poisson relationship, one more time, as a reminder that "
        "everything to its left was built on a flat hazard rate.",
    ],
    takeaway_label="KEY TAKEAWAY",
    takeaway_text="A single MTBF number silently assumes “no aging.” The next slides show why that assumption has an expiry date.",
    evidence_names={"Picture 4", "Picture 6"},
)
print("slide 34 done")

# ---- Slide 36 (orig 35, Kaeser): Case File - HP Compressors -------------------
build_right_half_slide(
    36, "Case File: HP Compressors",
    narrative=[
        "The HP screw compressors carry an assumed 350,000 h MTBF — long enough that, under "
        "the Poisson model, a bare 3-of-3 no-redundancy train still clears 90 days with 98.2% "
        "probability. That number is comforting only while the compressors stay inside their "
        "flat-hazard useful-life window.",
        "Screw compressors wear through bearings, seals and rotor packages on a calendar/duty-"
        "cycle clock that the Poisson model on the page knows nothing about. The 3+1 "
        "redundancy case (right) does not just buy short-window margin — it buys the time "
        "needed to pull a unit for scheduled service before it enters its own wear-out region.",
    ],
    takeaway_label="KEY TAKEAWAY",
    takeaway_text="Redundancy here is not only insurance against random failure — it is the maintenance window that makes predictable replacement possible.",
    evidence_names={"Picture 4"},
)
print("slide 36 done")

# ---- Slide 37 (orig 36, PVPS): Case File - PVPS -------------------------------
build_right_half_slide(
    37, "Case File: PVPS",
    narrative=[
        "PVPS multiplies the same logic across nine parallel units. At 9-of-9 with no spare, "
        "the 5-year probability of at least one trip is 67.6% — not because any single pump is "
        "unreliable, but because unreliability compounds as fleet size grows.",
        "N+1 and N+2 sparing claw back the 90-day and 1-year numbers dramatically, but the "
        "5-year figure stays stubborn (13–33%): over long horizons, “more components” "
        "re-introduces the wear-out problem this section is about. Sparing buys statistical "
        "margin against random failure — it does not buy the fleet time out of its wear-out clock.",
    ],
    takeaway_label="KEY TAKEAWAY",
    takeaway_text="For a 9-unit fleet, the honest 5-year reliability number is a scheduling problem, not just a redundancy problem.",
    evidence_names={"Picture 5"},
)
print("slide 37 done")

# ---- Slide 38 (orig 37, Turbines pt.1): Case File - Turbines ------------------
build_right_half_slide(
    38, "Case File: Turbines — Nowhere to Hide",
    narrative=[
        "Turbines are the clearest case in this compendium for why MTBF-only thinking runs "
        "out. There is no credible spare position: cooling topology fixes each turbine's "
        "location in the cold-box train, so a “redundant” unit is not interchangeable "
        "the way a compressor is.",
        "With 7-of-7 required and a 105,000 h assumed MTBF, the 5-year no-trip probability is "
        "already down to 94.6%; adding an 8th unit barely moves it. A fleet this exposed "
        "cannot be managed by waiting for an MTBF-driven alarm — it has to be managed by "
        "knowing, ahead of time, which unit is closest to end of usable life.",
    ],
    takeaway_label="KEY TAKEAWAY",
    takeaway_text="No spare position + no redundancy = the fleet where predictable, age-based replacement matters most, not least.",
    evidence_names={"Picture 5", "TextBox 6"},
    restyle_textbox_names={"TextBox 6"},
)
print("slide 38 done")

# ---- Slide 39 (orig 38, Turbines pt.2): The Turn ------------------------------
build_right_half_slide(
    39, "The Turn: From Predicting Failure to Preventing It",
    narrative=[
        "Every slide so far has used MTBF exactly as designed: to predict how often a "
        "well-behaved, randomly-failing fleet will trip. That is the right tool for "
        "budget-setting, architecture trade-offs, and acceptance testing.",
        "It is the wrong tool for one aging turbine, compressor or PVPS unit that has "
        "quietly moved from the flat part of the bathtub curve onto the rising one. Once "
        "wear-out dominates — as the three case files above all confirm it eventually does — "
        "the right question stops being “what's the MTBF” and becomes “how much "
        "usable life is left, and when do we act.” That is the shift from reliability "
        "prediction to reliability control.",
    ],
    takeaway_label="KEY TAKEAWAY",
    takeaway_text="MTBF tells you the odds. It does not tell you when to intervene — the remaining slides answer that question.",
    evidence_names={"Picture 5", "Picture 4", "TextBox 6", "Straight Arrow Connector 8"},
    restyle_textbox_names={"TextBox 6"},
)
print("slide 39 done")

# ---- Slide 40 (orig 39): Predictable Service Replacement ----------------------
build_right_half_slide(
    40, "Predictable Service Replacement",
    narrative=[
        "The control response to a rising hazard rate is not a better MTBF estimate — it is "
        "a scheduled intervention set before the wear-out region is reached, sized from each "
        "component's Weibull characteristic life (η) rather than its long-run MTBF.",
        "In practice: set overhaul/replacement intervals at a fraction of η (tightened "
        "further for Class B/C consequence items); track operating hours and Hz-band exposure "
        "per unit — data this deck already mandates the Applicant to log; and treat “still "
        "running at 340,000 hours” as a maintenance trigger, not a reliability success story.",
    ],
    takeaway_label="KEY TAKEAWAY",
    takeaway_text="A predictable replacement interval is a Weibull-based deadline set before the hazard curve turns upward — not a reaction to it.",
    evidence_names={"Picture 5", "Picture 6"},
)
print("slide 40 done")

# ---- Slide 35 (orig 34): When MTBF Lies (wear-out boundary) -------------------
s = slides[34]
set_banner(s, "When MTBF Lies — the Wear-Out Boundary")
add_narrative(
    s, Emu(580000), Emu(4000000), Emu(11000000), Emu(1350000),
    [
        "Slide 32's RCM/Weibull relation, h(t) = (β/η)·(t/η)^(β−1), makes the limit explicit: "
        "whenever β > 1, hazard rate rises with age, and a single constant MTBF understates "
        "real short-term risk near end-of-life. The worked example already in this deck shows "
        "the effect is not small: for a β = 2.5 asset (oil screw / cold compressor), the chance "
        "of failing in the next 90 days after 5 years in service is 200–240× higher than the "
        "same 90 days right after a good-as-new overhaul.",
    ], size=13,
)
add_takeaway(
    s, Emu(580000), Emu(5480000), Emu(11000000), Emu(760000),
    "KEY TAKEAWAY",
    "MTBF governs the plant's average reliability budget. Weibull shape (β) governs whether any one unit is quietly approaching a cliff.",
    size=12,
)
add_tag(s, "IMPROVE", style="banner")
print("slide 35 done")

# ---- Slide 41 (orig 40): Reset to New -----------------------------------------
s = slides[40]
set_banner(s, "Reset to New")
add_narrative(
    s, Emu(580000), Emu(970000), Emu(11000000), Emu(780000),
    [
        "A full overhaul or unit exchange does more than avoid one failure — it resets the "
        "clock. A good-as-new replacement returns the component to t = 0 on its Weibull curve, "
        "restoring the low, flat hazard rate that every Poisson/MTBF probability in this "
        "compendium (Slides 18–32) depends on.",
    ], size=12.5,
)
add_takeaway(
    s, Emu(580000), Emu(5680000), Emu(11000000), Emu(700000),
    "KEY TAKEAWAY",
    "“Reset to new” is what keeps the rest of this deck's mathematics honest — skip it, and every MTBF-based probability quoted earlier quietly becomes optimistic.",
    size=11.5,
)
add_tag(s, "CONTROL", style="banner")
print("slide 41 done")

# ---- Slide 42 (orig 41, duplicate charts): Reliability Governance -------------
s = slides[41]
set_banner(s, "Reliability Governance — Closing the Loop")
pics = [sh for sh in s.shapes if sh.name in ("Picture 4", "Picture 8")]
pics.sort(key=lambda sh: sh.left)
keep, drop = pics[0], pics[1]
drop._element.getparent().remove(drop._element)
rescale_group([keep], Emu(580000), Emu(1750000), Emu(5300000), Emu(4550000))

add_narrative(
    s, Emu(580000), Emu(970000), Emu(11000000), Emu(680000),
    [
        "This compendium opened with a budget (Slide 33: 0.35 events/year) and closes with "
        "the mechanism that keeps the plant inside it: a continuous RAMI governance loop.",
    ], size=12.5,
)

loop_x, loop_y, loop_cx = Emu(6250000), Emu(1830000), Emu(5300000)
steps = [
    ("1", "MEASURE exposure", "Log operating hours and Hz-band time per unit (already mandated, Slide 5)."),
    ("2", "TRACK age vs. η", "Compare accumulated life to each component's Weibull characteristic life."),
    ("3", "REPLACE before wear-out", "Trigger overhaul/exchange at a fraction of η — not at first symptom."),
    ("4", "RESET & feed back", "Log the renewal and roll it into next year's RAMI dossier (Slide 10)."),
]
step_h = 850000
for i, (num, head, body) in enumerate(steps):
    y = loop_y + i * (step_h + 60000)
    chip = s.shapes.add_shape(MSO_SHAPE.OVAL, loop_x, y, Emu(360000), Emu(360000))
    chip.fill.solid()
    chip.fill.fore_color.theme_color = MSO_THEME_COLOR_INDEX.ACCENT_1
    chip.line.fill.background()
    chip.shadow.inherit = False
    ctf = chip.text_frame
    ctf.margin_left = ctf.margin_right = ctf.margin_top = ctf.margin_bottom = 0
    ctf.vertical_anchor = MSO_ANCHOR.MIDDLE
    cp = ctf.paragraphs[0]
    cp.alignment = PP_ALIGN.CENTER
    cr = cp.add_run()
    set_run(cr, text=num, size=14, bold=True, color_theme=MSO_THEME_COLOR_INDEX.BACKGROUND_1)

    txt = s.shapes.add_textbox(loop_x + Emu(460000), y - Emu(40000), loop_cx - Emu(460000), Emu(step_h))
    ttf = txt.text_frame
    ttf.word_wrap = True
    tp = ttf.paragraphs[0]
    tp.line_spacing = 1.05
    tr = tp.add_run()
    set_run(tr, text=head, size=13, bold=True, color_theme=MSO_THEME_COLOR_INDEX.ACCENT_1)
    tp2 = ttf.add_paragraph()
    tp2.line_spacing = 1.05
    tp2.space_before = Pt(2)
    tr2 = tp2.add_run()
    set_run(tr2, text=body, size=11.5, color_theme=MSO_THEME_COLOR_INDEX.TEXT_1, brightness=0.15)

add_takeaway(
    s, Emu(580000), Emu(5760000), Emu(11000000), Emu(560000),
    "KEY TAKEAWAY",
    "MTBF opened the story; predictable, logged, reset-to-new replacement is how the plant keeps living up to it.",
    size=11.5,
)
add_tag(s, "CONTROL", style="banner")
print("slide 42 done")

prs.save(OUT)
print("SAVED", OUT)
