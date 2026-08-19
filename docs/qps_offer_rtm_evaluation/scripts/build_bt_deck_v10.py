"""
build_bt_deck_v10.py -- GBO: "plus additional slide(s) to appropriate deck to
stipulate relevance and importants" for the 6 top-level / 26 sub-item
contract Deliverables Dossier (Applicable Documentation, "AD_##") groups.

Appends ONE new slide (24) at the end of BT_Method_Evaluation_v9.pptx, badge
EVAL-S10 -- following the same append-only, no-renumbering convention used
for every prior deck addition (v6/v8/v9), since every cross-reference in
this deck uses METHOD-Sxx/EVAL-Sxx codes rather than hardcoded slide
numbers.

Content is pulled directly from DELIVERABLES_DOSSIER (via nav_data_v22.json,
itself exported straight from the workbook) -- the same 32-entry structure
and per-group linked-RTM counts already shown on the Navigator's Taxonomy
tab (flat table + new tree view this round). Nothing here is estimated.

IN:  BT_Method_Evaluation_v9.pptx
OUT: BT_Method_Evaluation_v10.pptx
"""
import warnings, json
warnings.filterwarnings("ignore")
from pptx import Presentation
from pptx.util import Emu, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_AUTO_SIZE

IN = "BT_Method_Evaluation_v9.pptx"
OUT = "BT_Method_Evaluation_v10.pptx"
NAV_DATA = "/tmp/nav_data_v22.json"

HEADER = RGBColor(0x44, 0x1F, 0x63)
FOOTER_BG = RGBColor(0xF3, 0xEC, 0xF8)
CALLOUT_BG = RGBColor(0xEF, 0xE5, 0xF5)
ACCENT = RGBColor(0x0B, 0x3D, 0x5C)
BODY = RGBColor(0x22, 0x22, 0x22)
FOOTER_TXT = RGBColor(0x55, 0x55, 0x55)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
CODE_TXT = RGBColor(0xBF, 0xD6, 0xE6)
ROW_ALT = RGBColor(0xF4, 0xF7, 0xF9)

SLIDE_W = Emu(12192000)
SLIDE_H = Emu(6858000)
HEADER_H = Emu(868680)
FOOTER_TOP = Emu(6547104)
LEFT_MARGIN = Emu(320040)
CONTENT_RIGHT = SLIDE_W - LEFT_MARGIN

data = json.load(open(NAV_DATA))
items = data["deliverablesDossier"]
note = data["deliverablesDossierNote"]
tops = [x for x in items if x["isTopLevel"]]
rows = []
for top in tops:
    kids = [x for x in items if not x["isTopLevel"] and x["id"].startswith(top["id"] + ".")]
    all_rtms = set(top["linkedRtmIds"])
    for k in kids:
        all_rtms.update(k["linkedRtmIds"])
    rows.append((top["id"], top["name"], len(kids), len(all_rtms)))

prs = Presentation(IN)
blank_layout = prs.slide_layouts[6]
assert blank_layout.name == "Blank"


def add_header_footer(slide, title, code, footer_note):
    rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Emu(0), Emu(0), SLIDE_W, HEADER_H)
    rect.fill.solid(); rect.fill.fore_color.rgb = HEADER
    rect.line.fill.background(); rect.shadow.inherit = False

    tb = slide.shapes.add_textbox(Emu(411480), Emu(109728), Emu(9875520), Emu(685800))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]
    r = p.add_run(); r.text = title
    r.font.size = Pt(25); r.font.bold = True; r.font.color.rgb = WHITE; r.font.name = "Aptos"

    cb = slide.shapes.add_textbox(Emu(10607040), Emu(292608), Emu(1371600), Emu(320040))
    ctf = cb.text_frame
    cp = ctf.paragraphs[0]; cp.alignment = PP_ALIGN.RIGHT
    cr = cp.add_run(); cr.text = code
    cr.font.size = Pt(11); cr.font.bold = True; cr.font.color.rgb = CODE_TXT; cr.font.name = "Aptos"

    frect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Emu(0), FOOTER_TOP, SLIDE_W, Emu(310896))
    frect.fill.solid(); frect.fill.fore_color.rgb = FOOTER_BG
    frect.line.fill.background(); frect.shadow.inherit = False

    fb = slide.shapes.add_textbox(Emu(411480), Emu(6565392), Emu(10515600), Emu(274320))
    ftf = fb.text_frame
    fp = ftf.paragraphs[0]
    fr = fp.add_run(); fr.text = footer_note
    fr.font.size = Pt(8.5); fr.font.color.rgb = FOOTER_TXT; fr.font.name = "Aptos"
    return rect, frect


slide = prs.slides.add_slide(blank_layout)
add_header_footer(
    slide,
    "Deliverables Dossier — Applicable Documentation",
    "EVAL-S10",
    "BT METHOD — Engineering Baseline (canonical RTM-722) · EVAL-S10 (Deliverables Dossier / AD relevance, new) | v3.8",
)

sub = slide.shapes.add_textbox(Emu(411480), Emu(940000), Emu(11400000), Emu(280000))
stf = sub.text_frame; stf.word_wrap = True
sp = stf.paragraphs[0]
sr = sp.add_run()
sr.text = "Contract Table 2: 32 named document bundles (6 groups + 26 sub-items), matched against every RTM's text for an AD_## citation."
sr.font.size = Pt(12.5); sr.font.color.rgb = BODY; sr.font.name = "Aptos"

# ---- left: table of 6 top-level groups ----
tbl_left = Emu(411480); tbl_top = Emu(1440000)
tbl_w = Emu(6900000); tbl_h = Emu(2500000)
gtable = slide.shapes.add_table(len(rows) + 1, 4, tbl_left, tbl_top, tbl_w, tbl_h).table
gtable.columns[0].width = Emu(900000)
gtable.columns[1].width = Emu(3900000)
gtable.columns[2].width = Emu(1050000)
gtable.columns[3].width = Emu(1050000)
headers = ["AD ID", "Document group", "Sub-items", "RTMs citing"]
for ci, h in enumerate(headers):
    cell = gtable.cell(0, ci)
    cell.fill.solid(); cell.fill.fore_color.rgb = ACCENT
    cell.text_frame.paragraphs[0].add_run()
    r = cell.text_frame.paragraphs[0].runs[0]
    r.text = h; r.font.size = Pt(10); r.font.bold = True; r.font.color.rgb = WHITE; r.font.name = "Aptos"
for ri, (aid, name, nkids, nrtm) in enumerate(rows, start=1):
    vals = [aid, name, str(nkids) if nkids else "—", str(nrtm) if nrtm else "0"]
    for ci, v in enumerate(vals):
        cell = gtable.cell(ri, ci)
        cell.fill.solid(); cell.fill.fore_color.rgb = ROW_ALT if ri % 2 == 0 else WHITE
        cell.text_frame.paragraphs[0].add_run()
        r = cell.text_frame.paragraphs[0].runs[0]
        r.text = v; r.font.size = Pt(9.5); r.font.color.rgb = BODY; r.font.name = "Aptos"
        if ci == 0:
            r.font.bold = True; r.font.color.rgb = ACCENT

# ---- right: relevance/importance callout ----
cb = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Emu(7498079), Emu(1440000), Emu(4343400), Emu(2500000))
cb.fill.solid(); cb.fill.fore_color.rgb = CALLOUT_BG
cb.line.color.rgb = ACCENT; cb.line.width = Pt(1)
cb.shadow.inherit = False
ctf = cb.text_frame; ctf.word_wrap = True
ctf.auto_size = MSO_AUTO_SIZE.NONE
ctf.margin_left = Emu(137160); ctf.margin_top = Emu(90000); ctf.margin_right = Emu(137160); ctf.margin_bottom = Emu(60000)
lines = [
    ("Why these matter: ", "the contract's own named execution deliverables (drawings, interface catalogues, H&S plan, 3D model, CAD conventions) — separate from the RTM shall-statement, required regardless of which RTM text names them.", True),
    ("Coverage today: ", "9 / 32 entries (28%) have at least one RTM explicitly citing the AD_## code — 13 distinct RTMs, concentrated in AD_05 (Controls/Interlocks/IT, 7 of 7 sub-items cited).", True),
    ("Read with care: ", "a low citation count is NOT low relevance — most RTMs describe the same deliverable in plain contract-section language without naming the AD code.", True),
]
first = True
for lead, text, bold in lines:
    p = ctf.paragraphs[0] if first else ctf.add_paragraph()
    first = False
    p.space_after = Pt(5)
    r0 = p.add_run(); r0.text = lead
    r0.font.bold = True; r0.font.size = Pt(10.5); r0.font.color.rgb = ACCENT; r0.font.name = "Aptos"
    r1 = p.add_run(); r1.text = text
    r1.font.size = Pt(10.5); r1.font.color.rgb = BODY; r1.font.name = "Aptos"

# ---- bottom: open items ----
ob = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, LEFT_MARGIN, Emu(4150000), CONTENT_RIGHT - LEFT_MARGIN, Emu(2150000))
ob.fill.solid(); ob.fill.fore_color.rgb = CALLOUT_BG
ob.line.fill.background(); ob.shadow.inherit = False
otf = ob.text_frame; otf.word_wrap = True
otf.margin_left = Emu(137160); otf.margin_top = Emu(80000); otf.margin_right = Emu(137160)
bullets = [
    "Navigator (Taxonomy tab, this round): the flat 32-row table now has a companion collapsible tree (6 top-level groups, real sub-items nested, linked RTMs clickable through to RTM Lookup).",
    "Open: the actual AD_01–AD_06 source documents (drawings, plans, catalogues) have not been ingested into this evaluation — only their names/reference numbers from the contract's Table 2. GBO has offered a local folder with these; once accessible, relevance/importance per document can move from name-only to content-based.",
    "This ties directly to roadmap item B1 (granular Applicant-specific OFFER documentation as a single bookmarked PDF) — the same ingestion gap applies to both the Dossier documents and the Applicants' own OFFER submissions.",
]
first = True
for b in bullets:
    p = otf.paragraphs[0] if first else otf.add_paragraph()
    first = False
    p.space_after = Pt(6)
    r0 = p.add_run(); r0.text = "•  "
    r0.font.bold = True; r0.font.size = Pt(11); r0.font.color.rgb = ACCENT; r0.font.name = "Aptos"
    r1 = p.add_run(); r1.text = b
    r1.font.size = Pt(11); r1.font.color.rgb = BODY; r1.font.name = "Aptos"

prs.save(OUT)
print(f"wrote {OUT}: {len(prs.slides._sldIdLst)} slides total (added slide 24, EVAL-S10)")
print(f"deliverables dossier note used: {note[:90]}...")
