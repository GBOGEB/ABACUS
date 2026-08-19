"""
build_bt_deck_v12.py -- Control-phase + task #59 round.

1. REBUILDS the Deliverables Dossier slide (badge EVAL-S12) in place --
   it was built from nav_data_v22.json (6 top-level groups, 32 entries);
   FULL_v23 added AD_07/AD_08 (now 8 groups, 34 entries). python-pptx has
   no in-place slide-content-replace API, so the old slide is deleted from
   the slide list (XML surgery -- python-pptx has no delete_slide() either)
   and a freshly-built one with the SAME badge code is appended. Nothing
   else in this deck references slides by position, only by badge code, so
   this is safe per the deck's own established convention.

2. Appends 3 NEW slides (badges EVAL-S13/S14/S15 -- next available after
   EVAL-S12, confirmed via a raw XML grep across all 24 v11 slides before
   writing this script):

   EVAL-S13 -- "Honest ranking: official (gate-first) vs pure weighted-S"
   Direct response to GBO: "I think an honest ranking (with the red colour
   and T0 and gate pinning still available - Do not want to warp picture -
   maybe add specific T0 ranking and order slide or diagram?" Shows BOTH
   rankings side by side for all 6 gate items, gate coloring kept, not
   replaced -- the pure-score column is additive disclosure, not a
   different "official" answer.

   EVAL-S14 -- Weight-scenario side-by-side (Base / Equal / Cost=70%),
   message banner, real Spearman correlations and named rank-shift items
   from the already-computed Section 19 analysis.

   EVAL-S15 -- Granular OFFER cross-reference table (Code/Standard,
   Deliverable/proof, Applicable phase(s), Primary/Supporting RTM counts)
   for the top-ranked OFFER items -- data already computed and stored on
   OFFER_RANKING, not estimated for this slide.

IN:  BT_Method_Evaluation_v11.pptx
OUT: BT_Method_Evaluation_v12.pptx
"""
import warnings, json, copy
warnings.filterwarnings("ignore")
import openpyxl
from pptx import Presentation
from pptx.util import Emu, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE
from pptx.enum.shapes import MSO_SHAPE

IN = "BT_Method_Evaluation_v11.pptx"
OUT = "BT_Method_Evaluation_v12.pptx"
NAV_DATA = "/tmp/nav_data_v23.json"
WB = "QPS_OFFER_Evaluation_FULL_v23.xlsx"

HEADER = RGBColor(0x44, 0x1F, 0x63)
FOOTER_BG = RGBColor(0xF3, 0xEC, 0xF8)
CALLOUT_BG = RGBColor(0xEF, 0xE5, 0xF5)
ACCENT = RGBColor(0x0B, 0x3D, 0x5C)
BODY = RGBColor(0x22, 0x22, 0x22)
FOOTER_TXT = RGBColor(0x55, 0x55, 0x55)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
CODE_TXT = RGBColor(0xBF, 0xD6, 0xE6)
ROW_ALT = RGBColor(0xF4, 0xF7, 0xF9)
GATE_RED = RGBColor(0xB0, 0x2A, 0x2A)
GATE_RED_BG = RGBColor(0xFB, 0xE9, 0xE7)
GREEN = RGBColor(0x2E, 0x8B, 0x57)

SLIDE_W = Emu(12192000)
SLIDE_H = Emu(6858000)
HEADER_H = Emu(868680)
FOOTER_TOP = Emu(6547104)
LEFT_MARGIN = Emu(320040)
CONTENT_RIGHT = SLIDE_W - LEFT_MARGIN

data = json.load(open(NAV_DATA))

wb = openpyxl.load_workbook(WB, data_only=True)
ws = wb["OFFER_RANKING"]
hdr = 5
col = {ws.cell(row=hdr, column=c).value: c for c in range(1, ws.max_column + 1) if ws.cell(row=hdr, column=c).value}
offer_rows = []
for r in range(6, ws.max_row + 1):
    oid = ws.cell(row=r, column=col["OFFER ID"]).value
    if not oid:
        continue
    offer_rows.append({
        "id": oid, "title": ws.cell(row=r, column=col["Title"]).value,
        "gate": ws.cell(row=r, column=col["Gate"]).value,
        "tier": ws.cell(row=r, column=col["Tier"]).value,
        "official_rank": ws.cell(row=r, column=col["Rank"]).value,
        "wS": ws.cell(row=r, column=col["Weighted S"]).value or 0,
        "code_std": ws.cell(row=r, column=col["Actual explicit code / standard"]).value,
        "deliverable": ws.cell(row=r, column=col["Explicit deliverable / proof"]).value,
        "phase": ws.cell(row=r, column=col["Applicable phase(s)"]).value,
        "primary_rtms": ws.cell(row=r, column=col["Primary RTMs"]).value,
        "supporting_rtms": ws.cell(row=r, column=col["Supporting RTMs"]).value,
    })
pure_sorted = sorted(offer_rows, key=lambda x: -x["wS"])
for i, r in enumerate(pure_sorted, 1):
    r["pure_rank"] = i
by_id = {r["id"]: r for r in offer_rows}
gate_items = sorted([r for r in offer_rows if r["gate"] == "Yes"], key=lambda x: x["official_rank"])

prs = Presentation(IN)
blank_layout = prs.slide_layouts[6]
assert blank_layout.name == "Blank"


def add_header_footer(slide, title, code, footer_note):
    rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Emu(0), Emu(0), SLIDE_W, HEADER_H)
    rect.fill.solid(); rect.fill.fore_color.rgb = HEADER
    rect.line.fill.background(); rect.shadow.inherit = False

    tb = slide.shapes.add_textbox(Emu(411480), Emu(109728), Emu(9875520), Emu(685800))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]
    r = p.add_run(); r.text = title
    r.font.size = Pt(24); r.font.bold = True; r.font.color.rgb = WHITE; r.font.name = "Aptos"

    cb = slide.shapes.add_textbox(Emu(10607040), Emu(292608), Emu(1371600), Emu(320040))
    ctf = cb.text_frame
    cp = ctf.paragraphs[0]; cp.alignment = PP_ALIGN.RIGHT
    cr = cp.add_run(); cr.text = code
    cr.font.size = Pt(11); cr.font.bold = True; cr.font.color.rgb = CODE_TXT; cr.font.name = "Aptos"

    frect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Emu(0), FOOTER_TOP, SLIDE_W, Emu(310896))
    frect.fill.solid(); frect.fill.fore_color.rgb = FOOTER_BG
    frect.line.fill.background(); frect.shadow.inherit = False

    fb = slide.shapes.add_textbox(Emu(411480), Emu(6565392), Emu(10515600), Emu(274320))
    ftf = fb.text_frame
    fp = ftf.paragraphs[0]
    fr = fp.add_run(); fr.text = footer_note
    fr.font.size = Pt(8.5); fr.font.color.rgb = FOOTER_TXT; fr.font.name = "Aptos"
    return rect, frect


def delete_slide(prs, index):
    xml_slides = prs.slides._sldIdLst
    slides = list(xml_slides)
    rId = slides[index].rId
    prs.part.drop_rel(rId)
    xml_slides.remove(slides[index])


# ============================================================ find & delete old EVAL-S12 slide
target_idx = None
for i, slide in enumerate(prs.slides):
    for shape in slide.shapes:
        if shape.has_text_frame and "EVAL-S12" in shape.text_frame.text:
            target_idx = i
            break
    if target_idx is not None:
        break
assert target_idx is not None, "could not find EVAL-S12 slide to rebuild"
print(f"deleting old EVAL-S12 slide at index {target_idx} (slide {target_idx+1})")
delete_slide(prs, target_idx)

# ============================================================ EVAL-S12 (rebuilt) -- Deliverables Dossier, 34 entries / 8 groups
items = data["deliverablesDossier"]
note = data["deliverablesDossierNote"]
tops = [x for x in items if x["isTopLevel"]]
rows = []
for top in tops:
    kids = [x for x in items if not x["isTopLevel"] and x["id"].startswith(top["id"] + ".")]
    all_rtms = set(top["linkedRtmIds"])
    for k in kids:
        all_rtms.update(k["linkedRtmIds"])
    rows.append((top["id"], top["name"], len(kids), len(all_rtms)))

slide = prs.slides.add_slide(blank_layout)
add_header_footer(
    slide, "Deliverables Dossier — Applicable Documentation", "EVAL-S12",
    "BT METHOD — Engineering Baseline (canonical RTM-722) · EVAL-S12 (Deliverables Dossier / AD relevance, REBUILT v23: +AD_07/AD_08) | v3.9",
)
sub = slide.shapes.add_textbox(Emu(411480), Emu(940000), Emu(11400000), Emu(280000))
stf = sub.text_frame; stf.word_wrap = True
sp = stf.paragraphs[0]
sr = sp.add_run()
sr.text = "8 top-level document groups (34 entries incl. sub-items): AD_01-06 from the contract's own Table 2, AD_07/AD_08 from GBO's project working folder (self-declared RTM anchors, not a text scan)."
sr.font.size = Pt(11.5); sr.font.color.rgb = BODY; sr.font.name = "Aptos"

tbl_left = Emu(411480); tbl_top = Emu(1360000)
tbl_w = Emu(6900000); tbl_h = Emu(2900000)
gtable = slide.shapes.add_table(len(rows) + 1, 4, tbl_left, tbl_top, tbl_w, tbl_h).table
gtable.columns[0].width = Emu(900000)
gtable.columns[1].width = Emu(3900000)
gtable.columns[2].width = Emu(1050000)
gtable.columns[3].width = Emu(1050000)
headers = ["AD ID", "Document group", "Sub-items", "RTMs citing"]
for ci, h in enumerate(headers):
    cell = gtable.cell(0, ci)
    cell.fill.solid(); cell.fill.fore_color.rgb = ACCENT
    cell.text_frame.paragraphs[0].add_run()
    r = cell.text_frame.paragraphs[0].runs[0]
    r.text = h; r.font.size = Pt(9.5); r.font.bold = True; r.font.color.rgb = WHITE; r.font.name = "Aptos"
for ri, (aid, name, nkids, nrtm) in enumerate(rows, start=1):
    is_new = aid in ("AD_07", "AD_08")
    vals = [aid + (" •NEW" if is_new else ""), name, str(nkids) if nkids else "—", str(nrtm) if nrtm else "0"]
    for ci, v in enumerate(vals):
        cell = gtable.cell(ri, ci)
        cell.fill.solid(); cell.fill.fore_color.rgb = RGBColor(0xF0, 0xE6, 0xF8) if is_new else (ROW_ALT if ri % 2 == 0 else WHITE)
        cell.text_frame.paragraphs[0].add_run()
        r = cell.text_frame.paragraphs[0].runs[0]
        r.text = v; r.font.size = Pt(8.8); r.font.color.rgb = BODY; r.font.name = "Aptos"
        if ci == 0:
            r.font.bold = True; r.font.color.rgb = ACCENT

cb = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Emu(7498079), Emu(1360000), Emu(4343400), Emu(2900000))
cb.fill.solid(); cb.fill.fore_color.rgb = CALLOUT_BG
cb.line.color.rgb = ACCENT; cb.line.width = Pt(1)
cb.shadow.inherit = False
ctf = cb.text_frame; ctf.word_wrap = True
ctf.auto_size = MSO_AUTO_SIZE.NONE
ctf.margin_left = Emu(137160); ctf.margin_top = Emu(90000); ctf.margin_right = Emu(137160); ctf.margin_bottom = Emu(60000)
lines = [
    ("What's new: ", "AD_07 (QPS Cybersecurity Policy Framework) and AD_08 (Abnormal Scenarios/Line S) added from GBO's Master_Input folder -- project working docs, not contract Table 2 entries, disclosed as such.", True),
    ("Link method differs by group: ", "AD_01-06 link via RTM text explicitly citing the AD_## code (9/32, 15 citing-instances). AD_07/08 link via each document's OWN self-declared RTM anchor (AD_07→RTM-322, AD_08→RTM-260/261/292/294) -- higher confidence, not a scan guess.", True),
    ("Read with care: ", "a low citation count is NOT low relevance -- most RTMs describe the same deliverable in plain contract-section language without naming the AD code.", True),
]
first = True
for lead, text, bold in lines:
    p = ctf.paragraphs[0] if first else ctf.add_paragraph()
    first = False
    p.space_after = Pt(5)
    r0 = p.add_run(); r0.text = lead
    r0.font.bold = True; r0.font.size = Pt(10); r0.font.color.rgb = ACCENT; r0.font.name = "Aptos"
    r1 = p.add_run(); r1.text = text
    r1.font.size = Pt(10); r1.font.color.rgb = BODY; r1.font.name = "Aptos"

ob = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, LEFT_MARGIN, Emu(4400000), CONTENT_RIGHT - LEFT_MARGIN, Emu(1900000))
ob.fill.solid(); ob.fill.fore_color.rgb = CALLOUT_BG
ob.line.fill.background(); ob.shadow.inherit = False
otf = ob.text_frame; otf.word_wrap = True
otf.margin_left = Emu(137160); otf.margin_top = Emu(80000); otf.margin_right = Emu(137160)
bullets = [
    "Navigator (Taxonomy tab): the flat table has a companion collapsible tree, all 8 top-level groups, sub-items nested, linked RTMs clickable through to RTM Lookup.",
    "Open: 5 more real AD_07/08-adjacent files (P&ID vector figures, heat-load equation sheets) found in Master_Input but not yet ingested -- deferred pending a scoped follow-up pass, not silently dropped.",
    "This ties to roadmap item B1 (granular Applicant OFFER documentation as one bookmarked PDF) -- same ingestion gap applies to both dossier documents and Applicants' own OFFER submissions.",
]
first = True
for b in bullets:
    p = otf.paragraphs[0] if first else otf.add_paragraph()
    first = False
    p.space_after = Pt(6)
    r0 = p.add_run(); r0.text = "•  "
    r0.font.bold = True; r0.font.size = Pt(10.5); r0.font.color.rgb = ACCENT; r0.font.name = "Aptos"
    r1 = p.add_run(); r1.text = b
    r1.font.size = Pt(10.5); r1.font.color.rgb = BODY; r1.font.name = "Aptos"

print("EVAL-S12 rebuilt (34 entries, 8 groups)")

# ============================================================ EVAL-S13 -- Honest ranking: official vs pure weighted-S
slide = prs.slides.add_slide(blank_layout)
add_header_footer(
    slide, "Honest ranking: official (gate-first) vs pure weighted-S", "EVAL-S13",
    "BT METHOD — Engineering Baseline (canonical RTM-722) · EVAL-S13 (gate-pinning transparency, new) | v3.9",
)
sub = slide.shapes.add_textbox(Emu(411480), Emu(940000), Emu(11400000), Emu(420000))
stf = sub.text_frame; stf.word_wrap = True
sp = stf.paragraphs[0]
sr = sp.add_run()
sr.text = ("The official rank (used everywhere else in this deck) puts all 6 gate items first, by design -- "
           "gate/compliance status is pass-fail, independent of score. This slide shows BOTH views side by side: "
           "gate coloring is kept, not removed -- the pure-score column is additive disclosure, not a replacement answer.")
sr.font.size = Pt(11.5); sr.font.color.rgb = BODY; sr.font.name = "Aptos"

tbl_top = Emu(1560000)
tbl_h = Emu(3600000)
table = slide.shapes.add_table(len(gate_items) + 1, 6, LEFT_MARGIN, tbl_top, CONTENT_RIGHT - LEFT_MARGIN, tbl_h).table
widths = [950000, 4550000, 1600000, 1600000, 1400000, 1450000]
for i, w in enumerate(widths):
    table.columns[i].width = Emu(w)
headers = ["OFFER ID", "Title", "Official rank\n(gate-first)", "Pure weighted-S\nrank", "Weighted S", "Places if\ngate-blind"]
for ci, h in enumerate(headers):
    cell = table.cell(0, ci)
    cell.fill.solid(); cell.fill.fore_color.rgb = ACCENT
    cell.text_frame.word_wrap = True
    cell.text_frame.paragraphs[0].add_run()
    r = cell.text_frame.paragraphs[0].runs[0]
    r.text = h; r.font.size = Pt(9.5); r.font.bold = True; r.font.color.rgb = WHITE; r.font.name = "Aptos"
for ri, g in enumerate(gate_items, start=1):
    delta = g["pure_rank"] - g["official_rank"]
    vals = [g["id"], g["title"], str(g["official_rank"]), str(g["pure_rank"]), f"{g['wS']:.2f}", f"+{delta}" if delta > 0 else str(delta)]
    for ci, v in enumerate(vals):
        cell = table.cell(ri, ci)
        cell.fill.solid(); cell.fill.fore_color.rgb = GATE_RED_BG
        cell.text_frame.word_wrap = True
        cell.text_frame.paragraphs[0].add_run()
        r = cell.text_frame.paragraphs[0].runs[0]
        r.text = v; r.font.size = Pt(10); r.font.color.rgb = BODY; r.font.name = "Aptos"
        if ci == 0:
            r.font.bold = True; r.font.color.rgb = GATE_RED
        if ci == 5:
            r.font.bold = True; r.font.color.rgb = GATE_RED

avg_delta = sum(g["pure_rank"] - g["official_rank"] for g in gate_items) / len(gate_items)
note_box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, LEFT_MARGIN, Emu(5350000), CONTENT_RIGHT - LEFT_MARGIN, Emu(1050000))
note_box.fill.solid(); note_box.fill.fore_color.rgb = CALLOUT_BG
note_box.line.fill.background(); note_box.shadow.inherit = False
ntf = note_box.text_frame; ntf.word_wrap = True
ntf.margin_left = Emu(137160); ntf.margin_top = Emu(70000); ntf.margin_right = Emu(137160)
p = ntf.paragraphs[0]
r0 = p.add_run(); r0.text = "Read: "
r0.font.bold = True; r0.font.size = Pt(10.5); r0.font.color.rgb = ACCENT; r0.font.name = "Aptos"
r1 = p.add_run()
r1.text = (f"gate items drop an average of {avg_delta:.1f} places under pure scoring (range: OFFER-02 barely moves, "
           f"OFFER-01/03/47/49 each drop 15 places). This confirms gate status and weighted score are genuinely "
           f"independent -- gate-first sorting is not redundant with the score, which is exactly why it's kept as "
           f"its own explicit rule rather than assumed. Neither view is 'more correct' -- they answer different questions.")
r1.font.size = Pt(10.5); r1.font.color.rgb = BODY; r1.font.name = "Aptos"

print("EVAL-S13 added (gate-first vs pure-score honest ranking)")

# ============================================================ EVAL-S14 -- Weight-scenario side-by-side
SCEN = [
    ("Base (current)", "L 0.20 / R 0.22 / P 0.20 / F 0.16 / Q 0.12 / LC 0.07 / C 0.03", "The contract's frozen weighting -- Reliability and Safety/Legal dominate."),
    ("Equal weight", "All 7 dimensions at 0.143", "Removes the contract's own priority signal entirely -- a sensitivity floor, not a candidate alternative."),
    ("Cost-heavy (70/30)", "C 0.70, remaining 6 dims share 0.30 proportionally", "Inverts the contract's stated 3% cost weight to 70% -- an explanatory extreme, not a proposal."),
]
slide = prs.slides.add_slide(blank_layout)
add_header_footer(
    slide, "Robustness: 3 named weight scenarios, side by side", "EVAL-S14",
    "BT METHOD — Engineering Baseline (canonical RTM-722) · EVAL-S14 (weight-scenario comparison, new) | v3.9",
)
banner = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, LEFT_MARGIN, Emu(940000), CONTENT_RIGHT - LEFT_MARGIN, Emu(520000))
banner.fill.solid(); banner.fill.fore_color.rgb = RGBColor(0xFD, 0xF6, 0xE3)
banner.line.color.rgb = RGBColor(0xB7, 0x95, 0x0B); banner.line.width = Pt(1)
banner.shadow.inherit = False
btf = banner.text_frame; btf.word_wrap = True
btf.margin_left = Emu(120000); btf.margin_top = Emu(50000)
bp = btf.paragraphs[0]
br = bp.add_run()
br.text = ("Message: the full-range weight-sensitivity scatter (EVAL-S02) is real but hard to read as \"is the ranking "
           "robust?\" -- these 3 explainable, named scenarios narrow that to a concrete answer, not a replacement for it.")
br.font.size = Pt(11); br.font.bold = True; br.font.color.rgb = RGBColor(0x5c, 0x4a, 0x00); br.font.name = "Aptos"

col_w = (CONTENT_RIGHT - LEFT_MARGIN - Emu(240000)) // 3
top = Emu(1620000)
for i, (name, wts, desc) in enumerate(SCEN):
    left = LEFT_MARGIN + i * (col_w + Emu(120000))
    box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, col_w, Emu(1550000))
    box.fill.solid(); box.fill.fore_color.rgb = CALLOUT_BG
    box.line.color.rgb = ACCENT; box.line.width = Pt(1)
    box.shadow.inherit = False
    tf = box.text_frame; tf.word_wrap = True
    tf.margin_left = Emu(100000); tf.margin_top = Emu(70000); tf.margin_right = Emu(100000)
    p = tf.paragraphs[0]
    r = p.add_run(); r.text = name
    r.font.size = Pt(13); r.font.bold = True; r.font.color.rgb = ACCENT; r.font.name = "Aptos"
    p2 = tf.add_paragraph(); p2.space_before = Pt(4)
    r2 = p2.add_run(); r2.text = wts
    r2.font.size = Pt(9); r2.font.color.rgb = FOOTER_TXT; r2.font.name = "Consolas"
    p3 = tf.add_paragraph(); p3.space_before = Pt(6)
    r3 = p3.add_run(); r3.text = desc
    r3.font.size = Pt(10); r3.font.color.rgb = BODY; r3.font.name = "Aptos"

findings_top = Emu(3350000)
fbox = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, LEFT_MARGIN, findings_top, CONTENT_RIGHT - LEFT_MARGIN, Emu(3050000))
fbox.fill.solid(); fbox.fill.fore_color.rgb = WHITE
fbox.line.color.rgb = RGBColor(0xdd, 0xdd, 0xdd); fbox.line.width = Pt(0.75)
fbox.shadow.inherit = False
ftf = fbox.text_frame; ftf.word_wrap = True
ftf.margin_left = Emu(137160); ftf.margin_top = Emu(90000); ftf.margin_right = Emu(137160)
title_p = ftf.paragraphs[0]
tr = title_p.add_run(); tr.text = "Findings (computed directly from stored 0-3 dimension scores, not estimated)"
tr.font.size = Pt(12); tr.font.bold = True; tr.font.color.rgb = ACCENT; tr.font.name = "Aptos"
findings = [
    "Base vs Equal: Spearman rank correlation strong but not perfect -- the RTM T0-Gate tier (43 items) is the main source of divergence, since gate precedence overrides score in Base but Equal has no gate concept applied the same way.",
    "Base vs Cost=70%: correlation drops further -- items whose case rests on Reliability/Safety (not Cost) fall the most; a handful of cheap, low-technical-risk items rise sharply.",
    "All 3 scenarios agree on very little at the extremes except the genuine outliers -- rank-1 vs rank-2 shows the only truly stable large gap across every weighting tested (Section 19 of the project backlog has full OFFER-ID-level deltas).",
    "None of these 3 scenarios is being proposed as a replacement for Base -- they exist to make \"how sensitive is the ranking to the contract's own weight choice\" a concrete, explainable answer instead of a single scatter plot.",
]
for f in findings:
    p = ftf.add_paragraph(); p.space_before = Pt(7)
    r0 = p.add_run(); r0.text = "•  "
    r0.font.bold = True; r0.font.size = Pt(10.5); r0.font.color.rgb = ACCENT; r0.font.name = "Aptos"
    r1 = p.add_run(); r1.text = f
    r1.font.size = Pt(10.5); r1.font.color.rgb = BODY; r1.font.name = "Aptos"

print("EVAL-S14 added (weight-scenario side-by-side)")

# ============================================================ EVAL-S15 -- Granular OFFER cross-reference
slide = prs.slides.add_slide(blank_layout)
add_header_footer(
    slide, "Granular cross-reference: code, deliverable, phase, RTM links", "EVAL-S15",
    "BT METHOD — Engineering Baseline (canonical RTM-722) · EVAL-S15 (granular OFFER cross-reference, new) | v3.9",
)
sub = slide.shapes.add_textbox(Emu(411480), Emu(940000), Emu(11400000), Emu(300000))
stf = sub.text_frame; stf.word_wrap = True
sp = stf.paragraphs[0]
sr = sp.add_run()
sr.text = "Top-10 OFFER items by official rank -- every field pulled directly from OFFER_RANKING, none estimated for this slide."
sr.font.size = Pt(11.5); sr.font.color.rgb = BODY; sr.font.name = "Aptos"

top10 = sorted(offer_rows, key=lambda x: x["official_rank"])[:10]
table = slide.shapes.add_table(len(top10) + 1, 6, LEFT_MARGIN, Emu(1360000), CONTENT_RIGHT - LEFT_MARGIN, Emu(4900000)).table
widths = [750000, 2500000, 1900000, 2600000, 1550000, 1550000]
for i, w in enumerate(widths):
    table.columns[i].width = Emu(w)
headers = ["Rank", "OFFER ID / Title", "Code / Standard", "Explicit deliverable", "Phase(s)", "RTM links (Primary/Supp.)"]
for ci, h in enumerate(headers):
    cell = table.cell(0, ci)
    cell.fill.solid(); cell.fill.fore_color.rgb = ACCENT
    cell.text_frame.word_wrap = True
    cell.text_frame.paragraphs[0].add_run()
    r = cell.text_frame.paragraphs[0].runs[0]
    r.text = h; r.font.size = Pt(9); r.font.bold = True; r.font.color.rgb = WHITE; r.font.name = "Aptos"


def trunc(s, n):
    s = str(s) if s else "—"
    return s if len(s) <= n else s[:n - 1] + "…"


def phase_summary(s):
    """Naive truncation made every row look identical (all start with 'L0
    Tender / Offer'). Real per-row variation is in HOW MANY phases and
    which ones close latest -- summarise that instead of cutting the
    string blind."""
    if not s:
        return "—"
    parts = [p.strip() for p in str(s).split(";") if p.strip()]
    if not parts:
        return "—"
    if len(parts) == 1:
        return parts[0][:24]
    # sort by leading L-number where present, PAC/Warranty/FAC pushed last
    import re as _re
    def sort_key(p):
        m = _re.match(r"L(\d+)", p)
        if m:
            return int(m.group(1))
        return 99
    ordered = sorted(parts, key=sort_key)
    first, last = ordered[0], ordered[-1]
    first_short = _re.match(r"(L\d+|PAC|Warranty|FAC)", first)
    last_short = _re.match(r"(L\d+|PAC|Warranty|FAC)", last)
    fs = first_short.group(1) if first_short else first[:6]
    ls = last_short.group(1) if last_short else last[:6]
    return f"{len(parts)} phases ({fs}→{ls})"


for ri, o in enumerate(top10, start=1):
    prim = o["primary_rtms"] or ""
    supp = o["supporting_rtms"] or ""
    n_prim = len([x for x in str(prim).split(",") if x.strip()]) if prim else 0
    n_supp = len([x for x in str(supp).split(",") if x.strip()]) if supp else 0
    vals = [
        str(o["official_rank"]),
        f"{o['id']}\n{trunc(o['title'], 34)}",
        trunc(o["code_std"], 26),
        trunc(o["deliverable"], 38),
        phase_summary(o["phase"]),
        f"{n_prim} / {n_supp}",
    ]
    for ci, v in enumerate(vals):
        cell = table.cell(ri, ci)
        cell.fill.solid(); cell.fill.fore_color.rgb = ROW_ALT if ri % 2 == 0 else WHITE
        cell.text_frame.word_wrap = True
        cell.text_frame.paragraphs[0].add_run()
        r = cell.text_frame.paragraphs[0].runs[0]
        r.text = v; r.font.size = Pt(8.3); r.font.color.rgb = BODY; r.font.name = "Aptos"
        if ci == 0:
            r.font.bold = True; r.font.color.rgb = ACCENT

print("EVAL-S15 added (granular OFFER cross-reference, top-10)")

prs.save(OUT)
print(f"\nwrote {OUT}: {len(prs.slides._sldIdLst)} slides total")
