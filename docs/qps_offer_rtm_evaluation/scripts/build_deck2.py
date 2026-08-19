#!/usr/bin/env python3
"""Pass 2: real content/visual redesign of the baseline slides, 3-way roadmap
split, and the new Unit Economics slide. Input: step2.pptx (structural
insertions already done). Output: final2.pptx."""
from pptx import Presentation
from pptx.util import Emu, Pt, Inches
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_THEME_COLOR_INDEX as TC
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
import copy as _copy

SRC = "step2.pptx"
OUT = "final2.pptx"
CH = "/home/claude/work/charts"

prs = Presentation(SRC)
slides = list(prs.slides)

# ---------------------------------------------------------------- helpers --

def set_run(run, text=None, size=None, bold=None, italic=None, color_theme=None,
            brightness=None, rgb=None, font=None):
    if text is not None:
        run.text = text
    f = run.font
    if size is not None:
        f.size = Pt(size)
    if bold is not None:
        f.bold = bold
    if italic is not None:
        f.italic = italic
    if font is not None:
        f.name = font
    if color_theme is not None:
        f.color.theme_color = color_theme
        if brightness is not None:
            f.color.brightness = brightness
    elif rgb is not None:
        f.color.rgb = RGBColor(*rgb)

def set_notes(slide, text):
    ns = slide.notes_slide
    tf = ns.notes_text_frame
    tf.text = text

def find_shape(slide, name):
    for sh in slide.shapes:
        if sh.name == name:
            return sh
    return None

def find_banner(slide):
    for sh in slide.shapes:
        if sh.shape_type == 1 and sh.name.startswith("Rectangle"):
            try:
                if sh.fill.fore_color.theme_color in (TC.ACCENT_4, TC.ACCENT_1, TC.ACCENT_2):
                    return sh
            except Exception:
                pass
    return None

def set_banner(slide, new_title=None, theme=None):
    banner = find_banner(slide)
    if banner is None:
        return None
    if theme is not None:
        banner.fill.fore_color.theme_color = theme
    if new_title is not None:
        tf = banner.text_frame
        p = tf.paragraphs[0]
        for extra in p.runs[1:]:
            extra._r.getparent().remove(extra._r)
        p.runs[0].text = new_title
    return banner

# Phase color system - reused for banners (Analyze section), tags, and diagrams
PHASE_COLOR = {
    "DEFINE": TC.ACCENT_1, "MEASURE": TC.ACCENT_2, "ANALYZE": TC.ACCENT_4,
    "IMPROVE": TC.ACCENT_1, "CONTROL": TC.ACCENT_2,
}

def add_tag(slide, label, style="layout"):
    color = PHASE_COLOR.get(label, TC.ACCENT_1)
    if style == "layout":
        x, y, cx, cy = Emu(9950000), Emu(160000), Emu(1750000), Emu(320000)
    else:
        x, y, cx, cy = Emu(10760000), Emu(300000), Emu(1330000), Emu(340000)
    box = slide.shapes.add_textbox(x, y, cx, cy)
    tf = box.text_frame
    tf.word_wrap = False
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    r = p.add_run()
    set_run(r, text=label, size=10.5, bold=True, color_theme=color, font="Calibri")
    r.font._rPr.set('spc', '120')
    return box

def add_takeaway(slide, x, y, cx, cy, label, text, size=12, color=TC.ACCENT_1):
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, cx, cy)
    box.adjustments[0] = 0.06
    box.fill.solid()
    box.fill.fore_color.theme_color = color
    box.fill.fore_color.brightness = 0.90
    box.line.fill.background()
    box.shadow.inherit = False
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(137160); tf.margin_right = Emu(137160)
    tf.margin_top = Emu(80000); tf.margin_bottom = Emu(80000)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.line_spacing = 1.1
    r1 = p.add_run(); set_run(r1, text=label + "  ", size=size, bold=True, color_theme=color)
    r2 = p.add_run(); set_run(r2, text=text, size=size, bold=False, color_theme=color)
    return box

def add_bullets(slide, x, y, cx, cy, items, size=13, color_theme=TC.TEXT_1,
                 brightness=0.1, bold_lead_word=None, space_after=8, line_spacing=1.12):
    box = slide.shapes.add_textbox(x, y, cx, cy)
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.line_spacing = line_spacing
        p.space_after = Pt(space_after)
        pPr = p._p.get_or_add_pPr()
        pPr.set('marL', '182880'); pPr.set('indent', '-182880')
        buFont = pPr.makeelement('{http://schemas.openxmlformats.org/drawingml/2006/main}buFont',
                                  {'typeface': 'Arial'})
        buChar = pPr.makeelement('{http://schemas.openxmlformats.org/drawingml/2006/main}buChar',
                                  {'char': '•'})
        pPr.append(buFont); pPr.append(buChar)
        r = p.add_run()
        set_run(r, text=item, size=size, color_theme=color_theme, brightness=brightness)
    return box

def add_picture_fit(slide, path, x, y, cx, cy):
    from PIL import Image
    im = Image.open(path)
    iw, ih = im.size
    scale = min(cx / iw, cy / ih)
    w, h = int(iw * scale), int(ih * scale)
    ox = x + (cx - w) // 2
    oy = y + (cy - h) // 2
    return slide.shapes.add_picture(path, ox, oy, w, h)

def add_step_chips(slide, x, y, cx, steps, color=TC.ACCENT_1, step_h=850000, gap=60000,
                    num_size=13, head_size=12.5, body_size=10.8):
    for i, (num, head, body) in enumerate(steps):
        sy = y + i * (step_h + gap)
        chip = slide.shapes.add_shape(MSO_SHAPE.OVAL, x, sy, Emu(340000), Emu(340000))
        chip.fill.solid(); chip.fill.fore_color.theme_color = color
        chip.line.fill.background(); chip.shadow.inherit = False
        ctf = chip.text_frame
        ctf.margin_left = ctf.margin_right = ctf.margin_top = ctf.margin_bottom = 0
        ctf.vertical_anchor = MSO_ANCHOR.MIDDLE
        cp = ctf.paragraphs[0]; cp.alignment = PP_ALIGN.CENTER
        cr = cp.add_run(); set_run(cr, text=str(num), size=num_size, bold=True, color_theme=TC.BACKGROUND_1)
        txt = slide.shapes.add_textbox(x + Emu(430000), sy - Emu(30000), cx - Emu(430000), Emu(step_h))
        ttf = txt.text_frame; ttf.word_wrap = True
        tp = ttf.paragraphs[0]; tp.line_spacing = 1.05
        tr = tp.add_run(); set_run(tr, text=head, size=head_size, bold=True, color_theme=color)
        if body:
            tp2 = ttf.add_paragraph(); tp2.line_spacing = 1.05; tp2.space_before = Pt(2)
            tr2 = tp2.add_run(); set_run(tr2, text=body, size=body_size, color_theme=TC.TEXT_1, brightness=0.15)

def add_stat_tile(slide, x, y, cx, cy, value, label, color=TC.ACCENT_1):
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, cx, cy)
    box.adjustments[0] = 0.08
    box.fill.solid(); box.fill.fore_color.theme_color = color; box.fill.fore_color.brightness = 0.92
    box.line.color.theme_color = color; box.line.width = Pt(0.75)
    box.shadow.inherit = False
    tf = box.text_frame; tf.word_wrap = True
    tf.margin_left = Emu(100000); tf.margin_right = Emu(100000)
    tf.margin_top = Emu(70000); tf.margin_bottom = Emu(70000)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER; p.line_spacing = 1.0
    r = p.add_run(); set_run(r, text=value, size=20, bold=True, color_theme=color)
    p2 = tf.add_paragraph(); p2.alignment = PP_ALIGN.CENTER; p2.line_spacing = 1.05
    r2 = p2.add_run(); set_run(r2, text=label, size=10.5, color_theme=TC.TEXT_1, brightness=0.2)
    return box

def add_arrow(slide, x1, y1, x2, y2, color=TC.ACCENT_1, weight=1.75):
    conn = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x1, y1, x2, y2)
    conn.line.color.theme_color = color
    conn.line.width = Pt(weight)
    ln = conn.line._get_or_add_ln()
    from pptx.oxml.ns import qn
    tail = ln.makeelement(qn('a:tailEnd'), {'type': 'triangle', 'w': 'med', 'len': 'med'})
    ln.append(tail)
    return conn

OLD_TAGS = {"DEFINE", "MEASURE", "ANALYZE", "IMPROVE", "CONTROL"}

def remove_old_tag(slide):
    for sh in list(slide.shapes):
        if sh.has_text_frame and sh.text_frame.text.strip() in OLD_TAGS:
            sh._element.getparent().remove(sh._element)

print("helpers2 loaded OK")

# ============================================================ ROADMAP 1/3-3/3 =
# pos2=slide2.xml, pos3=slide43.xml (Analyze), pos4=slide44.xml (Improve/Control)
r1, r2, r3 = slides[1], slides[2], slides[3]

ROADMAP_TOP = Emu(2260000)  # shared top-of-content line for all 3 roadmap slides

def rebuild_roadmap_shell(s, title, subtitle, banner_theme):
    t = find_shape(s, "Title 1")
    t.left, t.top, t.width, t.height = Emu(766763), Emu(700000), Emu(10650000), Emu(950000)
    t.text_frame.word_wrap = True
    tp = t.text_frame.paragraphs[0]
    tp.runs[0].text = title
    tp.runs[0].font.size = Pt(27)
    sub = find_shape(s, "Text Placeholder 2")
    sub.left, sub.top, sub.width, sub.height = Emu(766763), Emu(1720000), Emu(10650000), Emu(460000)
    sp = sub.text_frame.paragraphs[0]
    sp.runs[0].text = subtitle
    sp.runs[0].font.size = Pt(15)
    body = find_shape(s, "Text Placeholder 3")
    # wipe the body placeholder's existing paragraphs entirely
    txBody = body.text_frame._txBody
    for p in list(txBody.findall('{http://schemas.openxmlformats.org/drawingml/2006/main}p')):
        txBody.remove(p)
    body.left, body.top, body.width, body.height = Emu(766763), ROADMAP_TOP, Emu(5300000), Emu(4150000)
    # also recolor the title/subtitle to the phase color for visual identity
    for r in t.text_frame.paragraphs[0].runs:
        r.font.color.theme_color = banner_theme
    for r in sp.runs:
        r.font.color.theme_color = banner_theme
    return body

# ---- Roadmap 1/3: DEFINE + MEASURE --------------------------------------------
body = rebuild_roadmap_shell(
    r1, "DMAIC Roadmap (1/3) — Define & Measure",
    "How to read this master deck, part 1 — scope, requirements, and baseline data",
    TC.ACCENT_1,
)
tf = body.text_frame
items = [
    ("DEFINE", "Slides 5–7, 13, 20–21", "What the system must achieve: scope, RTM obligations, the operating-state hierarchy, and failure-class definitions with their MTBF targets."),
    ("MEASURE", "Slides 8–12, 16–18", "What the fleet can actually deliver: OEM flow data, N−1 capability, unit energy economics, acceptance-test criteria, and liquid-buffer coverage time."),
]
for i, (phase, rng, desc) in enumerate(items):
    p = tf.add_paragraph()
    p.line_spacing = 1.1; p.space_after = Pt(10)
    r1_ = p.add_run(); set_run(r1_, text=phase, size=16, bold=True, color_theme=TC.ACCENT_1)
    r2_ = p.add_run(); set_run(r2_, text=f"  ({rng})", size=12.5, italic=True, color_theme=TC.ACCENT_2)
    p2 = tf.add_paragraph(); p2.line_spacing = 1.12; p2.space_after = Pt(4)
    r3_ = p2.add_run(); set_run(r3_, text=desc, size=12.5, color_theme=TC.TEXT_1, brightness=0.1)

# right panel: 4 stat tiles = the baseline numbers this phase establishes
tiles = [
    ("307 g/s", "24 QM target flow\n(design envelope 320–350 g/s)"),
    ("N − 1", "sizing rule: fleet must\nsurvive one unit down"),
    ("A>5y · B>10y · C>15y", "MTBF targets by\nfailure-consequence class"),
    ("350 kW", "installed power per HP\ncompressor (validated Slide 10)"),
]
tx, ty, tw, th, gap = Emu(6350000), Emu(2260000), Emu(2560000), Emu(1200000), Emu(120000)
for i, (val, lab) in enumerate(tiles):
    col, row = i % 2, i // 2
    x = tx + col * (tw + gap); y = ty + row * (th + gap)
    add_stat_tile(r1, x, y, tw, th, val, lab, color=TC.ACCENT_1)
add_takeaway(r1, Emu(6350000), Emu(4950000), Emu(5240000), Emu(800000), "SAMPLE CALC",
             "N−1 check, FSD 575: 3 × 112.5 g/s @ 72 Hz = 337.5 g/s ≥ 307 g/s target → passes with +30 g/s margin (Slide 9).",
             size=11, color=TC.ACCENT_1)
p = tf.add_paragraph(); p.space_before = Pt(8)
r = p.add_run()
set_run(r, text="This is a MASTER slide set — each third of this roadmap can be lifted into its own briefing.",
        size=10.5, italic=True, color_theme=TC.TEXT_1, brightness=0.35)
add_tag(r1, "DEFINE", style="layout")

print("Roadmap 1/3 done")

# ---- Roadmap 2/3: ANALYZE ------------------------------------------------------
body = rebuild_roadmap_shell(
    r2, "DMAIC Roadmap (2/3) — Analyze",
    "How to read this master deck, part 2 — why the numbers behave the way they do",
    TC.ACCENT_4,
)
tf = body.text_frame
p = tf.add_paragraph(); p.line_spacing = 1.1; p.space_after = Pt(10)
r_ = p.add_run(); set_run(r_, text="ANALYZE", size=16, bold=True, color_theme=TC.ACCENT_4)
r2_ = p.add_run(); set_run(r2_, text="  (Slides 22–36)", size=12.5, italic=True, color_theme=TC.ACCENT_2)
p2 = tf.add_paragraph(); p2.line_spacing = 1.15; p2.space_after = Pt(6)
r3_ = p2.add_run()
set_run(r3_, text="This is the deck's math core: the Poisson/exponential campaign-success model, "
              "the 94–95% reliability ceiling, the Wrong / Better / Chosen architecture trade-off, "
              "and the Weibull wear-out relation that eventually limits any MTBF-only argument.",
        size=12.5, color_theme=TC.TEXT_1, brightness=0.1)
p3 = tf.add_paragraph(); p3.space_before = Pt(6)
r4_ = p3.add_run()
set_run(r4_, text="This section is intentionally chart-heavy — the technical-blue banner marks it "
              "as the deck's quantitative narrative, distinct from the definitional purple sections either side.",
        size=11, italic=True, color_theme=TC.TEXT_1, brightness=0.3)

# Right panel: sample calc callouts (native, not an image — matches Slide 19/21's own numbers)
add_takeaway(r2, Emu(6350000), Emu(2260000), Emu(5240000), Emu(1300000), "SAMPLE CALC — Poisson",
             "MTBF = 5 y → λ = 1/5 = 0.200 / y. P(zero trips, 90 d) = e^(−λt) = e^(−0.200 × 0.246) "
             "≈ 95.3%  →  the origin of the deck's ~95% campaign figure (Slides 22–24, 27–29).",
             size=11, color=TC.ACCENT_4)
add_takeaway(r2, Emu(6350000), Emu(3690000), Emu(5240000), Emu(1300000), "SAMPLE CALC — Weibull",
             "β = 2.5, η = 114,400 h (cold compressor). P(fail in next 90 d | age = 5 y) is "
             "≈ 235× the same-window probability right after a good-as-new overhaul (Slide 36) — "
             "the mechanism that IMPROVE (part 3) responds to.",
             size=11, color=TC.ACCENT_4)
add_stat_tile(r2, Emu(6350000), Emu(5120000), Emu(2560000), Emu(1050000), "94–95%", "90-day campaign\nsuccess ceiling (single-train)", color=TC.ACCENT_4)
add_stat_tile(r2, Emu(9030000), Emu(5120000), Emu(2560000), Emu(1050000), "β > 1", "hazard rises with age →\nMTBF alone understates risk", color=TC.ACCENT_4)
p4 = tf.add_paragraph(); p4.space_before = Pt(10)
r5_ = p4.add_run()
set_run(r5_, text="This is a MASTER slide set — the Analyze section alone can be lifted into a standalone reliability-math briefing.",
        size=10.5, italic=True, color_theme=TC.TEXT_1, brightness=0.35)
add_tag(r2, "ANALYZE", style="layout")
print("Roadmap 2/3 done")

# ---- Roadmap 3/3: IMPROVE + CONTROL --------------------------------------------
body = rebuild_roadmap_shell(
    r3, "DMAIC Roadmap (3/3) — Improve & Control",
    "How to read this master deck, part 3 — from prediction to a maintenance policy",
    TC.ACCENT_1,
)
tf = body.text_frame
items3 = [
    ("IMPROVE", "Slides 37–42", "Why MTBF matters, when it stops being trustworthy, and the component evidence (compressors, PVPS, turbines) that proves it."),
    ("CONTROL", "Slides 43–46", "Predictable, usage-based service replacement; the reset-to-new renewal principle; and the governance loop that keeps the system inside the assumptions MTBF requires."),
]
for i, (phase, rng, desc) in enumerate(items3):
    p = tf.add_paragraph(); p.line_spacing = 1.1; p.space_after = Pt(10)
    r_ = p.add_run(); set_run(r_, text=phase, size=16, bold=True, color_theme=TC.ACCENT_1)
    r2_ = p.add_run(); set_run(r2_, text=f"  ({rng})", size=12.5, italic=True, color_theme=TC.ACCENT_2)
    p2 = tf.add_paragraph(); p2.line_spacing = 1.12; p2.space_after = Pt(4)
    r3_ = p2.add_run(); set_run(r3_, text=desc, size=12.5, color_theme=TC.TEXT_1, brightness=0.1)

add_takeaway(r3, Emu(6350000), Emu(2260000), Emu(5240000), Emu(1150000), "SAMPLE CALC",
             "Same β=2.5, η=114,400 h cold-compressor pair as part 2 → replace at ≈ 0.65×η "
             "(≈ 74,000 h) instead of running to MTBF-implied end of life — the deadline is set "
             "before the hazard curve turns upward, not after (Slide 44).",
             size=11, color=TC.ACCENT_1)

loop_x, loop_y = Emu(6350000), Emu(3560000)
mini_steps = [("1", "MEASURE exposure", None), ("2", "TRACK age vs. η", None),
              ("3", "REPLACE before wear-out", None), ("4", "RESET & feed back", None)]
add_step_chips(r3, loop_x, loop_y, Emu(5240000), mini_steps, color=TC.ACCENT_1,
               step_h=560000, gap=40000, num_size=11, head_size=12, body_size=9)
p = tf.add_paragraph(); p.space_before = Pt(8)
r = p.add_run()
set_run(r, text="This is a MASTER slide set — Improve + Control together form a standalone maintenance-policy pack.",
        size=10.5, italic=True, color_theme=TC.TEXT_1, brightness=0.35)
add_tag(r3, "CONTROL", style="layout")
print("Roadmap 3/3 done")

# ============================================================ UNIT ECONOMICS =
ue = slides[9]  # pos10
remove_old_tag(ue)  # clone inherited orig6's flat-purple MEASURE tag; remove before adding the phase-colored one
t = find_shape(ue, "Title 1")
t.text_frame.paragraphs[0].runs[0].text = "Unit Economics — Single HP Compressor"
t.text_frame.paragraphs[0].runs[0].font.size = Pt(29)
sub = find_shape(ue, "Text Placeholder 2")
sp = sub.text_frame.paragraphs[0]
for extra in sp.runs[1:]:
    extra._r.getparent().remove(extra._r)
sp.runs[0].text = "Isothermal Compression Power & Energy Cost (new — not in baseline)"
body = find_shape(ue, "Text Placeholder 3")
body.left, body.top, body.width, body.height = Emu(645465), Emu(2260000), Emu(5350000), Emu(4060000)
txBody = body.text_frame._txBody
for p in list(txBody.findall('{http://schemas.openxmlformats.org/drawingml/2006/main}p')):
    txBody.remove(p)
tf = body.text_frame
lines = [
    ("Model", "Isothermal compression, Ẇ = ṁ·R·T·ln(P₂/P₁) — the thermodynamic best case for a water-cooled, intercooled screw compressor."),
    ("Assumptions (flagged)", "T=300 K inlet · P₂=14 bar(a) (Slide 6) · P₁≈1.1 bar(a) LP suction (assumed, not in RTMs) · η_isothermal=50% · helium, R=2077 J/(kg·K)"),
    ("Flow point used", "112.5 g/s = FSD 575 @ 72 Hz (Slide 8) — same per-unit number as the N−1 and Hz-requirement charts."),
    ("VFD & cooling note", "Water-cooling keeps compression near-isothermal; VFD scales power ~linearly with flow, not cubically."),
]
for i, (head, body_txt) in enumerate(lines):
    p = tf.add_paragraph(); p.line_spacing = 1.0; p.space_after = Pt(1)
    r = p.add_run(); set_run(r, text=head, size=11.5, bold=True, color_theme=TC.ACCENT_2)
    p2 = tf.add_paragraph(); p2.line_spacing = 1.02; p2.space_after = Pt(6)
    r2 = p2.add_run(); set_run(r2, text=body_txt, size=10.5, color_theme=TC.TEXT_1, brightness=0.12)

add_takeaway(ue, Emu(645465), Emu(5950000), Emu(11000000), Emu(470000), "KEY TAKEAWAY",
             "The 50%-isothermal model, with a plausible 1.1 bar suction assumption, reproduces the deck's own 350 kW nameplate rating to within 2% — a useful sanity check on both numbers.",
             size=11, color=TC.ACCENT_2)

add_picture_fit(ue, f"{CH}/s_unit_economics.png", Emu(6250000), Emu(1950000), Emu(5400000), Emu(3700000))
add_tag(ue, "MEASURE", style="layout")
set_notes(ue, "This slide is new — it is not derived from a baseline slide. It fills a gap the deck's "
              "compressor sizing sections did not cover: what a single HP compressor actually costs to "
              "run. Isothermal compression is the correct idealisation for a water-cooled, intercooled "
              "screw machine (it is the thermodynamic lower bound on required work; real machines sit "
              "between isothermal and adiabatic depending on cooling effectiveness). The 50% isothermal "
              "efficiency figure is a reasonable planning assumption for this class of machine but should "
              "be confirmed against the OEM's actual performance curve. The suction pressure (1.1 bar) is "
              "not stated anywhere in the RTMs excerpted in this deck and is flagged as an assumption — "
              "sensitivity: 1.05-1.5 bar suction brackets 313-363 kW, i.e. the 350 kW nameplate sits "
              "comfortably inside that range. The energy-cost figures use illustrative tariff bands "
              "(€0.12-0.20/kWh); replace with the actual SCK CEN industrial electricity contract rate "
              "before using this for a budget submission. This slide's real purpose is to make the case, "
              "in Euros, for why the CONTROL section's predictable-replacement policy matters: running "
              "cost this large dwarfs a single overhaul, so keeping the compressor in its efficient "
              "operating window is itself worth money, independent of the reliability argument.")
print("Unit Economics slide done")

# ============================================================ NUMERIC SLIDES =

def gather_and_wipe_body(slide):
    """Collect the full baseline text of Text Placeholder 3 (and Title/Subtitle) for
    notes, then wipe the body placeholder's paragraphs so it can be rebuilt."""
    body = find_shape(slide, "Text Placeholder 3")
    orig_lines = []
    for p in body.text_frame.paragraphs:
        t = p.text.strip()
        if t:
            orig_lines.append(t)
    txBody = body.text_frame._txBody
    for p in list(txBody.findall('{http://schemas.openxmlformats.org/drawingml/2006/main}p')):
        txBody.remove(p)
    return body, orig_lines

def rebuild_numeric_slide(slide, new_title, chart_path, condensed_items, phase, extra_note=None,
                           chart_box=(6400000, 1950000, 5250000, 4300000)):
    remove_old_tag(slide)
    t = find_shape(slide, "Title 1")
    if new_title:
        p0 = t.text_frame.paragraphs[0]
        for extra in p0.runs[1:]:
            extra._r.getparent().remove(extra._r)
        p0.runs[0].text = new_title
    body, orig_lines = gather_and_wipe_body(slide)
    body.left, body.top, body.width, body.height = Emu(645465), Emu(2000000), Emu(5350000), Emu(4350000)
    tf = body.text_frame
    for item in condensed_items:
        p = tf.add_paragraph(); p.line_spacing = 1.14; p.space_after = Pt(9)
        pPr = p._p.get_or_add_pPr(); pPr.set('marL', '182880'); pPr.set('indent', '-182880')
        ns = '{http://schemas.openxmlformats.org/drawingml/2006/main}'
        pPr.append(pPr.makeelement(f'{ns}buFont', {'typeface': 'Arial'}))
        pPr.append(pPr.makeelement(f'{ns}buChar', {'char': '•'}))
        r = p.add_run(); set_run(r, text=item, size=13.5, color_theme=TC.TEXT_1, brightness=0.08)
    if extra_note:
        p = tf.add_paragraph(); p.space_before = Pt(6); p.line_spacing = 1.1
        r = p.add_run(); set_run(r, text=extra_note, size=11, italic=True, color_theme=TC.ACCENT_2)
    cx, cy, cw, ch = chart_box
    add_picture_fit(slide, chart_path, Emu(cx), Emu(cy), Emu(cw), Emu(ch))
    add_tag(slide, phase, style="layout")
    notes = "Full baseline bullet text (verbatim), condensed on-slide for readability:\n\n" + \
            "\n".join(f"- {l}" for l in orig_lines)
    set_notes(slide, notes)
    return body

print("numeric-slide helper loaded")

# ---- pos7: Frequency and Upset Policy (orig4) ----------------------------------
rebuild_numeric_slide(
    slides[6], None, f"{CH}/s04_frequency_zones.png",
    [
        "Nominal: ≤ 65 Hz, sustainable continuous operation.",
        "Transient / upset: up to 72 Hz, ≤ 8 h per event (covers service or trip recovery, 8 h MTTR).",
        "Cumulative > 72 Hz: ≤ 24 h / compressor / month — logged, triggers monthly RAMI review if exceeded.",
        "Recovery: return to ≤ 60 Hz to re-stabilize after a repair.",
    ],
    "DEFINE",
    extra_note="Applicant shall quantify sustained MTBF vs. frequency band (≤ 60 Hz, 60–68 Hz, ≥ 68 Hz).",
    chart_box=(6100000, 2200000, 5650000, 2500000),
)
print("pos7 done")

# ---- pos8: Compressor Flow Capability (orig5) -----------------------------------
rebuild_numeric_slide(
    slides[7], "Compressor Flow Capability", f"{CH}/s05_flow_capability.png",
    [
        "OEM reference figures at 60 Hz and 72 Hz for the three candidate compressor types.",
        "FSD 575 gives the largest per-unit margin — carries directly into the N−1 case (Slide 9).",
        "HSD Combi figures are Option A basis, flagged in the baseline as an estimate ('guess').",
        "Controls principle: hold discharge pressure, use VFD turndown; avoid > 60 Hz except upsets.",
    ],
    "MEASURE",
)
print("pos8 done")

# ---- pos9: Total Flow and N-1 Snapshots (orig6) ---------------------------------
rebuild_numeric_slide(
    slides[8], None, f"{CH}/s06_total_flow_n1.png",
    [
        "Target: 307 g/s total flow (24 QM, FFT operating mode).",
        "FSD 575 is the only type whose N−1 case clears the target outright (+30 g/s margin).",
        "FSD 475's N−1 shortfall (≈19 g/s) is covered by a 5 m³ liquid buffer (≈12 h repair window, Slide 18).",
        "HSD Combi N−1 is marginal — usable only as a short-term check mode.",
    ],
    "MEASURE",
)
print("pos9 done")

# ---- pos11: 3-Only Envelope (orig7) ---------------------------------------------
rebuild_numeric_slide(
    slides[10], "3-Only Envelope (Service / MTTR)", f"{CH}/s07_three_only_envelope.png",
    [
        "Same N−1 question as Slide 9, reframed as the live 3-compressor operating envelope.",
        "FSD 575 clears 307 g/s at 72 Hz; FSD 475 does not.",
        "At ≤ 65 Hz, 3-unit flow is marginal for both types but may meet target with buffer support.",
        "Descent-rate check: 3-of-4 flow ≈ 0.4 m³/h → a 5 m³ buffer gives ≈ 12 h repair window (Slide 18).",
    ],
    "MEASURE",
    extra_note="Applicant shall simulate 2/3 and 3/4 flow cases (Fred/DBE loop analysis) and report restart capability.",
)
sub11 = find_shape(slides[10], "Text Placeholder 2")
sub11.text_frame.paragraphs[0].runs[0].text = "Flow vs Frequency, 3 Compressors Only"
print("pos11 done")

# ---- pos13: Reliability Targets and RAMI (orig9) --------------------------------
rebuild_numeric_slide(
    slides[12], None, f"{CH}/s09_mtbf_by_class.png",
    [
        "MTBF targets scale with the consequence of the failure — Class A / B / C (defined Slide 20).",
        "Planning assumption: rotary screw compressors run ≈ 40,000–80,000 h between major overhauls.",
        "These targets are the input to the Poisson campaign-success model built in Slides 22–36.",
        "Same three numbers are restated per failure-class in Slide 21 — same chart, different lens.",
    ],
    "DEFINE",
    extra_note="Applicant shall submit a RAMI dossier with component MTBFs, frequency bins, maintenance policy, and MTBF impact for ≥ 68 Hz operation.",
)
print("pos13 done")

# ---- pos16: Appendix A - Canonical Operating Scenarios (orig12) ----------------
rebuild_numeric_slide(
    slides[15], "Appendix A — Canonical Scenarios", f"{CH}/s12_appendixA_scenarios.png",
    [
        "Five canonical WCS HP-flow scenarios, cross-checked against Addendum II Table 3 / Fig. 6.",
        "24 QM Operation (307 g/s) is the real operations target used throughout this deck.",
        "30 QM Operation (344 g/s) is the design point — the margin above real operations.",
        "Standby scenarios (both QM counts) size the fleet's low-flow / warm-standby envelope.",
    ],
    "MEASURE",
    chart_box=(6250000, 1950000, 5350000, 4300000),
)
print("pos16 done")

# ---- pos17: Appendix B - N-1 Capability and Hz Requirement (orig13) -----------
rebuild_numeric_slide(
    slides[16], "Appendix B — N−1 Hz Requirement", f"{CH}/s13_appendixB_hz.png",
    [
        "Same N−1 question as Slides 9 and 11, now expressed as required frequency, not flow.",
        "FSD 575 needs only ≈ 65–66 Hz — comfortably inside the ≤ 72 Hz upset ceiling (Slide 7).",
        "FSD 475 needs ≈ 75 Hz — exceeds the limit outright; requires the liquid buffer or a 4th unit.",
        "HSD Combi sits right at the 72 Hz edge — usable only as a short-term check mode.",
    ],
    "MEASURE",
)
print("pos17 done")

# ---- pos18: Appendix C - Liquid Buffer Descent-Rate Mitigation (orig14) -------
rebuild_numeric_slide(
    slides[17], "Appendix C — Liquid Buffer Coverage", f"{CH}/s14_appendixC_buffer.png",
    [
        "5 m³ liquid buffer converts an N−1 flow shortfall into a bounded repair window.",
        "¾ nominal flow (the FSD 475 N−1 case, Slides 9/11): ≈ 12 h coverage before buffer depletion.",
        "⅔ and ½ flow cases are progressively less forgiving — down to ≈ 5 h at half flow.",
        "This coverage time is what sizes the acceptable MTTR in the acceptance-test criteria (Slide 12).",
    ],
    "MEASURE",
    extra_note="Applicant shall validate these loops (2/3, 3/4 nominal) and demonstrate restart criteria after flow recovery.",
    chart_box=(6100000, 2100000, 5600000, 3800000),
)
print("pos18 done")

# ============================================================ DIAGRAM SLIDES =

# ---- pos12: Acceptance and Test System (orig8) - 4-step test flow -------------
s = slides[11]
remove_old_tag(s)
t = find_shape(s, "Title 1")
p0 = t.text_frame.paragraphs[0]
for extra in p0.runs[1:]:
    extra._r.getparent().remove(extra._r)
p0.runs[0].text = "Acceptance & Test System"
body, orig_lines = gather_and_wipe_body(s)
body.left, body.top, body.width, body.height = Emu(645465), Emu(2000000), Emu(4600000), Emu(4300000)
tf = body.text_frame
p = tf.add_paragraph(); p.line_spacing = 1.15
r = p.add_run()
set_run(r, text="Four gated checks, run in sequence, before the plant is accepted for 24×2 K operation "
              "(SB and OP variants). Each step must pass before the next is credited.",
        size=12.5, color_theme=TC.TEXT_1, brightness=0.1)
add_step_chips(s, Emu(6250000), Emu(1950000), Emu(5350000), [
    ("1", "72 Hz stress (per unit)", "≥ 2 h continuous, no derate. Vibration / oil-temp / ΔT logged."),
    ("2", "3-only envelope check", "MTTR 8 h coverage; maintain 24 QM w/ FFT; flow + pressure trend recorded."),
    ("3", "CIS / MCS exchange", "Autonomy & data sync; RTM-0237–0252 pass; event logs + MTBF fields verified."),
    ("4", "Recovery", "Back to ≤ 60 Hz within ≤ 15 min post-upset; control loop verified."),
], color=TC.ACCENT_1, step_h=980000, gap=60000)
add_tag(s, "MEASURE", style="layout")
set_notes(s, "Applicant shall implement parametric scripts to estimate descent-rate and restart "
              "criteria with liquid inventory.\n\nFull baseline text:\n" + "\n".join(f"- {l}" for l in orig_lines))
print("pos12 done")

# ---- pos14: CIS Autonomy and MCS Exchange (orig10) - integration diagram ------
s = slides[13]
remove_old_tag(s)
t = find_shape(s, "Title 1")
p0 = t.text_frame.paragraphs[0]
for extra in p0.runs[1:]:
    extra._r.getparent().remove(extra._r)
p0.runs[0].text = "CIS Autonomy & MCS Exchange"
body, orig_lines = gather_and_wipe_body(s)
body.left, body.top, body.width, body.height = Emu(645465), Emu(2000000), Emu(4600000), Emu(4300000)
tf = body.text_frame
for head, txt in [
    ("RTM 0237–0252", "CIS Autonomy / MCS Exchange obligation (Addendum II verbatim extract)."),
    ("Applicant shall", "Implement real-time linking of event IDs, Hz-band hours, and MTBF tracking per compressor."),
    ("Cross-references", "Fig. 9, Addendum II §3.5.4, RTM-047 (VFD target ≤ 65 Hz nominal)."),
]:
    p = tf.add_paragraph(); p.line_spacing = 1.1; p.space_after = Pt(3)
    r = p.add_run(); set_run(r, text=head, size=12.5, bold=True, color_theme=TC.ACCENT_2)
    p2 = tf.add_paragraph(); p2.line_spacing = 1.1; p2.space_after = Pt(10)
    r2 = p2.add_run(); set_run(r2, text=txt, size=11.5, color_theme=TC.TEXT_1, brightness=0.12)

# integration diagram: QPLANT compressors -> CIS -> MCS, bidirectional
bx, by, bw, bh = Emu(6450000), Emu(2450000), Emu(1500000), Emu(1000000)
gap = Emu(500000)
boxes = [("QPLANT\ncompressors", TC.ACCENT_2), ("CIS", TC.ACCENT_1), ("MCS", TC.ACCENT_2)]
centers = []
for i, (label, color) in enumerate(boxes):
    x = bx + i * (bw + gap)
    box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, by, bw, bh)
    box.adjustments[0] = 0.12
    box.fill.solid(); box.fill.fore_color.theme_color = color
    box.line.fill.background(); box.shadow.inherit = False
    tfb = box.text_frame; tfb.word_wrap = True; tfb.vertical_anchor = MSO_ANCHOR.MIDDLE
    pb = tfb.paragraphs[0]; pb.alignment = PP_ALIGN.CENTER
    rb = pb.add_run(); set_run(rb, text=label, size=13, bold=True, color_theme=TC.BACKGROUND_1)
    centers.append((x, x + bw))
for i in range(2):
    y_mid = by + bh // 2
    add_arrow(s, Emu(centers[i][1]), Emu(y_mid - 60000), Emu(centers[i+1][0]), Emu(y_mid - 60000), color=TC.ACCENT_1, weight=1.5)
    add_arrow(s, Emu(centers[i+1][0]), Emu(y_mid + 60000), Emu(centers[i][1]), Emu(y_mid + 60000), color=TC.ACCENT_2, weight=1.5)
lbl = s.shapes.add_textbox(bx, by + bh + Emu(180000), Emu(5500000), Emu(600000))
ltf = lbl.text_frame; ltf.word_wrap = True
lp = ltf.paragraphs[0]; lp.alignment = PP_ALIGN.CENTER
lr = lp.add_run()
set_run(lr, text="event IDs · Hz-band hours · MTBF tracking, exchanged continuously in both directions",
        size=10.5, italic=True, color_theme=TC.TEXT_1, brightness=0.3)
add_takeaway(s, Emu(6450000), Emu(4600000), Emu(5500000), Emu(700000), "WHY IT MATTERS",
             "This live link is what makes Slide 13's MTBF dossier and the CONTROL section's governance loop possible — the data has to flow before it can be tracked.",
             size=10.5, color=TC.ACCENT_2)
add_tag(s, "DEFINE", style="layout")
set_notes(s, "Full baseline text:\n" + "\n".join(f"- {l}" for l in orig_lines))
print("pos14 done")

# ---- pos15: Utilities and Interfaces (orig11) - hub & spoke -------------------
s = slides[14]
remove_old_tag(s)
t = find_shape(s, "Title 1")
p0 = t.text_frame.paragraphs[0]
for extra in p0.runs[1:]:
    extra._r.getparent().remove(extra._r)
p0.runs[0].text = "Utilities & Interfaces"
body, orig_lines = gather_and_wipe_body(s)
body.left, body.top, body.width, body.height = Emu(645465), Emu(2000000), Emu(4600000), Emu(4300000)
tf = body.text_frame
p = tf.add_paragraph(); p.line_spacing = 1.15
r = p.add_run()
set_run(r, text="Four utility interfaces size and constrain the compressor fleet's operating envelope "
              "used throughout this deck (Slides 8–18). Applicant shall design the mechanical & process tie-in kit.",
        size=12.5, color_theme=TC.TEXT_1, brightness=0.1)

cx, cy = Emu(9000000), Emu(3950000)
hub = s.shapes.add_shape(MSO_SHAPE.OVAL, cx - Emu(750000), cy - Emu(750000), Emu(1500000), Emu(1500000))
hub.fill.solid(); hub.fill.fore_color.theme_color = TC.ACCENT_1
hub.line.fill.background(); hub.shadow.inherit = False
htf = hub.text_frame; htf.word_wrap = True; htf.vertical_anchor = MSO_ANCHOR.MIDDLE
hp = htf.paragraphs[0]; hp.alignment = PP_ALIGN.CENTER
hr = hp.add_run(); set_run(hr, text="QPLANT", size=14, bold=True, color_theme=TC.BACKGROUND_1)

spokes = [
    ("Electrical", "4 × 350 kW\nmetered + aggregated log", -1900000, -2350000),
    ("Cooling Water", "27→37°C (ΔT 10K)\n42°C max (ΔT 15K, PS01)", 1900000, -2350000),
    ("Process Headers", "14 bar supply setpoint\nVFD + trim valves", -1900000, 1650000),
    ("Future 4.5K/2K Tie-in", "QRB provision\npilot valve pads + IO stubs", 1900000, 1650000),
]
for label, detail, dx, dy in spokes:
    bx, by = cx + Emu(dx) - Emu(1150000), cy + Emu(dy)
    box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, bx, by, Emu(2300000), Emu(950000))
    box.adjustments[0] = 0.1
    box.fill.solid(); box.fill.fore_color.theme_color = TC.ACCENT_2; box.fill.fore_color.brightness = 0.88
    box.line.color.theme_color = TC.ACCENT_2; box.line.width = Pt(0.75); box.shadow.inherit = False
    btf = box.text_frame; btf.word_wrap = True; btf.vertical_anchor = MSO_ANCHOR.MIDDLE
    btf.margin_left = Emu(80000); btf.margin_right = Emu(80000)
    bp = btf.paragraphs[0]; bp.alignment = PP_ALIGN.CENTER; bp.line_spacing = 1.0
    br = bp.add_run(); set_run(br, text=label, size=11.5, bold=True, color_theme=TC.ACCENT_2)
    bp2 = btf.add_paragraph(); bp2.alignment = PP_ALIGN.CENTER; bp2.line_spacing = 1.0
    br2 = bp2.add_run(); set_run(br2, text=detail, size=9.5, color_theme=TC.TEXT_1, brightness=0.2)
    bcx, bcy = bx + Emu(1150000), by + Emu(475000)  # box center
    ddx, ddy = bcx - cx, bcy - cy
    sx_, sy_ = cx + int(ddx * 0.16), cy + int(ddy * 0.16)   # start just outside the hub circle
    ex, ey = cx + int(ddx * 0.80), cy + int(ddy * 0.80)     # end just short of the box, not on its text
    add_arrow(s, Emu(sx_), Emu(sy_), Emu(ex), Emu(ey), color=TC.ACCENT_1, weight=1.25)
add_tag(s, "DEFINE", style="layout")
set_notes(s, "Full baseline text:\n" + "\n".join(f"- {l}" for l in orig_lines))
print("pos15 done")

# ---- pos19: Operational Philosophy (orig15) - state-flow diagram --------------
s = slides[18]
remove_old_tag(s)
banner_title = find_shape(s, "TextBox 7")  # actually find by content below
# find the title-like shape robustly: any shape whose text starts with "2. Operational"
title_shape = None
for sh in s.shapes:
    if sh.has_text_frame and sh.text_frame.text.strip().startswith("2. Operational"):
        title_shape = sh
        break
if title_shape is not None:
    p0 = title_shape.text_frame.paragraphs[0]
    for extra in p0.runs[1:]:
        extra._r.getparent().remove(extra._r)
    p0.runs[0].text = "Operational Philosophy: Cold ≠ Off"
# remove the old plain-text hierarchy line and the old 4-bullet baseline list (separate shapes)
for sh in list(s.shapes):
    if sh.has_text_frame and sh is not title_shape:
        txt = sh.text_frame.text
        if ("2K-OP" in txt or "Preferred operational state hierarchy" in txt
                or "expected to operate continuously" in txt
                or "Warm Stop is the last resort" in txt):
            sh._element.getparent().remove(sh._element)
add_bullets(s, Emu(645465), Emu(1650000), Emu(11000000), Emu(1150000), [
    "QPLANT is expected to operate continuously — nominal and standby modes are both valid operating states.",
    "Standby ≠ unavailability. Warm Stop is the last resort, not a routine state.",
], size=14, color_theme=TC.TEXT_1, brightness=0.08, space_after=8)

states = [
    ("2K-OP", "Nominal", "#2E7D32"), ("2K-SB", "Standby, valid", "#7A4C82"),
    ("4.5K-SB", "Reduced avail.", "#C77700"), ("TS-SB", "Warm standby", "#C77700"),
    ("WS", "Last resort", "#C62828"),
]
sx, sy, sw, sh_ = Emu(645465), Emu(3350000), Emu(2100000), Emu(1150000)
gapx = Emu(150000)
prev_right = None
for i, (name, desc, hexcol) in enumerate(states):
    x = sx + i * (sw + gapx)
    box = s.shapes.add_shape(MSO_SHAPE.CHEVRON if 0 < i < 4 else MSO_SHAPE.PENTAGON, x, sy, sw, sh_)
    if 0 < i < 4:
        box.adjustments[0] = 0.2  # shrink the notch so text has more usable width
    box.fill.solid(); box.fill.fore_color.rgb = RGBColor.from_string(hexcol.lstrip('#'))
    box.line.fill.background(); box.shadow.inherit = False
    btf = box.text_frame; btf.word_wrap = True; btf.vertical_anchor = MSO_ANCHOR.MIDDLE
    btf.margin_left = Emu(20000); btf.margin_right = Emu(20000)
    bp = btf.paragraphs[0]; bp.alignment = PP_ALIGN.CENTER; bp.line_spacing = 1.0
    br = bp.add_run(); set_run(br, text=name, size=13, bold=True, rgb=(255, 255, 255))
    bp2 = btf.add_paragraph(); bp2.alignment = PP_ALIGN.CENTER; bp2.line_spacing = 1.0
    br2 = bp2.add_run(); set_run(br2, text=desc, size=9, rgb=(255, 255, 255))
add_takeaway(s, Emu(645465), Emu(4850000), Emu(11000000), Emu(650000), "READING THIS DIAGRAM",
             "Preferred hierarchy runs left → right. Every step right of 2K-OP is a deliberate, monitored "
             "trade-down — not a failure — until Warm Stop, which is.", size=11.5, color=TC.ACCENT_1)
add_tag(s, "DEFINE", style="banner")
print("pos19 done")

# ---- pos20/21: split "Failure Classification" + "Reliability Targets" ---------
# BSLN Slide 16 crammed two RTM items ("3. Failure Classification" and
# "4. Reliability Targets by Failure Class") onto one slide, banner-on-banner,
# with an orphan leftover textbox (the pos19 state-hierarchy note, present in
# the original file too — dead content, not related to either topic). pos21
# was structurally cloned from this same slide so each half gets full room,
# a proper top banner, and (on pos21) the s09 chart reused from Slide 13.

# -- pos20: Failure Classification (keeps the table; everything else removed) --
s = slides[19]
remove_old_tag(s)
banner = find_shape(s, "Rectangle 1")
p0 = banner.text_frame.paragraphs[0]
for extra in p0.runs[1:]:
    extra._r.getparent().remove(extra._r)
p0.runs[0].text = "Failure Classification"
for sh in list(s.shapes):
    if sh.has_text_frame and sh is not banner:
        txt = sh.text_frame.text
        if ("Preferred operational state hierarchy" in txt
                or "Reliability Targets by Failure Class" in txt
                or "MTBF requirements scale" in txt
                or txt.strip().startswith("Class A:")):
            sh._element.getparent().remove(sh._element)
subtitle = find_shape(s, "Rectangle 3")  # "Failure Classification (RTM-034 to RTM-036)"
subtitle.text_frame.paragraphs[0].runs[0].text = "RTM-034 to RTM-036 — three consequence-based classes"
table_shape = None
for sh in s.shapes:
    if sh.has_table:
        table_shape = sh
        break
table_shape.left, table_shape.top = Emu(1928793), Emu(1750000)
tbl = table_shape.table
for ci, cell in enumerate(tbl.rows[0].cells):
    cell.fill.solid(); cell.fill.fore_color.theme_color = TC.ACCENT_1
    for p in cell.text_frame.paragraphs:
        for r in p.runs:
            r.font.bold = True; r.font.color.theme_color = TC.BACKGROUND_1
for ri, row in enumerate(tbl.rows):
    if ri == 0:
        continue
    for cell in row.cells:
        cell.fill.solid()
        if ri % 2 == 0:
            cell.fill.fore_color.theme_color = TC.ACCENT_1; cell.fill.fore_color.brightness = 0.92
        else:
            cell.fill.fore_color.theme_color = TC.BACKGROUND_1
        for p in cell.text_frame.paragraphs:
            for r in p.runs:
                if p.text.strip() in ("A", "B", "C"):
                    r.font.bold = True; r.font.color.theme_color = TC.ACCENT_1
intro = s.shapes.add_textbox(Emu(1928793), Emu(3450000), Emu(8229600), Emu(700000))
itf = intro.text_frame; itf.word_wrap = True
ip = itf.paragraphs[0]; ip.line_spacing = 1.15
ir = ip.add_run()
set_run(ir, text="Every unplanned transition out of 2 K is classified by how far it falls and how long "
              "recovery takes — the classification alone doesn't set a number, it decides which MTBF "
              "target applies (Slide 21).", size=13, color_theme=TC.TEXT_1, brightness=0.08)
add_takeaway(s, Emu(1928793), Emu(4550000), Emu(8229600), Emu(700000), "WHY IT MATTERS",
             "Class is assigned by consequence, not cause — a Class C warm-up from a sensor fault "
             "counts the same as one from a compressor failure. This keeps the MTBF budget honest.",
             size=11.5, color=TC.ACCENT_1)
add_tag(s, "DEFINE", style="banner")
set_notes(s, "BSLN Slide 16 combined this table with the Reliability Targets content below it on one "
              "crowded slide (plus an unrelated leftover textbox duplicated from the Operational "
              "Philosophy slide, removed here as dead content). Split for readability — targets now on "
              "Slide 21.\n\nFull baseline table:\n- Class A: Exit from 2 K, recoverable ≤ 24 h\n"
              "- Class B: Transition to 4.5 K / TS standby\n- Class C: Warm-up > 4.5 K or > 24 h recovery")
print("pos20 done")

# -- pos21: Reliability Targets by Failure Class (repositioned + s09 chart) ----
s = slides[20]
remove_old_tag(s)
banner = find_shape(s, "Rectangle 1")
p0 = banner.text_frame.paragraphs[0]
for extra in p0.runs[1:]:
    extra._r.getparent().remove(extra._r)
p0.runs[0].text = "Reliability Targets by Failure Class"
bullets_box, classvals_box, old_subtitle_box, second_banner = None, None, None, None
for sh in list(s.shapes):
    if not sh.has_text_frame or sh is banner:
        continue
    txt = sh.text_frame.text
    if sh.has_table:
        continue
    if "Preferred operational state hierarchy" in txt:
        sh._element.getparent().remove(sh._element); continue
    if txt.strip() == "Failure Classification (RTM-034 to RTM-036)":
        old_subtitle_box = sh; continue
    if "Reliability Targets by Failure Class" in txt and sh is not banner:
        second_banner = sh; continue
    if "MTBF requirements scale" in txt:
        bullets_box = sh; continue
    if txt.strip().startswith("Class A:"):
        classvals_box = sh; continue
if old_subtitle_box is not None:
    old_subtitle_box._element.getparent().remove(old_subtitle_box._element)
if second_banner is not None:
    second_banner._element.getparent().remove(second_banner._element)
for sh in list(s.shapes):
    if sh.has_table:
        sh._element.getparent().remove(sh._element)

subtitle = s.shapes.add_textbox(Emu(1928793), Emu(1209106), Emu(4617290), Emu(369332))
stf = subtitle.text_frame; stf.word_wrap = True
sp = stf.paragraphs[0]
sr = sp.add_run()
set_run(sr, text="Same three classes as Slide 20 — MTBF target per consequence tier", size=13,
        bold=True, color_theme=TC.ACCENT_2)

bullets_box.left, bullets_box.top, bullets_box.width, bullets_box.height = \
    Emu(1928793), Emu(1900000), Emu(4600000), Emu(1400000)
for p in bullets_box.text_frame.paragraphs:
    for r in p.runs:
        r.font.size = Pt(13)

classvals_box._element.getparent().remove(classvals_box._element)  # redundant w/ chart, baseline highlight looked broken here
note = s.shapes.add_textbox(Emu(1928793), Emu(3450000), Emu(4600000), Emu(500000))
ntf = note.text_frame; ntf.word_wrap = True
np = ntf.paragraphs[0]; np.line_spacing = 1.1
nr = np.add_run()
set_run(nr, text="Target values (> 5 / > 10 / > 15 years) are read directly off the chart, right →",
        size=11.5, italic=True, color_theme=TC.TEXT_1, brightness=0.3)

add_picture_fit(s, f"{CH}/s09_mtbf_by_class.png", Emu(6850000), Emu(1750000), Emu(4700000), Emu(3450000))
cap = s.shapes.add_textbox(Emu(6850000), Emu(5200000), Emu(4700000), Emu(400000))
ctf = cap.text_frame; ctf.word_wrap = True
cp = ctf.paragraphs[0]; cp.alignment = PP_ALIGN.CENTER
cr = cp.add_run()
set_run(cr, text="Same chart, different lens — first shown on Slide 13", size=10.5, italic=True,
        color_theme=TC.TEXT_1, brightness=0.3)
add_takeaway(s, Emu(1928793), Emu(4750000), Emu(4600000), Emu(850000), "CROSS-CHECK",
             "These are the same three numbers Slide 13 feeds into the Poisson campaign-success "
             "model (Slides 22–36) — defined once here, used twice.",
             size=11, color=TC.ACCENT_1)
add_tag(s, "DEFINE", style="banner")
set_notes(s, "Second half of BSLN Slide 16 (see Slide 20's notes for the split rationale).\n\n"
              "Full baseline text:\n- MTBF requirements scale with severity of consequence\n"
              "- Applies to aggregated QPLANT contribution\n"
              "- Prevents trading rare catastrophic events against frequent minor ones\n"
              "- Class A: > 5 years\n- Class B: > 10 years\n- Class C: > 15 years")
print("pos21 done")

# ---- pos5/pos6: light-touch tidy (already-reviewed baseline slides) -----------
def add_head_body_bullets(slide, x, y, cx, cy, pairs, head_size=12.5, body_size=11.5, gap=9):
    box = slide.shapes.add_textbox(x, y, cx, cy)
    tf = box.text_frame; tf.word_wrap = True
    first = True
    for head, body in pairs:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.line_spacing = 1.08; p.space_after = Pt(2)
        r = p.add_run(); set_run(r, text=head, size=head_size, bold=True, color_theme=TC.ACCENT_2)
        p2 = tf.add_paragraph(); p2.line_spacing = 1.1; p2.space_after = Pt(gap)
        r2 = p2.add_run(); set_run(r2, text=body, size=body_size, color_theme=TC.TEXT_1, brightness=0.12)
    return box

# -- pos5: Objectives & Scope (orig2) --
s = slides[4]
remove_old_tag(s)
body, orig_lines = gather_and_wipe_body(s)
body._element.getparent().remove(body._element)
add_head_body_bullets(s, Emu(766763), Emu(2200000), Emu(6100000), Emu(4300000), [
    ("Fleet baseline", "Consolidate compressor fleet assumptions — ≥ 3 units minimum, Option A baseline = 4 × FSD 575."),
    ("RTM alignment", "Align with Addendum II canonical RTMs, incl. 'To be offered…' statements."),
    ("Acceptance & integration", "Define acceptance tests (24×2 K SB & 24×2 K OP) and CIS/MCS integration checkpoints (Slides 12, 14)."),
    ("Operating scope", "24 QM with FFT operating mode — below the 307 g/s nominal design target."),
    ("Design envelope", "320–350 g/s total flow — the envelope every capability check in this deck is measured against (Slide 9)."),
])
tx, ty, tw, th, gap = Emu(7200000), Emu(2200000), Emu(2160000), Emu(1150000), Emu(120000)
add_stat_tile(s, tx, ty, tw, th, "307 g/s", "24 QM nominal\ntarget flow", color=TC.ACCENT_1)
add_stat_tile(s, tx + tw + gap, ty, tw, th, "320–350", "g/s design\nenvelope (Slide 9)", color=TC.ACCENT_1)
add_takeaway(s, Emu(7200000), Emu(3550000), Emu(4440000), Emu(750000), "SCOPE IN ONE LINE",
             "Everything from here to the closing governance loop exists to keep the fleet inside this envelope, provably.",
             size=11, color=TC.ACCENT_1)
add_tag(s, "DEFINE", style="layout")
set_notes(s, "Full baseline bullet text (verbatim), condensed on-slide for readability:\n\n" +
              "\n".join(f"- {l}" for l in orig_lines))
print("pos5 done")

# -- pos6: Configuration Baseline (orig3) --
s = slides[5]
remove_old_tag(s)
title6 = find_shape(s, "Title 1")
for r in title6.text_frame.paragraphs[0].runs:
    r.font.size = Pt(29)  # fixes baseline title-wrap-over-subtitle bug (title too long for 1 line at default size)
body, orig_lines = gather_and_wipe_body(s)
body._element.getparent().remove(body._element)
add_head_body_bullets(s, Emu(766763), Emu(2200000), Emu(6600000), Emu(4300000), [
    ("Fleet power", "4 × 350 kW HP compressors ≈ 1.4 MW total installed."),
    ("Cooling water", "27→ 37°C supply/return (ΔT 10 K) nominal; up to 42°C return (ΔT 15 K) for PS01 at reduced flow."),
    ("Piping", "Shared header, looped feed/return."),
    ("4th-unit deferral", "Fully stubbed install position — enforceable by SCK CEN at grant or EoCD."),
    ("Pressure control", "≈ 14 bar discharge maintained; VFD trims flow (feeds the Unit Economics model, Slide 10)."),
    ("Applicant shall", "Deliver layout & piping stubs if the 4th unit is deferred."),
])
note = s.shapes.add_textbox(Emu(766763), Emu(5900000), Emu(6600000), Emu(400000))
ntf = note.text_frame; ntf.word_wrap = True
np_ = ntf.paragraphs[0]; np_.line_spacing = 1.05
nr = np_.add_run()
set_run(nr, text="Open baseline question: pressure-constant vs. pressure-swing control — not yet decided.",
        size=10.5, italic=True, color_theme=TC.ACCENT_2)
tx, ty, tw, th, gap = Emu(7650000), Emu(2200000), Emu(1330000), Emu(1150000), Emu(90000)
add_stat_tile(s, tx, ty, tw, th, "1.4 MW", "installed\n(4×350 kW)", color=TC.ACCENT_1)
add_stat_tile(s, tx + tw + gap, ty, tw, th, "14 bar", "discharge\npressure", color=TC.ACCENT_1)
add_stat_tile(s, tx + 2*(tw + gap), ty, tw, th, "10–15 K", "cooling\nΔT range", color=TC.ACCENT_1)
add_takeaway(s, Emu(7650000), Emu(3550000), Emu(4170000), Emu(950000), "WHY IT MATTERS",
             "This power/water/pressure envelope is the input to the Unit Economics energy-cost model (Slide 10) and the Utilities hub-and-spoke diagram (Slide 15).",
             size=10.5, color=TC.ACCENT_1)
add_tag(s, "DEFINE", style="layout")
set_notes(s, "Full baseline bullet text (verbatim), condensed on-slide for readability:\n\n" +
              "\n".join(f"- {l}" for l in orig_lines))
print("pos6 done")

# ---- deck-wide pass: recolor any remaining Phase-1-era flat-purple tags -------
# Phase 1's add_tag always used flat ACCENT_1 regardless of phase. The ANALYZE
# section (blue) and CONTROL section (magenta) tags still carry that flat
# purple; DEFINE/MEASURE/IMPROVE already resolve to ACCENT_1 under the new
# PHASE_COLOR system so this is a no-op for them. Runs deck-wide, idempotent.
tag_fix_count = 0
for _s in slides:
    for _sh in list(_s.shapes):
        if _sh.has_text_frame:
            _t = _sh.text_frame.text.strip()
            if _t in OLD_TAGS:
                _correct = PHASE_COLOR[_t]
                for _p in _sh.text_frame.paragraphs:
                    for _r in _p.runs:
                        try:
                            if _r.font.color.type is not None and _r.font.color.theme_color != _correct:
                                _r.font.color.theme_color = _correct
                                tag_fix_count += 1
                        except Exception:
                            pass
print(f"deck-wide tag recolor pass: {tag_fix_count} runs fixed")

# ---- fix stale cross-references in the Phase-1 narrative slides ---------------
# These 10 slides (pos37-46) were authored in Phase 1 against the OLD 42-slide
# numbering. Phase 2 inserted 4 slides ahead of them (2 extra roadmap slides,
# the new Unit Economics slide, and the slide-16 split), shifting every
# absolute slide reference inside their body text by +4. Fixed by direct string
# replacement since each is a single run.
_xref_fixes = [
    (37, "Slides 18–32", "Slides 22–36"),
    (38, "Slide 32's", "Slide 36's"),
    (44, "Slides 18–32", "Slides 22–36"),
    (45, "Slide 33:", "Slide 37:"),
    (45, "Slide 5).", "Slide 7)."),
    (45, "Slide 10).", "Slide 13)."),
]
for _idx, _old, _new in _xref_fixes:
    _fixed = False
    for _sh in slides[_idx].shapes:
        if _sh.has_text_frame and _old in _sh.text_frame.text:
            for _p in _sh.text_frame.paragraphs:
                for _r in _p.runs:
                    if _old in _r.text:
                        _r.text = _r.text.replace(_old, _new)
                        _fixed = True
    if not _fixed:
        print(f"WARNING: cross-ref fix not applied — pos{_idx+1} '{_old}' not found")
print("narrative cross-reference fixes applied")

prs.save(OUT)
print(f"saved {OUT}")
