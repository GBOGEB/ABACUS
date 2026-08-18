#!/usr/bin/env python3
"""Pass 3: replace the 15-slide screenshot-heavy old Analyze block (Slides
22-36) with 7 consolidated, fully-editable Appendix D-I slides. Input:
final2.pptx (Phase 2 output). Output: final3.pptx."""
from pptx import Presentation
from pptx.util import Emu, Pt, Inches
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_THEME_COLOR_INDEX as TC
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
import copy as _copy

SRC = "final2.pptx"
OUT = "final3.pptx"
CH = "/home/claude/work/charts"

prs = Presentation(SRC)

# ---------------------------------------------------------------- helpers --

def set_run(run, text=None, size=None, bold=None, italic=None, color_theme=None,
            brightness=None, rgb=None, font=None):
    if text is not None:
        run.text = text
    f = run.font
    if size is not None:
        f.size = Pt(size)
    if bold is not None:
        f.bold = bold
    if italic is not None:
        f.italic = italic
    if font is not None:
        f.name = font
    if color_theme is not None:
        f.color.theme_color = color_theme
        if brightness is not None:
            f.color.brightness = brightness
    elif rgb is not None:
        f.color.rgb = RGBColor(*rgb)

def set_notes(slide, text):
    ns = slide.notes_slide
    tf = ns.notes_text_frame
    tf.text = text

def find_shape(slide, name):
    for sh in slide.shapes:
        if sh.name == name:
            return sh
    return None

PHASE_COLOR = {
    "DEFINE": TC.ACCENT_1, "MEASURE": TC.ACCENT_2, "ANALYZE": TC.ACCENT_4,
    "IMPROVE": TC.ACCENT_1, "CONTROL": TC.ACCENT_2,
}

def add_tag(slide, label, style="layout"):
    color = PHASE_COLOR.get(label, TC.ACCENT_1)
    if style == "layout":
        x, y, cx, cy = Emu(9950000), Emu(160000), Emu(1750000), Emu(320000)
    else:
        x, y, cx, cy = Emu(10760000), Emu(300000), Emu(1330000), Emu(340000)
    box = slide.shapes.add_textbox(x, y, cx, cy)
    tf = box.text_frame
    tf.word_wrap = False
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    r = p.add_run()
    set_run(r, text=label, size=10.5, bold=True, color_theme=color, font="Calibri")
    r.font._rPr.set('spc', '120')
    return box

def add_takeaway(slide, x, y, cx, cy, label, text, size=12, color=TC.ACCENT_1):
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, cx, cy)
    box.adjustments[0] = 0.06
    box.fill.solid()
    box.fill.fore_color.theme_color = color
    box.fill.fore_color.brightness = 0.90
    box.line.fill.background()
    box.shadow.inherit = False
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(137160); tf.margin_right = Emu(137160)
    tf.margin_top = Emu(80000); tf.margin_bottom = Emu(80000)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.line_spacing = 1.1
    r1 = p.add_run(); set_run(r1, text=label + "  ", size=size, bold=True, color_theme=color)
    r2 = p.add_run(); set_run(r2, text=text, size=size, bold=False, color_theme=color)
    return box

def add_bullets(slide, x, y, cx, cy, items, size=13, color_theme=TC.TEXT_1,
                 brightness=0.1, space_after=8, line_spacing=1.12):
    box = slide.shapes.add_textbox(x, y, cx, cy)
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.line_spacing = line_spacing
        p.space_after = Pt(space_after)
        pPr = p._p.get_or_add_pPr()
        pPr.set('marL', '182880'); pPr.set('indent', '-182880')
        ns_ = '{http://schemas.openxmlformats.org/drawingml/2006/main}'
        pPr.append(pPr.makeelement(f'{ns_}buFont', {'typeface': 'Arial'}))
        pPr.append(pPr.makeelement(f'{ns_}buChar', {'char': '•'}))
        r = p.add_run()
        set_run(r, text=item, size=size, color_theme=color_theme, brightness=brightness)
    return box

def add_head_body(slide, x, y, cx, cy, pairs, head_size=12.5, body_size=11,
                   head_color=TC.ACCENT_2, gap=8):
    box = slide.shapes.add_textbox(x, y, cx, cy)
    tf = box.text_frame; tf.word_wrap = True
    first = True
    for head, body in pairs:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.line_spacing = 1.05; p.space_after = Pt(1)
        r = p.add_run(); set_run(r, text=head, size=head_size, bold=True, color_theme=head_color)
        p2 = tf.add_paragraph(); p2.line_spacing = 1.08; p2.space_after = Pt(gap)
        r2 = p2.add_run(); set_run(r2, text=body, size=body_size, color_theme=TC.TEXT_1, brightness=0.12)
    return box

def add_picture_fit(slide, path, x, y, cx, cy):
    from PIL import Image
    im = Image.open(path)
    iw, ih = im.size
    scale = min(cx / iw, cy / ih)
    w, h = int(iw * scale), int(ih * scale)
    ox = x + (cx - w) // 2
    oy = y + (cy - h) // 2
    return slide.shapes.add_picture(path, ox, oy, w, h)

def add_stat_tile(slide, x, y, cx, cy, value, label, color=TC.ACCENT_1, value_size=18):
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, cx, cy)
    box.adjustments[0] = 0.08
    box.fill.solid(); box.fill.fore_color.theme_color = color; box.fill.fore_color.brightness = 0.92
    box.line.color.theme_color = color; box.line.width = Pt(0.75)
    box.shadow.inherit = False
    tf = box.text_frame; tf.word_wrap = True
    tf.margin_left = Emu(90000); tf.margin_right = Emu(90000)
    tf.margin_top = Emu(60000); tf.margin_bottom = Emu(60000)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER; p.line_spacing = 1.0
    r = p.add_run(); set_run(r, text=value, size=value_size, bold=True, color_theme=color)
    p2 = tf.add_paragraph(); p2.alignment = PP_ALIGN.CENTER; p2.line_spacing = 1.05
    r2 = p2.add_run(); set_run(r2, text=label, size=10, color_theme=TC.TEXT_1, brightness=0.2)
    return box

def style_title_subtitle(slide, title_text, subtitle_text, banner_theme=TC.ACCENT_4, title_size=27):
    t = find_shape(slide, "Title 1")
    t.left, t.top, t.width, t.height = Emu(645465), Emu(700000), Emu(10900000), Emu(900000)
    t.text_frame.word_wrap = True
    tp = t.text_frame.paragraphs[0]
    for extra in tp.runs[1:]:
        extra._r.getparent().remove(extra._r)
    if not tp.runs:
        tp.add_run()
    tp.runs[0].text = title_text
    tp.runs[0].font.size = Pt(title_size)
    tp.runs[0].font.color.theme_color = banner_theme
    sub = find_shape(slide, "Text Placeholder 2")
    sub.left, sub.top, sub.width, sub.height = Emu(645465), Emu(1620000), Emu(10900000), Emu(460000)
    sp = sub.text_frame.paragraphs[0]
    for extra in sp.runs[1:]:
        extra._r.getparent().remove(extra._r)
    if not sp.runs:
        sp.add_run()
    sp.runs[0].text = subtitle_text
    sp.runs[0].font.size = Pt(14)
    sp.runs[0].font.color.theme_color = banner_theme
    return t, sub

def wipe_body(slide):
    body = find_shape(slide, "Text Placeholder 3")
    if body is None:
        return None
    # text_frame.clear() keeps one empty paragraph (valid OOXML requires >=1
    # <a:p> child) instead of removing every paragraph, which PowerPoint
    # rejects as corrupt.
    body.text_frame.clear()
    return body

def add_table(slide, x, y, cx, cy, headers, rows, col_widths=None, header_color=TC.ACCENT_1,
              font_size=10.5, header_size=10.5, row_band=True):
    n_rows = len(rows) + 1
    n_cols = len(headers)
    gshape = slide.shapes.add_table(n_rows, n_cols, x, y, cx, cy)
    tbl = gshape.table
    if col_widths:
        total = sum(col_widths)
        for i, w in enumerate(col_widths):
            tbl.columns[i].width = Emu(int(cx * w / total))
    for ci, htext in enumerate(headers):
        cell = tbl.cell(0, ci)
        cell.fill.solid(); cell.fill.fore_color.theme_color = header_color
        cell.margin_left = Emu(45720); cell.margin_right = Emu(45720)
        cell.margin_top = Emu(20000); cell.margin_bottom = Emu(20000)
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf = cell.text_frame; tf.word_wrap = True
        p = tf.paragraphs[0]
        r = p.add_run(); set_run(r, text=htext, size=header_size, bold=True, color_theme=TC.BACKGROUND_1)
    for ri, row in enumerate(rows, start=1):
        for ci, val in enumerate(row):
            cell = tbl.cell(ri, ci)
            cell.margin_left = Emu(45720); cell.margin_right = Emu(45720)
            cell.margin_top = Emu(15000); cell.margin_bottom = Emu(15000)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            cell.fill.solid()
            if row_band and ri % 2 == 0:
                cell.fill.fore_color.theme_color = header_color; cell.fill.fore_color.brightness = 0.92
            else:
                cell.fill.fore_color.theme_color = TC.BACKGROUND_1
            tf = cell.text_frame; tf.word_wrap = True
            p = tf.paragraphs[0]
            r = p.add_run(); set_run(r, text=str(val), size=font_size, color_theme=TC.TEXT_1, brightness=0.1)
    return tbl

def delete_slide_by_index(prs_, index):
    xml_slides = prs_.slides._sldIdLst
    slides = list(xml_slides)
    rId = slides[index].rId
    prs_.part.drop_rel(rId)
    xml_slides.remove(slides[index])

def move_block_to(prs_, start_index, count, dest_index):
    xml_slides = prs_.slides._sldIdLst
    slides = list(xml_slides)
    block = slides[start_index:start_index + count]
    for e in block:
        xml_slides.remove(e)
    for i, e in enumerate(block):
        xml_slides.insert(dest_index + i, e)

print("helpers3 loaded OK")

# ============================================================ STEP 1: ADD 7 ==
# Added BEFORE deletion so python-pptx assigns fresh, non-colliding slide part
# numbers (it fills the lowest free gap, which would otherwise collide with
# kept slides 37/38 once 22-36 are freed up).
TEXT_LAYOUT = prs.slide_masters[0].slide_layouts[12]
assert TEXT_LAYOUT.name == "Text Slide", TEXT_LAYOUT.name
new_slides = []
for _ in range(7):
    s = prs.slides.add_slide(TEXT_LAYOUT)
    new_slides.append(s)
print(f"added 7 new slides, deck now {len(prs.slides)}")

# ============================================================ STEP 2: DELETE =
# old pos22-36 (0-indexed 21-35), the screenshot-heavy Analyze block.
n_before = len(prs.slides)
for idx in range(35, 20, -1):
    delete_slide_by_index(prs, idx)
print(f"deleted 15 slides: {n_before} -> {len(prs.slides)}")

# move them from the end (indices 31-37) to right after old pos21 (index 21..27)
move_block_to(prs, len(prs.slides) - 7, 7, 21)
print("moved appendix slides into position 22-28")

slides = list(prs.slides)
D1, D2, E, F, G, H, I = slides[21:28]

# ============================================================ APPENDIX D (1/2) =
style_title_subtitle(D1, "Appendix D (1/2) — Poisson Reliability Model",
                      "Methodology behind every campaign-success number in this compendium")
wipe_body(D1)
add_head_body(D1, Emu(645465), Emu(2260000), Emu(5350000), Emu(3200000), [
    ("The model", "Failures are modeled as a Poisson process with constant rate λ (events/year) — "
                   "equivalent to assuming a constant hazard rate, i.e. MTBF = 1/λ."),
    ("What the chart shows", "P₀ = probability of zero trips within a time window, as MTBF varies. "
                              "Three windows are plotted: the 90-day campaign (the operational requirement), "
                              "1 year, and 5 years (long-term intuition check)."),
    ("Formula", "P(N=0) = e^(−λt),  with  t in years and  λ = 1 / MTBF."),
], gap=10)
add_takeaway(D1, Emu(645465), Emu(5750000), Emu(11000000), Emu(620000), "WHERE THIS GETS USED",
             "Appendix D (2/2) applies this to three specific cases; Appendix F uses it to derive "
             "RTM-05's 0.26 events/year limit; Appendix I applies the same logic with a Weibull "
             "(age-dependent) hazard instead of a constant one.",
             size=10.5, color=TC.ACCENT_4)
add_picture_fit(D1, f"{CH}/appendixD_poisson.png", Emu(6250000), Emu(2100000), Emu(5350000), Emu(3450000))
add_tag(D1, "ANALYZE", style="layout")
set_notes(D1, "This slide replaces 5 near-identical Poisson-curve screenshots from the original "
              "working deck (all plotting the same P0=e^-lambda*t relationship with progressively more "
              "annotation layered on) with one clean chart and the underlying formula stated once. "
              "Full original chart axis/curve legend: X-axis = trips per year (lambda), the Poisson "
              "event rate, equivalent to MTBF = 1/lambda. Y-axis = probability of zero trips over a "
              "given time window. Curves = 90-day campaign (physics continuity requirement), 1 year "
              "(operational intuition), 5 years (long-term cold-operation intuition). All curves follow "
              "P(N=0) = exp(-lambda*t).")
print("Appendix D (1/2) done")

# ============================================================ APPENDIX D (2/2) =
style_title_subtitle(D2, "Appendix D (2/2) — Worked Examples",
                      "Three specific readings of the same P(N=0) = e^(−λt) relationship")
wipe_body(D2)
add_head_body(D2, Emu(645465), Emu(2260000), Emu(11000000), Emu(950000), [
    ("Why three examples", "The same curve answers three different questions this compendium has to "
                            "settle: what does the fleet realistically deliver, what does the contract "
                            "require, and where does a naive reading of the requirement break down."),
], gap=6)

ex_y = Emu(3300000)
ex_h = Emu(950000)
ex_gap = Emu(100000)
examples = [
    ("1 — REALISTIC CASE", TC.ACCENT_4,
     "CC + PVPS single-train, λ ≈ 0.25/yr (MTBF ≈ 4 y): 90-day P₀ ≈ 94%, 1-year ≈ 78%, 5-year ≈ 29%. "
     "Consistent with MTBF ≈ 4 y and “good but not magical” reliability — see Appendix I for where "
     "this number comes from."),
    ("2 — CONTRACTUAL CASE", TC.ACCENT_1,
     "Class-A target MTBF = 5 y (λ = 0.20/yr): 90-day P₀ ≈ 95% (RTM-031 compliant), 1-year ≈ 82%, "
     "5-year ≈ 37%. This is the number RTM-05 (Appendix F) is built on."),
    ("3 — A COMMON MISREADING", TC.ACCENT_2,
     "Reading “≥99% success” as a literal requirement over the full mission window (not per 90-day "
     "campaign) implies MTBF ≈ 387 years — impossible for a single-train cryoplant. This is why RTM-05 "
     "is written as an annual event-rate limit, not a lifetime success probability."),
]
for i, (label, color, body) in enumerate(examples):
    y = ex_y + i * (ex_h + ex_gap)
    add_takeaway(D2, Emu(645465), y, Emu(11000000), ex_h, label, body, size=10.5, color=color)
add_tag(D2, "ANALYZE", style="layout")
set_notes(D2, "Consolidates three separate worked-example screenshots from the original working deck "
              "(the CC+PVPS realistic case, the contractual Class-A case, and the '100 occurrences' / "
              "'99% over 90 days' literal-reading walkthrough) into three compact takeaway boxes. Full "
              "original derivation for example 3: mission time ~46.7 months = ~34,090 hours; setting "
              "P(>=1 failure) = 1% over that window and solving for MTBF gives ~3.39 million hours "
              "(~387 years) -- confirming that a literal 99%-over-full-window reading is not the intended "
              "requirement, which is why RTM-05 (Appendix F) instead caps the annual SAE rate.")
print("Appendix D (2/2) done")

# ============================================================ APPENDIX E ======
style_title_subtitle(E, "Appendix E — Maintenance Activity Matrix",
                      "What's permitted in each operating state (states defined on Slide 19)")
wipe_body(E)
add_bullets(E, Emu(645465), Emu(2100000), Emu(11000000), Emu(500000), [
    "Same five states as the Slide 19 chevron diagram and the Slide 20 failure-consequence classes — "
    "this table is what makes those states operationally concrete.",
], size=11.5, space_after=0)

GOOD_RGB = (46, 125, 50); WARN_RGB = (199, 119, 0); CRIT_RGB = (198, 40, 40)
MATRIX_HEADERS = ["Maintenance activity", "2K-OP", "2K-SB", "4.5K-SB", "TS-SB", "WS", "Notes"]
MATRIX_ROWS = [
    ("Monitoring / diagnostics", "G", "G", "G", "G", "G", "No operational impact"),
    ("Software / control updates (non-intrusive)", "W", "G", "G", "G", "G", "With rollback + MOC"),
    ("Instrument checks (non-isolating)", "W", "G", "G", "G", "G", "No venting or isolation"),
    ("Helium inventory adjustment", "W", "G", "G", "G", "G", "No pressure excursion"),
    ("Partial valve tests", "W", "G", "G", "G", "G", "Proven transient envelope"),
    ("Rotating machinery minor service", "C", "W", "G", "G", "G", "Depends on redundancy"),
    ("Cold compressor interventions", "C", "C", "W", "C", "G", "Below TS temperature"),
    ("Cryogenic valve intrusive work", "C", "C", "W", "C", "G", "Boundary opening"),
    ("Safety valve / pressure tests", "C", "C", "C", "C", "G", "Warm Stop only"),
    ("Cryostat / cold box opening", "C", "C", "C", "C", "G", "Warm Stop only"),
]
ICON = {"G": "✓", "W": "⚠", "C": "✕"}
ICON_RGB = {"G": GOOD_RGB, "W": WARN_RGB, "C": CRIT_RGB}

tbl_shape = E.shapes.add_table(len(MATRIX_ROWS) + 1, len(MATRIX_HEADERS),
                                Emu(645465), Emu(2650000), Emu(11000000), Emu(3350000))
tbl = tbl_shape.table
col_w = [3050000, 950000, 950000, 950000, 950000, 950000, 2250000]
for i, w in enumerate(col_w):
    tbl.columns[i].width = Emu(w)
for ci, htext in enumerate(MATRIX_HEADERS):
    cell = tbl.cell(0, ci)
    cell.fill.solid(); cell.fill.fore_color.theme_color = TC.ACCENT_1
    cell.margin_top = Emu(20000); cell.margin_bottom = Emu(20000)
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = cell.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.CENTER if ci else PP_ALIGN.LEFT
    r = p.add_run(); set_run(r, text=htext, size=10.5, bold=True, color_theme=TC.BACKGROUND_1)
for ri, row in enumerate(MATRIX_ROWS, start=1):
    for ci, val in enumerate(row):
        cell = tbl.cell(ri, ci)
        cell.margin_top = Emu(12000); cell.margin_bottom = Emu(12000)
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        cell.fill.solid()
        cell.fill.fore_color.theme_color = TC.ACCENT_1 if ri % 2 == 0 else TC.BACKGROUND_1
        if ri % 2 == 0:
            cell.fill.fore_color.brightness = 0.94
        p = cell.text_frame.paragraphs[0]
        r = p.add_run()
        if ci in (1, 2, 3, 4, 5):
            p.alignment = PP_ALIGN.CENTER
            set_run(r, text=ICON[val], size=12, bold=True, rgb=ICON_RGB[val])
        else:
            set_run(r, text=val, size=9.5, color_theme=TC.TEXT_1, brightness=0.1)

legend = E.shapes.add_textbox(Emu(645465), Emu(6120000), Emu(11000000), Emu(300000))
ltf = legend.text_frame; ltf.word_wrap = True
lp = ltf.paragraphs[0]
for sym, col, txt in [("✓", GOOD_RGB, " permitted    "), ("⚠", WARN_RGB, " conditional    "), ("✕", CRIT_RGB, " not permitted")]:
    r = lp.add_run(); set_run(r, text=sym, size=11, bold=True, rgb=col)
    r2 = lp.add_run(); set_run(r2, text=txt, size=10.5, color_theme=TC.TEXT_1, brightness=0.3)
add_tag(E, "ANALYZE", style="layout")
set_notes(E, "Recreated as an editable table from the original screenshot (same data, unchanged). "
              "Legend: check = permitted, triangle = conditional (see Notes column), cross = not "
              "permitted. Columns are the same five operating states as the Slide 19 diagram: "
              "2K-OP (nominal), 2K-SB (standby, valid), 4.5K-SB (reduced availability), TS-SB (warm "
              "standby), WS (Warm Stop, last resort).")
print("Appendix E done")

# ============================================================ APPENDIX F ======
style_title_subtitle(F, "Appendix F — Deriving the SAE Frequency Limit",
                      "How the 90-day continuity requirement becomes RTM-05's annual rate cap")
wipe_body(F)
add_head_body(F, Emu(645465), Emu(2260000), Emu(5350000), Emu(4300000), [
    ("1. Continuity requirement", "QPLANT shall support ≥ 90 consecutive days at 2 K, with no unplanned "
                                    "operational-state transition — the user-relevant continuity metric."),
    ("2. Model it as Poisson", "Class-A Service-Affecting Events (SAE) are modeled as a Poisson process, "
                                 "constant rate λ. For a 90-day window, t = 90/365 = 0.2466 y, and "
                                 "P(zero SAE) = e^(−λt) (Appendix D)."),
    ("3. Invert for the rate cap", "Requiring P(zero SAE) ≥ 0.95 over 90 days gives "
                                     "λ ≤ −ln(0.95)/0.2466 ≈ 0.208 SAE/year, i.e. MTBF_A ≥ 4.8 years — "
                                     "not an unrealistic number, once read correctly."),
    ("4. Duty-factor conversion", "SAE statistics are gathered during active 2 K time. For the nominal "
                                    "duty factor d = 9/14, calendar time = active 2 K time / d — this is "
                                    "how a per-active-time rate becomes a calendar-year rate."),
], gap=8, head_size=12, body_size=10.5)

add_takeaway(F, Emu(6250000), Emu(2260000), Emu(5350000), Emu(1150000), "A COMMON MISREADING",
             "Reading the success target as ≤1% failure probability over the entire multi-year mission "
             "(not per 90-day campaign) implies MTBF ≈ 387 years — see Appendix D (2/2), example 3. "
             "RTM-05 is deliberately written as an annual rate, not a lifetime probability.",
             size=10.5, color=TC.ACCENT_2)
add_takeaway(F, Emu(6250000), Emu(3540000), Emu(5350000), Emu(1500000), "RTM-05 (NORMATIVE)",
             "The QPLANT design shall limit the frequency of Service-Affecting Events (SAE) attributable "
             "to QPLANT — unplanned transitions out of 2 K Operation — to λ_SAE ≤ 0.26 events per "
             "calendar year, averaged over representative operating periods. ≈ one SAE every 46.7 "
             "calendar months.",
             size=10.5, color=TC.ACCENT_4)
add_takeaway(F, Emu(6250000), Emu(5170000), Emu(5350000), Emu(1150000), "ANNEX X — METHOD",
             "Compliance is assessed with the same Poisson event-rate model (Appendix D), applied to the "
             "declared operational duty cycle. Worked check: λ ≈ 12/46.7 ≈ 0.257 SAE/yr → "
             "P(N=0, 90 d) ≈ e^(−0.257×0.2466) ≈ 0.94 — satisfies RTM-05a.",
             size=10.5, color=TC.ACCENT_1)
add_tag(F, "ANALYZE", style="layout")
set_notes(F, "Condenses 6 screenshot-heavy slides from the original working deck (the '100 occurrences' "
              "scratch derivation, the 'Wrong' literal-reading demonstration, the 'Better' consistency-"
              "rule derivation, the RTM-05 text -- repeated 3 times in the original -- and Annex X) into "
              "one slide with the logical chain plus the two normative end-states. No numbers were "
              "changed: RTM-05's 0.26 events/calendar-year cap, the ~4.8-year equivalent MTBF, the 9/14 "
              "duty factor, and the Annex X worked check (lambda ~0.257/yr, P0~94%) all match the "
              "original document verbatim.")
print("Appendix F done")

# ============================================================ APPENDIX G ======
style_title_subtitle(G, "Appendix G — Component Reliability Reference",
                      "MTBF / MTTR by component, from CERN feedback and published literature")
wipe_body(G)
G_HEADERS = ["Equipment", "MTBF (h)", "Source", "Warm-up (h)", "MTTR (h)", "Recover (h)", "MDT (h)", "MTTR comment"]
G_ROWS = [
    ("Oil screw compressors", "333,450", "CERN + other projects", "0", "20", "6", "26", "Bare compressor available"),
    ("Turbines", "150,000", "REx AL-AT", "0", "8", "3", "11", "Spare cartridge available"),
    ("Cold compressors", "105,000", "REx AL-AT", "0", "8", "3", "11", "Spare cartridge available"),
    ("Control / on-off valves", "824,487", "[1] CERN", "0", "3", "3", "6", "Spare part kits available"),
    ("Cryovalves (positioner only)", "2,298,333", "[1] CERN", "0", "3", "3", "6", "Spare part kits available"),
    ("Cryovalves (tightness + positioner)", "482,927", "[1] CERN", "12", "3", "12", "27", "Spare part kits available"),
    ("Pressure sensors", "2,035,667", "[1] CERN", "0", "2", "2", "4", "Spare transmitters available"),
    ("Temperature sensors", "1,250,000", "EIReDA 98, p.239", "0", "2", "2", "4", "Spare transmitter available"),
]
add_table(G, Emu(645465), Emu(2180000), Emu(11000000), Emu(3350000), G_HEADERS, G_ROWS,
          col_widths=[19, 9, 13, 8, 7, 8, 7, 16], header_color=TC.ACCENT_1, font_size=9.5, header_size=9.5)
foot = G.shapes.add_textbox(Emu(645465), Emu(5600000), Emu(11000000), Emu(280000))
ftf = foot.text_frame; ftf.word_wrap = True
fp = ftf.paragraphs[0]
fr = fp.add_run()
set_run(fr, text="[1] CERN — \"First assessment of reliability data for the LHC accelerator and detector "
              "cryogenic system components\", Perinić et al. Impurities analyzers and the vacuum system "
              "are not listed: not impacting cryo availability.",
        size=9.5, italic=True, color_theme=TC.TEXT_1, brightness=0.3)
add_takeaway(G, Emu(645465), Emu(5980000), Emu(11000000), Emu(420000), "WHERE THESE NUMBERS FEED IN",
             "Oil screw compressor MTBF drives the Case File on Slide 32; cold compressor MTBF drives "
             "the cold-compressor-train calculation on Appendix I; turbine MTBF is cross-checked against "
             "OEM data on Appendix H.",
             size=10, color=TC.ACCENT_1)
add_tag(G, "ANALYZE", style="layout")
set_notes(G, "Recreated as an editable table from a screenshot of the original source table (all values "
              "unchanged). This is the base reference table — component MTBF/MTTR figures used everywhere "
              "else in this compendium trace back to this table.")
print("Appendix G done")

# ============================================================ APPENDIX H ======
style_title_subtitle(H, "Appendix H — Turbine OEM & Contractor Data",
                      "Linde Kryotechnik TED-series field data, and a check against Contractor A's claim",
                      title_size=27)
wipe_body(H)
add_bullets(H, Emu(645465), Emu(2100000), Emu(6900000), Emu(650000), [
    "Linde Kryotechnik AG field data, TED-series turbines in operation 2009–2019 (improved-design "
    "cartridges only) — cross-checks the 150,000 h turbine planning figure used on Appendix G.",
], size=11, space_after=0)
H_HEADERS = ["Turbine type", "In service since 2009", "Operating hours", "Damages (turbine-related)", "MTBF (h)"]
H_ROWS = [
    ("TED16", "182", "3,876,000", "11", "352,364"),
    ("TED22", "94", "1,797,000", "6", "299,500"),
    ("TED32", "23", "1,050,000", "2", "525,000"),
    ("TED45", "21", "833,000", "1", "833,000"),
]
add_table(H, Emu(645465), Emu(2820000), Emu(6900000), Emu(1650000), H_HEADERS, H_ROWS,
          col_widths=[16, 21, 21, 21, 14], header_color=TC.ACCENT_4, font_size=10, header_size=10)
add_bullets(H, Emu(645465), Emu(4600000), Emu(6900000), Emu(500000), [
    "MTTR by spare-cartridge exchange: 2–3 h. By express repair at manufacturer site: 2 weeks. "
    "By standard repair at manufacturer site: 4 weeks (all four types).",
], size=10, space_after=0)

add_head_body(H, Emu(7750000), Emu(2100000), Emu(3900000), Emu(3200000), [
    ("Contractor A's claim", "≥96% availability with maintenance intervals after 8,400 operating hours "
                              "(1 year). Planned time 8,400 h/yr ⇒ allowed unavailability = 4% = 336 h/yr "
                              "≈ 14 days/yr."),
    ("This deck's RTM-032", "If read as 99% over 8,000 h: allowed unavailability = 1% = 80 h/yr "
                             "≈ 3.3 days/yr — about 4× stricter than Contractor A's figure, and stricter "
                             "still against a full 8,760 h year."),
], head_size=11, body_size=10)
add_takeaway(H, Emu(7750000), Emu(5450000), Emu(3900000), Emu(1050000), "READING THIS RIGHT",
             "The gap isn't a contradiction — it means RTM-032 should be read either with a lower annual "
             "% target, a reduced planned-time basis, or (the best fit with what users actually feel) "
             "moderate annual availability paired with a stringent 2 K campaign-success requirement "
             "(Appendix F).",
             size=10, color=TC.ACCENT_2)
add_tag(H, "ANALYZE", style="layout")
set_notes(H, "Recreated as editable tables/text from Linde Kryotechnik AG document screenshots and the "
              "original 'Contrast with Contractor A' analysis (all figures unchanged). The comparison is "
              "presented as a reading-clarification, not a contradiction, per instruction to not let "
              "appendix content contradict the rest of the deck.")
print("Appendix H done")

# ============================================================ APPENDIX I ======
style_title_subtitle(I, "Appendix I — RCM/Weibull Applied: Cold-Compressor Train",
                      "The worked example behind Slide 31's wear-out argument and Roadmap 2/3's sample calc")
wipe_body(I)
add_head_body(I, Emu(645465), Emu(2260000), Emu(5350000), Emu(2650000), [
    ("The system", "3 cold compressors in series, no redundancy — losing one means the plant cannot "
                    "maintain 27 mbar. A pure series system: λ_train = 3 × λ_stage."),
    ("Train MTBF", "Single-stage MTBF = 105,000 h (Appendix G) ⇒ "
                    "MTBF_train = 105,000 / 3 ≈ 35,000 h ≈ 4.0 years — the number used for Class A "
                    "failure frequency and RTM-030/031 aggregation."),
    ("Why Weibull, not just MTBF", "Slide 31 explains the general relation h(t) = (β/η)·(t/η)^(β−1): "
                                     "whenever β > 1, hazard rises with age. The table (right) applies it "
                                     "to real components instead of restating the formula."),
], gap=8, head_size=12, body_size=10.5)
add_takeaway(I, Emu(645465), Emu(5150000), Emu(5350000), Emu(950000), "MTBF_CC-TRAIN ≈ 35,000 h ≈ 4.0 y",
             "This is the correct number for Class A failure frequency, RTM-030 aggregation, and the "
             "RTM-031 90-day campaign probability — not the single-compressor 105,000 h figure.",
             size=11, color=TC.ACCENT_4)

I_HEADERS = ["Equipment", "Weibull (β, η)", "P(fail 90d) good-as-new", "P(fail 90d) at age 5y", "Increase"]
I_ROWS = [
    ("Oil screw compressor", "β=2.5, η≈363,300 h", "2.7×10⁻⁶ (0.00027%)", "6.45×10⁻⁴ (0.0645%)", "×237"),
    ("Turbine", "β=2.0, η≈169,300 h", "1.63×10⁻⁴ (0.0163%)", "6.74×10⁻³ (0.674%)", "×41"),
    ("Cold compressor", "β=2.5, η≈114,400 h", "4.90×10⁻⁵ (0.0049%)", "1.15×10⁻² (1.153%)", "×235"),
]
add_table(I, Emu(6250000), Emu(2260000), Emu(5350000), Emu(1450000), I_HEADERS, I_ROWS,
          col_widths=[18, 15, 20, 20, 14], header_color=TC.ACCENT_4, font_size=8.5, header_size=8.5)
add_picture_fit(I, f"{CH}/appendixI_weibull_increase.png", Emu(6250000), Emu(3850000), Emu(5350000), Emu(2250000))
add_tag(I, "ANALYZE", style="layout")
set_notes(I, "Consolidates two slides from the original working deck (the cold-compressor-train series "
              "calculation and the Weibull 'fresh vs 5-years-old' results table, which was duplicated "
              "across two slides in the original) into one. All figures unchanged. Interpretation "
              "unchanged: if the 90-day 2 K run starts right after a good-as-new overhaul, the wear-out "
              "contribution is tiny; if it starts at the end of the 5-year interval, the risk is "
              "dramatically higher, especially for beta=2.5 assets.")
print("Appendix I done")

# ============================================================ CROSS-REF FIX ===
# Consolidating 15 slides into 7 shifts every slide from old-pos37 onward by -8,
# and old-pos22-36 is now old-pos22-28. Fix every "Slide N" reference that
# points at the block we just rebuilt or at slides that moved.

def replace_run_text(slide, shape_name, old, new):
    sh = find_shape(slide, shape_name)
    assert sh is not None, f"shape {shape_name!r} not found on slide"
    found = False
    for p in sh.text_frame.paragraphs:
        for r in p.runs:
            if old in r.text:
                r.text = r.text.replace(old, new)
                found = True
    assert found, f"text {old!r} not found in {shape_name!r}"

slides = list(prs.slides)

# pos13 (index 12) and pos21 (index 20) -- untouched slides that still cite the
# old 15-slide ANALYZE range by number.
replace_run_text(slides[12], "Text Placeholder 3", "built in Slides 22–36.", "built in Slides 22–28.")
replace_run_text(slides[20], "Rounded Rectangle 20",
                  "campaign-success model (Slides 22–36)", "campaign-success model (Slides 22–28)")

# Roadmap 2/3 (pos3, index 2) — ANALYZE range, and the two sample-calc citations
replace_run_text(slides[2], "Text Placeholder 3", "(Slides 22–36)", "(Slides 22–28)")
replace_run_text(slides[2], "Rounded Rectangle 4",
                  "(Slides 22–24, 27–29)", "(Slides 22–23, 25)")
replace_run_text(slides[2], "Rounded Rectangle 5", "(Slide 36)", "(Slide 28)")

# Roadmap 3/3 (pos4, index 3) — IMPROVE / CONTROL ranges, and the replacement-interval citation
replace_run_text(slides[3], "Text Placeholder 3", "(Slides 37–42)", "(Slides 29–34)")
replace_run_text(slides[3], "Text Placeholder 3", "(Slides 43–46)", "(Slides 35–38)")
replace_run_text(slides[3], "Rounded Rectangle 4", "(Slide 44)", "(Slide 36)")

# old pos38 "The Hidden Assumption" -> new pos30 (index 29)
replace_run_text(slides[29], "TextBox 7", "(Slides 22–36)", "(Slides 22–28)")

# old pos39 "When MTBF Lies" -> new pos31 (index 30): the general-formula slide it
# cited (old pos36) no longer exists on its own -- that formula is now stated on
# Appendix I (new Slide 28), so re-point the citation there instead of deleting it.
replace_run_text(slides[30], "TextBox 10",
                  "Slide 36's RCM/Weibull relation, h(t)",
                  "Appendix I's (Slide 28) RCM/Weibull relation, h(t)")
replace_run_text(slides[30], "TextBox 10",
                  "The worked example already in this deck shows",
                  "Appendix I's worked example shows")

# old pos45 "Reset to New" -> new pos37 (index 36)
replace_run_text(slides[36], "TextBox 9", "(Slides 22–36)", "(Slides 22–28)")

# old pos46 "Reliability Governance -- Closing the Loop" -> new pos38 (index 37);
# "Slide 7" and "Slide 13" are both <=21, unaffected by this shift.
replace_run_text(slides[37], "TextBox 5", "(Slide 37:", "(Slide 29:")

print("cross-reference fix pass done")

# ============================================================ SAVE ============
prs.save(OUT)
print(f"saved {OUT}, {len(prs.slides)} slides")
