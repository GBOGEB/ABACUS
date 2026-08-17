"""
extract_content.py -- dumps every text-bearing shape on every slide of
final4.pptx (== QPS_MTBF_WCS_DMAIC_v4.pptx) into content.yaml, decoupled
from all styling. Shape order follows z-order (document order), which is
also natural reading order for this deck's layouts.

Speaker notes (if present) are included per slide. DMAIC tag text is
called out separately since it drives STEP 4's colour-badge logic in
build_deck4.py.
"""
import yaml
from pptx import Presentation

SRC = "final4.pptx"
OUT = "content.yaml"

TAGS = {"DEFINE", "MEASURE", "ANALYZE", "IMPROVE", "CONTROL"}


def shape_text(sh):
    if not sh.has_text_frame:
        return None
    txt = sh.text_frame.text.replace("\x0b", "\n").strip()
    return txt if txt else None


def is_table(sh):
    return sh.has_table


def dump_table(tbl):
    rows = []
    for row in tbl.rows:
        rows.append([cell.text.strip() for cell in row.cells])
    return rows


prs = Presentation(SRC)
deck = {"source_file": SRC, "slide_count": len(prs.slides.__iter__.__self__._sldIdLst), "slides": []}

for i, slide in enumerate(prs.slides, 1):
    entry = {"index": i, "elements": []}
    dmaic = None
    for sh in slide.shapes:
        if is_table(sh):
            entry["elements"].append({
                "shape": sh.name,
                "type": "table",
                "rows": dump_table(sh.table),
            })
            continue
        txt = shape_text(sh)
        if txt is None:
            continue
        if txt.upper() in TAGS:
            dmaic = txt.upper()
            continue
        entry["elements"].append({
            "shape": sh.name,
            "type": "text",
            "text": txt,
        })
    if dmaic:
        entry["dmaic_phase"] = dmaic
    # speaker notes
    if slide.has_notes_slide:
        notes_txt = slide.notes_slide.notes_text_frame.text.strip()
        if notes_txt:
            entry["speaker_notes"] = notes_txt
    deck["slides"].append(entry)

with open(OUT, "w") as f:
    yaml.dump(deck, f, sort_keys=False, allow_unicode=True, width=100)

print(f"wrote {OUT}: {len(deck['slides'])} slides")
